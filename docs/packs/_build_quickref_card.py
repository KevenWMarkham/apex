"""Build the printable 1-page Sellers Quick Reference Card.

Output: PPTX (single slide, letter portrait) → soffice converts to PDF.
Use: fold in half, slip into a notebook.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy

OUT = r"C:\Stage\Clients\Industries\APEX\docs\reference\Sellers-Quick-Reference-Card.pptx"

# Deloitte palette
GREEN = RGBColor(0x86, 0xBC, 0x25)
NAVY  = RGBColor(0x1A, 0x1F, 0x2E)
DEEP  = RGBColor(0x0E, 0x14, 0x23)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
BG    = RGBColor(0xFA, 0xFA, 0xFA)
INK   = RGBColor(0x12, 0x18, 0x25)
GREY  = RGBColor(0x53, 0x5B, 0x6B)
ACCENT= RGBColor(0xD9, 0x73, 0x06)

prs = Presentation()
prs.slide_width  = Inches(8.5)
prs.slide_height = Inches(11.0)

s = prs.slides.add_slide(prs.slide_layouts[6])

def box(x, y, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background() if line is None else None
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def tb(x, y, w, h, text, size=10, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.10):
    tb_ = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb_.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb_

# ============ Header band ============
box(0, 0, 8.5, 0.85, fill=DEEP)
box(0, 0.85, 8.5, 0.06, fill=GREEN)
tb(0.35, 0.12, 7.8, 0.4, "APEX  ·  CFMP  ·  CHC", size=22, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
tb(0.35, 0.50, 7.8, 0.32, "Sellers Quick Reference  ·  Deloitte's Microsoft Practice (DMTSP)",
   size=11, color=GREEN)

# ============ Section: 60-second pitch ============
y = 1.05
box(0.30, y, 7.90, 1.15, fill=BG, line=NAVY)
tb(0.45, y+0.05, 7.6, 0.25, "THE 60-SECOND PITCH", size=9, bold=True, color=GREEN)
tb(0.45, y+0.30, 7.6, 0.85,
   ("Deloitte has built an agentic-commerce delivery framework called APEX. We've shipped a working merchandise "
    "demo on Microsoft Foundry — 800 SKUs in a vector catalog, voice-driven shopping with Azure Speech, a "
    "14-field cryptographic audit trail, customer journey from CHOOSE through pickup. We've extended it into "
    "the home: cameras, Sonos, kitchen tablet, privacy regime tougher than Alexa by construction. Scoped "
    "engagement to your team: 4-hour BVA workshop → 4-6w Pack Lite live with 3 scenarios for $150-400K. "
    "Demo runs in a browser. Want me to share my screen?"),
   size=9.5, color=INK, line_spacing=1.18)

# ============ Two-column body ============
y = 2.30
COL_W = 3.85
COL2_X = 4.45

# --- LEFT COLUMN ---

# Whiteboard moves
box(0.30, y, COL_W, 1.55, fill=BG, line=NAVY)
tb(0.45, y+0.05, COL_W-0.30, 0.25, "5 WHITEBOARD MOVES", size=9, bold=True, color=GREEN)
moves = [
    ("Alexa contrast", "Amazon (commercial ToS) vs Telco/Merch (CPNI / LEDGER)"),
    ("APEX pancake",   "Scenarios → VVs → 6-Agent → LEDGER → Cloud Profile"),
    ("Journey spine",  "NEED → CHOOSE → SELECT → BUY → SERVICES"),
    ("Two variants",   "B2B2C Telco-fronted   |   B2C Merch-direct"),
    ("Envelope ladder","T1 Foundation → T2 Pack → T3 Scenario → T4 Custom → T5 Operate"),
]
for i, (name, desc) in enumerate(moves):
    yy = y + 0.32 + i*0.23
    tb(0.45, yy, 1.2, 0.22, name, size=8.5, bold=True, color=NAVY)
    tb(1.62, yy, COL_W-1.40, 0.22, desc, size=8.5, color=INK)

# Buyer personas
y2 = y + 1.65
box(0.30, y2, COL_W, 2.10, fill=BG, line=NAVY)
tb(0.45, y2+0.05, COL_W-0.30, 0.25, "BUYER ONE-LINERS", size=9, bold=True, color=GREEN)
personas = [
    ("CMO/CX-VP", "Voice replenishment in YOUR voice; basket routed to you, not Amazon."),
    ("CIO/CDO (Telco)","CPNI is your moat; add $25-40/mo Home Concierge tier on 5G/fiber."),
    ("DPO/CPO", "Constitution blocks at runtime, not policy; consent-hash on every row."),
    ("CISO", "Beelink on-prem privacy floor; raw frames never leave the home."),
    ("Architect","14-interface profile contract; pack manifests touch only the abstraction."),
    ("Procurement","Three additive sub-tiers, no rework. ISV burndown for Microsoft money."),
]
for i, (role, line) in enumerate(personas):
    yy = y2 + 0.32 + i*0.28
    tb(0.45, yy, 0.95, 0.26, role, size=8.5, bold=True, color=NAVY)
    tb(1.40, yy, COL_W-1.20, 0.26, line, size=8.5, color=INK, line_spacing=1.10)

# --- RIGHT COLUMN ---

# Price ladder
box(COL2_X, y, COL_W, 1.85, fill=BG, line=NAVY)
tb(COL2_X+0.15, y+0.05, COL_W-0.30, 0.25, "PRICE LADDER", size=9, bold=True, color=GREEN)
prices = [
    ("BVA Workshop",   "$0-50K",       "1 day"),
    ("T1 Foundation",  "$400-700K",    "6-10w"),
    ("CFMP Pack Lite", "$150-250K",    "4-6w"),
    ("CHC Pack Lite (B2C)","$250-400K","6-8w"),
    ("Pack Standard",  "$500K-1.5M",   "12-16w"),
    ("Pack Enterprise","$1.5-3.5M",    "6-9mo"),
    ("Cross-pack (A)", "$250-400K",    "4-6w"),
    ("T5 Operate",     "$25-110K/mo",  "ongoing"),
]
# Header row
hy = y + 0.32
tb(COL2_X+0.15, hy, 1.65, 0.22, "Tier", size=7.5, bold=True, color=GREY)
tb(COL2_X+1.82, hy, 1.10, 0.22, "Price", size=7.5, bold=True, color=GREY)
tb(COL2_X+2.95, hy, 0.80, 0.22, "Timeline", size=7.5, bold=True, color=GREY)
for i,(t, p, tt) in enumerate(prices):
    yy = y + 0.52 + i*0.18
    tb(COL2_X+0.15, yy, 1.65, 0.18, t, size=8.5, bold=True, color=NAVY)
    tb(COL2_X+1.82, yy, 1.10, 0.18, p, size=8.5, color=INK)
    tb(COL2_X+2.95, yy, 0.80, 0.18, tt, size=8.5, color=GREY)

# Independence
y3 = y + 1.95
box(COL2_X, y3, COL_W, 1.80, fill=BG, line=NAVY)
tb(COL2_X+0.15, y3+0.05, COL_W-0.30, 0.25, "INDEPENDENCE LANGUAGE", size=9, bold=True, color=GREEN)

tb(COL2_X+0.15, y3+0.34, 0.85, 0.22, "NEVER", size=8.5, bold=True, color=ACCENT)
never_items = [
    "“Microsoft partner / alliance”",
    "“Microsoft funds Deloitte”",
    "“ECIF” in client convo",
    "“Joint Microsoft commercial”",
]
for i, t in enumerate(never_items):
    tb(COL2_X+0.95, y3+0.34+i*0.18, COL_W-1.10, 0.18, "• "+t, size=8, color=INK)

tb(COL2_X+0.15, y3+1.10, 0.85, 0.22, "ALWAYS", size=8.5, bold=True, color=GREEN)
always_items = [
    "“Deloitte's Microsoft practice / DMTSP”",
    "“Microsoft money via ISV Marketplace”",
    "“Cloud profile: Microsoft / APEX-M”",
    "“Independent of Microsoft”",
]
for i, t in enumerate(always_items):
    tb(COL2_X+0.95, y3+1.10+i*0.16, COL_W-1.10, 0.16, "• "+t, size=8, color=INK)

# ============ Bottom: top objections + close ============
y4 = 6.55

# Objections
box(0.30, y4, 4.50, 3.10, fill=BG, line=NAVY)
tb(0.45, y4+0.05, 4.30, 0.25, "TOP OBJECTIONS — ONE-LINE HANDLERS", size=9, bold=True, color=GREEN)
objs = [
    ("Just Alexa with a logo?", "Alexa is a commercial provider; we deliver a regulated framework with statutory custodianship."),
    ("Microsoft does this with Foundry.", "Foundry IS the runtime under APEX. We deliver the engagement layer above."),
    ("Framework lock-in?",   "Cloud-neutral. Pack manifests touch only the 14-interface abstraction."),
    ("How long to production?","BVA 4hr → Lite 4-6w live with 3 scenarios → Standard 12-16w."),
    ("Sounds expensive.",     "BVA filters scope; Lite then converts the workshop to a live engagement."),
    ("Privacy / children?",   "Constitution hard-blocks at runtime; raw frames never leave home; <13 not enrolled."),
    ("Switch off Azure later?","Same pack runs on G or A — config change, not re-platform."),
    ("Existing Alexa Skill?", "Skill routes 30% to Amazon AND competes vs Fresh by default. We replace the backend."),
    ("Independence rules?",   "Microsoft money via ISV Marketplace + SI Teaming. Never direct ECIF."),
]
for i, (q, a) in enumerate(objs):
    yy = y4 + 0.34 + i*0.30
    tb(0.45, yy, 4.20, 0.14, "Q: "+q, size=8, bold=True, color=NAVY)
    tb(0.45, yy+0.13, 4.20, 0.18, "A: "+a, size=8, color=INK, line_spacing=1.05)

# Close + URLs
box(4.95, y4, 3.25, 3.10, fill=DEEP, line=None)
tb(5.10, y4+0.05, 3.00, 0.25, "THE STANDARD CLOSE", size=9, bold=True, color=GREEN)
tb(5.10, y4+0.33, 3.00, 1.10,
   ("“I'd like to set up a 4-hour BVA workshop with your CMO/CX team and our delivery lead. "
    "We come with the demo on a laptop, shortlist 3-5 scenarios, and produce a Pack Lite SOW "
    "skeleton in the same session. Half a day in the next two weeks?”"),
   size=8.5, color=RGBColor(0xFF,0xFF,0xFF), line_spacing=1.18)

tb(5.10, y4+1.55, 3.00, 0.22, "5 PITCH LENGTHS", size=9, bold=True, color=GREEN)
five = [
    ("30s","Elevator / LinkedIn DM"),
    ("60s","Discovery-call opener"),
    ("5m","Meeting slot, no demo yet"),
    ("15m","Demo = the pitch"),
    ("45m","Exec briefing w/ v5 deck"),
]
for i,(t, d) in enumerate(five):
    yy = y4 + 1.83 + i*0.18
    tb(5.10, yy, 0.40, 0.16, t, size=8.5, bold=True, color=GREEN)
    tb(5.55, yy, 2.55, 0.16, d, size=8.5, color=RGBColor(0xFF,0xFF,0xFF))

tb(5.10, y4+2.78, 3.00, 0.16, "DEMO URL",       size=7, bold=True, color=GREEN)
tb(5.10, y4+2.92, 3.00, 0.16, "ca-visionkit-portal...azurecontainerapps.io", size=7.5,
   color=RGBColor(0xFF,0xFF,0xFF), font="Consolas")

# ============ Footer ============
box(0, 10.55, 8.5, 0.45, fill=SLATE)
tb(0.30, 10.60, 8.0, 0.18, "Internal · Deloitte's Microsoft Technology & Services Practice",
   size=7.5, color=RGBColor(0xCC,0xD0,0xD9))
tb(0.30, 10.78, 8.0, 0.18, "Prepared by Keven Markham, VP  ·  2026-05-23  ·  Companion to APEX-CFMP-CHC-Sellers-Guide-for-DMTSP.md",
   size=7, color=RGBColor(0x90,0x96,0xA3))

prs.save(OUT)
print(f"saved {OUT}")
