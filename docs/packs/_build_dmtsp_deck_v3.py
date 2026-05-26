"""DMTSP seller deck v3 — photo-forward, minimal text.

Design intent:
  • Speaker delivers content by heart; slides are visual anchors.
  • Every slide has ONE dominant image area + 3–5 short text elements max.
  • Body words per slide stays well under 50 (target 20–35).
  • Image placeholders are clearly marked with a suggested search term so
    the user can drop in stock photos via PowerPoint → Change Picture →
    Stock Images (Microsoft's built-in Unsplash-style library).

Palette is unchanged from v2 (parent APEX-Design spec):
  green #86BC25  navy #1A1F2E  bottom #1E293B  bg #FAFAFA  Calibri throughout
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\APEX-Walkthrough-Deck-for-DMTSP-Sellers-v3.pptx")

# ─── Palette ──────────────────────────────────────────────────────────────
BG          = RGBColor(0xFA, 0xFA, 0xFA)
HEAD_NAVY   = RGBColor(0x1A, 0x1F, 0x2E)
GREEN       = RGBColor(0x86, 0xBC, 0x25)
MS_BLUE     = RGBColor(0x00, 0x78, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
STRONG_INK  = RGBColor(0x1A, 0x1F, 0x2E)
MUTED       = RGBColor(0x47, 0x55, 0x69)
LIGHT_TEXT  = RGBColor(0xCB, 0xD5, 0xE1)
PLACEHOLDER_BG = RGBColor(0xE2, 0xE8, 0xF0)
PLACEHOLDER_FG = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, *, fill, line=None, corner_radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    if corner_radius:
        shp.adjustments[0] = corner_radius
    return shp


def set_text(shape, text, *, size=10, bold=False, color=STRONG_INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36_000); tf.margin_right = Emu(36_000)
    tf.margin_top = Emu(18_000);  tf.margin_bottom = Emu(18_000)
    lines = [text] if isinstance(text, str) else list(text)
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, **kw)
    return box


def photo_placeholder(slide, left, top, width, height, *, prompt: str,
                      caption: str = ""):
    """Renders a styled placeholder rectangle.

    The rendered area looks like a real picture frame and includes the
    suggested stock-image search term. To drop in a real photo: right-click
    the placeholder → 'Change Picture' → 'From Stock Images' (Microsoft 365)
    or 'From a File' / 'From Online Pictures'.
    """
    # Frame
    add_rect(slide, left, top, width, height, fill=PLACEHOLDER_BG)
    # Camera icon (a centered emoji renders as a glyph in PPT)
    add_textbox(slide, left, top + height * 0.18, width, height * 0.30,
                "📷", size=72, color=PLACEHOLDER_FG, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # Prompt — what to search for
    add_textbox(slide, left + Inches(0.25),
                top + height * 0.52,
                width - Inches(0.50),
                Inches(0.45),
                "INSERT IMAGE", size=10, bold=True, color=PLACEHOLDER_FG,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.25),
                top + height * 0.62,
                width - Inches(0.50),
                Inches(0.85),
                "“" + prompt + "”",
                size=11, italic=True, color=PLACEHOLDER_FG,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        add_textbox(slide, left + Inches(0.25),
                    top + height - Inches(0.45),
                    width - Inches(0.50),
                    Inches(0.30),
                    caption, size=8.5, color=PLACEHOLDER_FG,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def header_minimal(slide, title: str, subtitle: str = ""):
    """Thin header — just title + tiny subtitle on slide bg. No big navy strip."""
    # Small green accent dot
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.20), Inches(0.20),
             fill=GREEN)
    add_textbox(slide, Inches(0.85), Inches(0.30), Inches(11), Inches(0.55),
                title, size=28, bold=True, color=STRONG_INK)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(0.90), Inches(11.5), Inches(0.40),
                    subtitle, size=13, color=MUTED, italic=True)


def page_footer(slide, slide_no: int, total: int):
    """Lean footer — just a thin line + page number, no busy attribution strip."""
    # Hairline
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.02),
             fill=PLACEHOLDER_BG)
    add_textbox(slide, Inches(0.5), Inches(7.10), Inches(9), Inches(0.30),
                "APEX  ·  Deloitte Microsoft Technology & Services Practice",
                size=9, color=PLACEHOLDER_FG)
    add_textbox(slide, Inches(11.00), Inches(7.10), Inches(2.0), Inches(0.30),
                f"{slide_no} / {total}", size=9, bold=True, color=GREEN,
                align=PP_ALIGN.RIGHT)


def keyword_chip(slide, left, top, width, height, *, text, color=GREEN,
                 text_color=WHITE):
    add_rect(slide, left, top, width, height, fill=color, corner_radius=0.5)
    add_textbox(slide, left, top, width, height,
                text, size=11, bold=True, color=text_color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─── Build the deck ──────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 15

    # ── Slide 1: Title hero (full-bleed photo) ─────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    photo_placeholder(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
                      prompt="modern retail store interior, customer using smartphone, warm light, wide angle",
                      caption="Half-bleed; subtle dark gradient overlay (Format Picture → 50% dark transparency) before placing the title.")
    # Dark overlay strip across bottom-left for text legibility (over future photo)
    add_rect(s, Inches(0), Inches(4.5), Inches(8.5), Inches(3.0), fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(4.5), Inches(0.20), Inches(3.0), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(4.75), Inches(8), Inches(0.7),
                "APEX", size=72, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(5.65), Inches(8), Inches(0.5),
                "Deloitte's agentic-AI delivery framework",
                size=22, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(6.20), Inches(8), Inches(0.40),
                "Walkthrough for DMTSP  ·  what it is · how to demo · how to sell",
                size=13, color=LIGHT_TEXT, italic=True)
    add_textbox(s, Inches(0.5), Inches(6.80), Inches(8), Inches(0.40),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=10, color=GREEN, bold=True)
    page_footer(s, 1, TOTAL)

    # ── Slide 2: Big quote (left photo / right quote) ──────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    # Photo left half-bleed
    photo_placeholder(s, Inches(0), Inches(0), Inches(6.5), SLIDE_H,
                      prompt="executive boardroom, two people in conversation, neutral light, candid",
                      caption="(half-bleed left)")
    # Quote right
    add_textbox(s, Inches(7.0), Inches(1.8), Inches(0.5), Inches(1.0),
                "“", size=120, bold=True, color=GREEN)
    add_textbox(s, Inches(7.0), Inches(2.6), Inches(6), Inches(2.2),
                ["The engagement layer", "cloud vendors don't ship."],
                size=34, bold=True, color=STRONG_INK)
    add_textbox(s, Inches(7.0), Inches(5.0), Inches(6), Inches(1.0),
                "APEX  in  one  sentence.",
                size=14, color=MUTED, italic=True)
    page_footer(s, 2, TOTAL)

    # ── Slide 3: The Problem — two photos, two captions ────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Without APEX  /  With APEX",
                   "Same Foundry runtime. Different deal shape.")
    # Two image columns
    photo_placeholder(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(4.6),
                      prompt="single developer at messy desk surrounded by code, looks tired, dim light",
                      caption="(without APEX)")
    photo_placeholder(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(4.6),
                      prompt="confident team presenting architecture to executive, bright modern office",
                      caption="(with APEX)")
    # Three keyword chips per side, below the photos
    for i, t in enumerate(["STAFF AUG", "BESPOKE DISCOVERY", "MARGIN COMPRESSION"]):
        keyword_chip(s, Inches(0.5 + i * 2.1), Inches(6.35), Inches(1.95), Inches(0.4),
                     text=t, color=HEAD_NAVY)
    for i, t in enumerate(["FIXED-FEE PACKS", "WORM LEDGER", "REPEATABLE MOTION"]):
        keyword_chip(s, Inches(6.85 + i * 2.05), Inches(6.35), Inches(1.95), Inches(0.4),
                     text=t, color=GREEN)
    page_footer(s, 3, TOTAL)

    # ── Slide 4: The 5 pieces — 5 icon cards ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Five pieces. Every APEX engagement.",
                   "Core · Profiles · Packs · Envelopes · LEDGER")
    pieces = [
        ("🏛️",   "CORE",       "Cloud-neutral runtime"),
        ("☁️",   "PROFILES",   "M  /  G  /  A"),
        ("📦",   "PACKS",      "7 in the catalog"),
        ("📑",   "ENVELOPES",  "5 sellable tiers"),
        ("🔗",   "LEDGER",     "WORM hash chain"),
    ]
    cw = Inches(2.45); gap = Inches(0.10); x0 = Inches(0.50)
    for i, (icon, label, sub) in enumerate(pieces):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, Inches(1.9), cw, Inches(4.6), fill=WHITE,
                 line=PLACEHOLDER_BG)
        # Green top stripe
        add_rect(s, cx, Inches(1.9), cw, Inches(0.20), fill=GREEN)
        # Big icon
        add_textbox(s, cx, Inches(2.4), cw, Inches(2.0),
                    icon, size=96, color=STRONG_INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Number
        add_textbox(s, cx, Inches(4.55), cw, Inches(0.32),
                    f"0{i+1}", size=11, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER)
        # Label
        add_textbox(s, cx, Inches(4.90), cw, Inches(0.40),
                    label, size=18, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Sub
        add_textbox(s, cx, Inches(5.45), cw, Inches(0.80),
                    sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    page_footer(s, 4, TOTAL)

    # ── Slide 5: APEX family — 3 huge logos ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same core. Different cloud.",
                   "Three implementations of one 14-interface contract.")
    families = [
        ("APEX-M", "MICROSOFT  /  AZURE", "GA",            MS_BLUE,                                "data center, blue lighting, server racks"),
        ("APEX-G", "GOOGLE  /  GCP",      "BETA Q3 FY27", RGBColor(0x0F, 0x9D, 0x58),             "modern Google-style office, glass walls, daylight"),
        ("APEX-A", "AMAZON  /  AWS",      "BETA Q4 FY27", RGBColor(0xFF, 0x99, 0x00),             "industrial-scale data center, orange accents"),
    ]
    cw = Inches(4.05); gap = Inches(0.15); x0 = Inches(0.50)
    for i, (tag, name, status, color, prompt) in enumerate(families):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Photo
        photo_placeholder(s, cx, Inches(1.9), cw, Inches(3.0),
                          prompt=prompt, caption="")
        # Card under photo
        add_rect(s, cx, Inches(4.95), cw, Inches(1.55), fill=WHITE,
                 line=PLACEHOLDER_BG)
        add_rect(s, cx, Inches(4.95), cw, Inches(0.40), fill=color)
        add_textbox(s, cx, Inches(4.95), cw, Inches(0.40),
                    tag, size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, cx, Inches(5.40), cw, Inches(0.35),
                    name, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        # Status pill
        add_rect(s, cx + Inches(1.0), Inches(5.85),
                 cw - Inches(2.0), Inches(0.35),
                 fill=color, corner_radius=0.5)
        add_textbox(s, cx + Inches(1.0), Inches(5.85),
                    cw - Inches(2.0), Inches(0.35),
                    status, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 5, TOTAL)

    # ── Slide 6: 7 Packs — image grid ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Seven Industry Packs",
                   "Same 10-asset bundle. Different industry per Pack.")
    packs = [
        ("Manufacturing",        "factory floor, production line, workers in safety gear"),
        ("Retail Merch",         "store merchandiser checking planogram with tablet"),
        ("Finance",              "trading desk, multiple monitors with charts"),
        ("Risk & Compliance",    "audit professional with documents, neutral office"),
        ("Customer Experience",  "call center agent with headset, warm light"),
        ("ESG",                  "wind turbine and solar panels, blue sky"),
        ("CFMP",                 "shopper using smartphone in grocery aisle, bright"),
    ]
    # 4×2 grid; last cell is empty
    cw = Inches(3.05); ch = Inches(2.45); gap = Inches(0.10)
    for i, (name, prompt) in enumerate(packs):
        col = i % 4; row = i // 4
        cx = Inches(0.40 + col * (cw / Inches(1) + gap / Inches(1)))
        cy = Inches(1.9 + row * (ch / Inches(1) + gap / Inches(1)))
        # Photo
        photo_placeholder(s, cx, cy, cw, ch - Inches(0.45),
                          prompt=prompt, caption="")
        # Label band
        is_cfmp = name == "CFMP"
        add_rect(s, cx, cy + ch - Inches(0.40), cw, Inches(0.40),
                 fill=GREEN if is_cfmp else HEAD_NAVY)
        add_textbox(s, cx, cy + ch - Inches(0.40), cw, Inches(0.40),
                    (f"⭐  {name}  (live demo)" if is_cfmp else name),
                    size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Last cell — total
    cx = Inches(0.40 + 3 * (3.05 + 0.10))
    cy = Inches(1.9 + 1 * (2.45 + 0.10))
    add_rect(s, cx, cy, cw, ch, fill=HEAD_NAVY)
    add_rect(s, cx, cy, Inches(0.20), ch, fill=GREEN)
    add_textbox(s, cx + Inches(0.30), cy + Inches(0.30),
                cw - Inches(0.50), Inches(0.50),
                "Total catalog", size=12, color=LIGHT_TEXT)
    add_textbox(s, cx + Inches(0.30), cy + Inches(0.80),
                cw - Inches(0.50), Inches(1.2),
                "7", size=96, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.30), cy + Inches(1.95),
                cw - Inches(0.50), Inches(0.40),
                "1 more on the roadmap (Q3 FY28)", size=10, color=LIGHT_TEXT,
                align=PP_ALIGN.CENTER)
    page_footer(s, 6, TOTAL)

    # ── Slide 7: 10-Asset Bundle (visual only — toolbox metaphor) ──────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same engineering structure. Every Pack.",
                   "Ten assets. Same slots. New Packs land in predictable time.")
    # Background photo wide
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(5.5), Inches(4.7),
                      prompt="organized workshop toolbox or kitchen, top-down view, neat compartments",
                      caption="(visual metaphor — neat toolbox)")
    # 10 mini-chips on the right
    assets = [
        "VV manifests", "Scenarios", "Source adapters", "Adaptive cards", "Persona map",
        "Demo data", "BVA worksheet", "Sample SOW", "Acceptance tests", "Runbook + training",
    ]
    x0 = Inches(6.4); y0 = Inches(1.95)
    cw = Inches(3.20); ch = Inches(0.48); cgap = Inches(0.12)
    for i, name in enumerate(assets):
        col = i % 2; row = i // 2
        cx = x0 + Inches(col * (cw / Inches(1) + cgap / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + cgap / Inches(1)))
        # Card with number + name
        add_rect(s, cx, cy, cw, ch, fill=WHITE, line=PLACEHOLDER_BG)
        add_rect(s, cx, cy, Inches(0.50), ch, fill=GREEN)
        add_textbox(s, cx, cy, Inches(0.50), ch,
                    f"{i+1:02d}", size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, cx + Inches(0.60), cy, cw - Inches(0.70), ch,
                    name, size=12, bold=True, color=STRONG_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 7, TOTAL)

    # ── Slide 8: CFMP spotlight (single hero photo + journey ribbon) ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Spotlight  —  CFMP",
                   "Customer Focused Merchandise Pack v0.2  ·  the live worked example")
    # Full-width hero photo
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(12.33), Inches(3.4),
                      prompt="modern grocery store customer using smartphone in dairy or beverage aisle, candid, bright daylight",
                      caption="(full-width hero — the customer at the moment of decision)")
    # Journey ribbon (4 phases)
    phases = [
        ("CHOOSE",   "What should I buy?",   MS_BLUE),
        ("SELECT",   "Where is it?",         GREEN),
        ("BUY",      "How do I pay?",        RGBColor(0xFF, 0x99, 0x00)),
        ("SERVICES", "What's next?",         RGBColor(0xB4, 0x42, 0xCC)),
    ]
    cw = Inches(3.05); gap = Inches(0.10); x0 = Inches(0.5)
    for i, (label, sub, color) in enumerate(phases):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        add_rect(s, cx, Inches(5.45), cw, Inches(1.10), fill=color)
        add_textbox(s, cx, Inches(5.55), cw, Inches(0.50),
                    label, size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, cx, Inches(6.00), cw, Inches(0.45),
                    sub, size=11, color=WHITE, italic=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 8, TOTAL)

    # ── Slide 9: Demo URL — minimal, prominent ─────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=HEAD_NAVY)
    photo_placeholder(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
                      prompt="presenter at laptop showing screen to two clients, modern conference room",
                      caption="(full-bleed; apply 65% dark transparency over photo for text legibility)")
    # Overlay text
    add_rect(s, Inches(0), Inches(2.0), SLIDE_W, Inches(3.5), fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(2.0), Inches(0.20), Inches(3.5), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(2.15), Inches(12.5), Inches(0.5),
                "Open the demo. Talk through it.",
                size=32, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.85), Inches(12.5), Inches(0.45),
                "Eight surfaces · ten minutes · let the agent answer for you.",
                size=15, color=LIGHT_TEXT, italic=True)
    # The URL
    add_rect(s, Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.85), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.85),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 8 step pills along bottom
    steps = ["1. Demo /", "2. Where is X?", "3. Dietary chip", "4. Cart click",
             "5. /architecture", "6. Details ⓘ", "7. APEX value", "8. Tour flyout"]
    cw = Inches(1.50); gap = Inches(0.05); x0 = Inches(0.55)
    for i, t in enumerate(steps):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        keyword_chip(s, cx, Inches(5.9), cw, Inches(0.55),
                     text=t, color=GREEN, text_color=WHITE)
    page_footer(s, 9, TOTAL)

    # ── Slide 10: Engagement ladder — single visual ────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Engagement ladder",
                   "Five tiers. Stackable. Each tier earns the next.")
    # Left photo
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(5.0), Inches(4.6),
                      prompt="ascending modern staircase, glass-walled office building, looking upward, warm light",
                      caption="(visual metaphor — the ladder)")
    # Right: 5 ascending tier chips
    tiers = [
        ("T1", "FOUNDATION",  "$400–700K  ·  6–10 wks"),
        ("T2", "PACK",        "$100K–$3.5M  ·  4 wks–9 mo"),
        ("T3", "SCENARIO",    "$150–400K  ·  4–8 wks"),
        ("T4", "CUSTOM",      "T&M  ·  Sprints"),
        ("T5", "OPERATE",     "$/mo subscription"),
    ]
    x0 = Inches(6.0); y0 = Inches(1.9); ch = Inches(0.85)
    for i, (tag, name, fund) in enumerate(tiers):
        cy = y0 + Inches(i * (ch / Inches(1) + 0.05))
        # Tag chip
        add_rect(s, x0, cy, Inches(0.85), ch, fill=HEAD_NAVY)
        add_textbox(s, x0, cy, Inches(0.85), ch,
                    tag, size=16, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bar
        add_rect(s, x0 + Inches(0.95), cy, Inches(6.4), ch,
                 fill=WHITE, line=PLACEHOLDER_BG)
        add_textbox(s, x0 + Inches(1.1), cy + Inches(0.05),
                    Inches(6.0), Inches(0.40),
                    name, size=16, bold=True, color=STRONG_INK)
        add_textbox(s, x0 + Inches(1.1), cy + Inches(0.45),
                    Inches(6.0), Inches(0.40),
                    fund, size=11, color=MUTED)
    page_footer(s, 10, TOTAL)

    # ── Slide 11: Independence — single dominant visual ────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Independence  —  non-negotiable",
                   "Deloitte audits Microsoft. SEC rules shape funding.")
    # Left: photo
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(5.0), Inches(4.6),
                      prompt="scales of justice statue on legal desk, soft natural light",
                      caption="(visual metaphor)")
    # Right: two chip stacks
    do_items = ["BVA", "DCIF", "ISV Marketplace", "SI Teaming POC", "T&M  ·  Client direct"]
    dont_items = ["Direct ECIF", "“Partnership”", "Microsoft co-sell credit", "Selling APEX itself"]
    # DO header
    add_textbox(s, Inches(6.0), Inches(1.95), Inches(3.4), Inches(0.4),
                "✓  DO", size=20, bold=True, color=GREEN)
    for i, t in enumerate(do_items):
        keyword_chip(s, Inches(6.0), Inches(2.40 + i * 0.55),
                     Inches(3.4), Inches(0.45),
                     text=t, color=GREEN)
    # DO NOT header
    add_textbox(s, Inches(9.5), Inches(1.95), Inches(3.4), Inches(0.4),
                "✕  DO NOT", size=20, bold=True, color=RGBColor(0xA0, 0x35, 0x35))
    for i, t in enumerate(dont_items):
        keyword_chip(s, Inches(9.5), Inches(2.40 + i * 0.55),
                     Inches(3.4), Inches(0.45),
                     text=t, color=RGBColor(0xA0, 0x35, 0x35))
    page_footer(s, 11, TOTAL)

    # ── Slide 12: NOT a product (binary visual) ────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "What we sell  /  What we don't",
                   "The single most important framing.")
    # Two halves
    add_rect(s, Inches(0.5), Inches(1.95), Inches(6.1), Inches(4.7),
             fill=RGBColor(0xF8, 0xEE, 0xEE), line=RGBColor(0xA0, 0x35, 0x35))
    add_textbox(s, Inches(0.5), Inches(2.15), Inches(6.1), Inches(0.8),
                "✕", size=80, bold=True, color=RGBColor(0xA0, 0x35, 0x35),
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.3), Inches(6.1), Inches(0.7),
                "NOT a product",
                size=28, bold=True, color=RGBColor(0xA0, 0x35, 0x35),
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(4.1), Inches(5.7), Inches(2.4),
                ["Not a SKU.",
                 "Not a license.",
                 "Not a subscription a client can buy.",
                 "Not Microsoft co-sell."],
                size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(6.85), Inches(1.95), Inches(6.0), Inches(4.7),
             fill=RGBColor(0xEE, 0xF8, 0xEE), line=GREEN)
    add_textbox(s, Inches(6.85), Inches(2.15), Inches(6.0), Inches(0.8),
                "✓", size=80, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.85), Inches(3.3), Inches(6.0), Inches(0.7),
                "Deloitte delivery framework",
                size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.05), Inches(4.1), Inches(5.6), Inches(2.4),
                ["We sell  agentic-AI delivery services.",
                 "Scoped per Pack. Per tier.",
                 "Per SOW under a Service Envelope.",
                 "Framework is Deloitte IP. Always."],
                size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)
    page_footer(s, 12, TOTAL)

    # ── Slide 13: BVA wedge — 5 visual steps ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The wedge  —  BVA  →  Pack Lite  →  Standard  →  Enterprise  →  Operate",
                   "Each step earns the next. Nothing gets redone.")
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(12.33), Inches(2.6),
                      prompt="workshop sticky notes on glass wall, hands gesturing, post-its in clusters",
                      caption="(visual metaphor — the BVA workshop kickoff)")
    waves = [
        ("BVA",        "4 hrs",  "free"),
        ("Pack Lite",  "4-6 wk", "$150-250K"),
        ("Standard",   "12-16w", "$500K-1.5M"),
        ("Enterprise", "6-9 mo", "$1.5-3.5M"),
        ("Operate",    "ongoing","$/mo"),
    ]
    cw = Inches(2.45); gap = Inches(0.10); x0 = Inches(0.5)
    for i, (name, time, fund) in enumerate(waves):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Big number ball
        ball = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx + Inches(1.00), Inches(4.7),
                                   Inches(0.45), Inches(0.45))
        ball.fill.solid(); ball.fill.fore_color.rgb = GREEN
        ball.line.fill.background()
        tfc = ball.text_frame; tfc.clear()
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc.margin_left = Emu(0); tfc.margin_right = Emu(0)
        pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run(); rc.text = str(i + 1)
        rc.font.name = "Calibri"; rc.font.size = Pt(15); rc.font.bold = True
        rc.font.color.rgb = WHITE
        # Card
        add_rect(s, cx, Inches(5.25), cw, Inches(1.5), fill=WHITE,
                 line=PLACEHOLDER_BG)
        add_textbox(s, cx, Inches(5.30), cw, Inches(0.45),
                    name, size=18, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cx, Inches(5.78), cw, Inches(0.35),
                    time, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        add_rect(s, cx + Inches(0.2), Inches(6.18),
                 cw - Inches(0.4), Inches(0.40), fill=GREEN, corner_radius=0.3)
        add_textbox(s, cx + Inches(0.2), Inches(6.18),
                    cw - Inches(0.4), Inches(0.40),
                    fund, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 13, TOTAL)

    # ── Slide 14: Q&A — questions only, you answer ────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The questions you'll hear",
                   "Six openers cover ~90% of the first ten minutes. You know the answers.")
    photo_placeholder(s, Inches(0.5), Inches(1.9), Inches(4.5), Inches(4.6),
                      prompt="client meeting, two executives leaning forward asking question, modern office",
                      caption="")
    # 6 question pills on the right
    questions = [
        "Is this Microsoft's product?",
        "Can we buy APEX?",
        "How is this different from Foundry?",
        "What if we're on AWS or Google?",
        "How fast can we see something working?",
        "What does the audit story look like?",
    ]
    x0 = Inches(5.3); y0 = Inches(1.95); ch = Inches(0.70); cgap = Inches(0.10)
    for i, q in enumerate(questions):
        cy = y0 + Inches(i * (ch / Inches(1) + cgap / Inches(1)))
        # Question chip
        add_rect(s, x0, cy, Inches(7.55), ch, fill=WHITE, line=PLACEHOLDER_BG)
        add_rect(s, x0, cy, Inches(0.55), ch, fill=GREEN)
        add_textbox(s, x0, cy, Inches(0.55), ch,
                    "Q", size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x0 + Inches(0.70), cy, Inches(6.75), ch,
                    q, size=14, color=STRONG_INK, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 14, TOTAL)

    # ── Slide 15: This week — three actions ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    photo_placeholder(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
                      prompt="business professionals walking forward through bright lobby, motion blur, golden hour",
                      caption="(full-bleed; apply 65% dark transparency overlay)")
    # Dark band for text
    add_rect(s, Inches(0), Inches(1.0), SLIDE_W, Inches(5.4), fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(1.0), Inches(0.20), Inches(5.4), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.7),
                "This week", size=48, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(1.95), Inches(12), Inches(0.45),
                "Three actions. Every DMTSP seller. Every week.",
                size=16, color=LIGHT_TEXT, italic=True)

    actions = [
        ("01", "OPEN  the  demo",       "60 minutes  ·  one client meeting"),
        ("02", "BOOK  the  BVA",        "4 hours  ·  client-funded"),
        ("03", "SIGN  Pack Lite SOW",   "DCIF  ·  $150-250K  ·  4-6 wks"),
    ]
    y0 = Inches(2.7)
    for i, (n, t, sub) in enumerate(actions):
        cy = y0 + Inches(i * 1.15)
        # Big number
        add_textbox(s, Inches(0.7), cy, Inches(2), Inches(1.0),
                    n, size=72, bold=True, color=GREEN,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Title + sub
        add_textbox(s, Inches(2.7), cy + Inches(0.10),
                    Inches(10), Inches(0.55),
                    t, size=28, bold=True, color=WHITE)
        add_textbox(s, Inches(2.7), cy + Inches(0.65),
                    Inches(10), Inches(0.40),
                    sub, size=13, color=LIGHT_TEXT, italic=True)

    # URL at very bottom
    add_rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.45), fill=GREEN)
    add_textbox(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.45),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 15, TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"size: {OUT.stat().st_size:,} bytes · {TOTAL} slides")


if __name__ == "__main__":
    main()
