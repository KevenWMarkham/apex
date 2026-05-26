"""Inspect the parent APEX-Design-v3.pptx so the DMTSP walkthrough deck
can match its visual style exactly."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

p = Presentation(r"C:\Users\kmarkham\Downloads\APEX-Design-v3.pptx")
print(f"Slides:     {len(p.slides)}")
print(f"Slide size: {p.slide_width/914400:.2f}in x {p.slide_height/914400:.2f}in")
print(f"Masters:    {len(p.slide_masters)}")
print()

for sm_i, sm in enumerate(p.slide_masters):
    print(f"=== Master {sm_i} ===")
    for sl in sm.slide_layouts:
        print(f"   layout: {sl.name}")
    print()

# Look at slides 2-4 for header/body patterns
for idx in [1, 2, 3, 6, 8, 10]:
    if idx >= len(p.slides): break
    s = p.slides[idx]
    print(f"=== Slide {idx+1} · layout={s.slide_layout.name} ===")
    for sh in s.shapes:
        info = f"  [{sh.shape_type}] '{sh.name}'"
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n", " | ")[:60]
            info += f" txt='{txt}'"
            # Get font from first run
            for para in sh.text_frame.paragraphs[:1]:
                for run in para.runs[:1]:
                    if run.font.name:
                        info += f" font='{run.font.name}'"
                    if run.font.size:
                        info += f" size={run.font.size.pt}pt"
                    if run.font.bold:
                        info += f" BOLD"
                    if run.font.color and run.font.color.type is not None:
                        try:
                            if run.font.color.rgb:
                                info += f" color=#{run.font.color.rgb}"
                        except Exception:
                            pass
                    break
                break
        # Position
        if sh.left is not None and sh.top is not None:
            info += f" @ {sh.left/914400:.2f},{sh.top/914400:.2f}"
            info += f" {sh.width/914400:.2f}x{sh.height/914400:.2f}"
        try:
            if hasattr(sh, "fill") and sh.fill.type == 1:
                if sh.fill.fore_color.rgb is not None:
                    info += f" fill=#{sh.fill.fore_color.rgb}"
        except Exception:
            pass
        print(info)
    print()
