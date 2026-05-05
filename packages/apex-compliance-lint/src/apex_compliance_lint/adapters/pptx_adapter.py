"""PPTX adapter — Sprint 29 Task 29.9.2.

Stdlib-only PPTX text extractor — pulls text-frames from each slide's
XML. Optional ``python-pptx`` upgrade path retained for richer
extraction (slide titles, notes pages, group shapes) when the dep is
installed.
"""

from __future__ import annotations

import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from apex_compliance_lint.types import FileAdapter


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class PPTXAdapter(FileAdapter):
    """Extract slide text from a .pptx file."""

    file_extensions = (".pptx",)

    def extract(self, path: Path) -> list[tuple[int, str]]:
        try:
            return self._extract_via_python_pptx(path)
        except ImportError:
            return self._extract_via_zip(path)

    @staticmethod
    def _extract_via_python_pptx(path: Path) -> list[tuple[int, str]]:
        from pptx import Presentation  # type: ignore[import-not-found]

        prs = Presentation(str(path))
        out: list[tuple[int, str]] = []
        line_no = 1
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append((line_no, text))
                        line_no += 1
            # Notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    out.append((line_no, notes_text))
                    line_no += 1
        return out

    @staticmethod
    def _extract_via_zip(path: Path) -> list[tuple[int, str]]:
        """Stdlib-only fallback — walks ppt/slides/slideN.xml entries."""
        out: list[tuple[int, str]] = []
        try:
            with zipfile.ZipFile(path) as zf:
                slide_names = sorted(
                    n for n in zf.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                )
                line_no = 1
                for slide_name in slide_names:
                    try:
                        data = zf.read(slide_name)
                    except KeyError:
                        continue
                    try:
                        root = ET.fromstring(data)
                    except ET.ParseError:
                        continue
                    for t in root.iter(f"{{{A_NS}}}t"):
                        text = (t.text or "").strip()
                        if text:
                            out.append((line_no, text))
                            line_no += 1
        except zipfile.BadZipFile:
            pass
        return out
