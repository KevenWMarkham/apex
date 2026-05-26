"""Build the BVA Workshop Facilitator Deck (4 hours).

PPTX deck the LCSP/facilitator uses to run a 4-hour BVA workshop with the client's
CMO/CX/DPO team. Includes agenda, thesis, demo flow, two interactive exercises
(current-state map + scenario shortlisting), ROI walkthrough, SOW skeleton review,
and the close.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

OUT = r"C:\Stage\Clients\Industries\APEX\docs\reference\BVA-Workshop-Facilitator-Deck.pptx"

# Deloitte palette
GREEN = RGBColor(0x86, 0xBC, 0x25)
NAVY  = RGBColor(0x1A, 0x1F, 0x2E)
DEEP  = RGBColor(0x0E, 0x14, 0x23)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
BG    = RGBColor(0xFA, 0xFA, 0xFA)
INK   = RGBColor(0x12, 0x18, 0x25)
GREY  = RGBColor(0x53, 0x5B, 0x6B)
LITE  = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT= RGBColor(0xD9, 0x73, 0x06)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE= RGBColor(0xF0, 0x93, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width.inches; H = prs.slide_height.inches

def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(s, x, y, w, h, fill=None, line=None, line_w=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rounded(s, x, y, w, h, fill=None, line=None, line_w=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.10
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def tb(s, x, y, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
       font="Calibri", line_spacing=1.15, italic_=False):
    sh = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic_
        r.font.color.rgb = color
    return sh

def header(s, eyebrow, title):
    box(s, 0, 0, W, 0.05, fill=GREEN)
    tb(s, 0.5, 0.20, W-1.0, 0.34, eyebrow.upper(), size=10, bold=True, color=GREEN, font="Calibri")
    tb(s, 0.5, 0.50, W-1.0, 0.85, title, size=28, bold=True, color=NAVY)
    box(s, 0.5, 1.32, 0.6, 0.04, fill=NAVY)  # short underline accent

def footer(s, slide_num, total, section):
    box(s, 0, H-0.40, W, 0.40, fill=SLATE)
    tb(s, 0.40, H-0.32, W-0.8, 0.25,
       f"BVA Workshop · {section}",
       size=9, color=WHITE)
    tb(s, W-1.5, H-0.32, 1.2, 0.25, f"{slide_num} / {total}", size=9, color=GREY, align=PP_ALIGN.RIGHT)

TOTAL = 30  # final slide count

# ============================================================
# Slide 1: Cover
# ============================================================
s = slide()
box(s, 0, 0, W, H, fill=DEEP)
box(s, 0, 5.6, W, 0.10, fill=GREEN)
tb(s, 0.7, 1.0, W-1.4, 0.6, "APEX · BVA WORKSHOP", size=14, bold=True, color=GREEN)
tb(s, 0.7, 1.55, W-1.4, 1.5, "Business Value Assessment", size=48, bold=True, color=WHITE)
tb(s, 0.7, 2.85, W-1.4, 0.6, "A 4-Hour Conversation with Your CMO, CX, Privacy, and Cloud Architecture Leads",
   size=18, color=GREEN)
tb(s, 0.7, 4.0, W-1.4, 0.6, "Outcome of this session:", size=12, color=GREY)
tb(s, 0.7, 4.35, W-1.4, 1.1,
   "Shortlisted 3-5 in-scope scenarios  ·  Populated ROI calculator  ·  Pack Lite SOW skeleton ready for redline",
   size=14, bold=True, color=WHITE)
tb(s, 0.7, 6.0, 6.0, 0.30, "Deloitte's Microsoft Technology & Services Practice", size=10, color=GREY)
tb(s, 0.7, 6.30, 6.0, 0.30, "Facilitator: [LCSP_NAME]  ·  Client: [CLIENT_NAME]  ·  [DATE]", size=10, color=GREY)

# ============================================================
# Slide 2: Welcome & Introductions  (10 min)
# ============================================================
s = slide()
header(s, "00:00 – 00:10  (10 minutes)", "Welcome & introductions")
tb(s, 0.5, 1.6, W-1.0, 0.4, "Around the room:", size=14, bold=True, color=NAVY)
intros = [
    ("Your name", "Role + how long in seat"),
    ("Your function", "Who depends on you · who you depend on"),
    ("Your priority", "One thing you'd most like this engagement to move"),
    ("Your concern", "One thing that worries you about agentic AI in your business"),
]
for i,(k,v) in enumerate(intros):
    y = 2.10 + i*0.78
    rounded(s, 0.7, y, 5.0, 0.65, fill=LITE, line=GREEN)
    tb(s, 0.85, y+0.10, 4.7, 0.25, k, size=12, bold=True, color=NAVY)
    tb(s, 0.85, y+0.34, 4.7, 0.28, v, size=10.5, color=INK)

# Facilitator notes panel
box(s, 7.0, 1.6, W-7.5, 5.3, fill=BG, line=NAVY)
tb(s, 7.15, 1.70, W-7.7, 0.30, "FACILITATOR NOTES", size=10, bold=True, color=GREEN)
notes = [
    ("Listen for", "Who in the room actually owns the customer in-the-moment relationship. Often it's NOT the most senior person."),
    ("Avoid", "Letting one exec dominate. Use a timer (90s each)."),
    ("Capture", "Each person's stated priority on the easel — these become the BVA scoring criteria later (slide 16)."),
    ("Body language", "Note who sits forward / pulls out a notebook when 'agentic AI' is mentioned. They're your champion."),
]
for i,(k,v) in enumerate(notes):
    y = 2.05 + i*1.20
    tb(s, 7.15, y, 0.80, 0.30, k, size=10, bold=True, color=ACCENT)
    tb(s, 7.15, y+0.30, W-7.7, 0.80, v, size=10, color=INK)

footer(s, 2, TOTAL, "Welcome")

# ============================================================
# Slide 3: Agenda  (5 min)
# ============================================================
s = slide()
header(s, "00:10 – 00:15  (5 minutes)", "Today's agenda")

agenda = [
    ("00:00", "Welcome & introductions",                "10m"),
    ("00:10", "Today's agenda",                          "5m"),
    ("00:15", "The thesis: replacing Alexa for commerce","15m"),
    ("00:30", "What APEX is (and isn't)",                "10m"),
    ("00:40", "The customer-journey spine",              "10m"),
    ("00:50", "LIVE DEMO — 7 chapters",                  "25m"),
    ("01:15", "— Break —",                                "15m"),
    ("01:30", "Exercise 1: Current-state mapping",       "30m"),
    ("02:00", "The CFMP / CHC scenario catalog",         "10m"),
    ("02:10", "Exercise 2: Scenario shortlisting",       "30m"),
    ("02:40", "— Break —",                                "10m"),
    ("02:50", "ROI calculator walkthrough",              "20m"),
    ("03:10", "Pack Lite / Standard / Enterprise sizing","15m"),
    ("03:25", "Funding paths (independence-safe)",       "10m"),
    ("03:35", "Exercise 3: Pack Lite scoping",            "15m"),
    ("03:50", "Next-step commitments + SOW skeleton",    "10m"),
]
for i,(t, name_, dur) in enumerate(agenda):
    col = 0 if i < 8 else 1
    row = i if col == 0 else i - 8
    x0 = 0.5 + col * 6.4
    y0 = 1.65 + row * 0.55
    rounded(s, x0, y0, 6.0, 0.45, fill=BG if i % 2 == 0 else LITE, line=NAVY)
    tb(s, x0+0.10, y0+0.05, 0.65, 0.35, t, size=11, bold=True, color=GREEN)
    tb(s, x0+0.85, y0+0.05, 4.4, 0.35, name_, size=11, color=NAVY, bold=True)
    tb(s, x0+5.25, y0+0.05, 0.7, 0.35, dur, size=11, color=GREY, align=PP_ALIGN.RIGHT)

footer(s, 3, TOTAL, "Agenda")

# ============================================================
# Slide 4-5: The Alexa-Replacement Thesis (15 min, 2 slides)
# ============================================================
s = slide()
header(s, "00:15 – 00:30  (15 minutes)", "The Alexa problem — and why it's solvable")
tb(s, 0.5, 1.6, W-1.0, 0.5,
   "Amazon Alexa routes 70M US households' agentic-commerce queries to Amazon Fresh by default.",
   size=16, bold=True, color=NAVY)

problems = [
    ("MERCHANDISER", "Amazon competes for the basket. Skill placement = pay Amazon AND lose share.", ORANGE),
    ("TELCO",        "Alexa rides the pipe and bypasses every customer-facing ARPU opportunity.", ORANGE),
    ("HOUSEHOLD",    "Always-listening · no statutory privacy floor · ToS-based recourse only.", ORANGE),
    ("REGULATOR",    "FTC 2023 settlement. Ongoing COPPA findings. No statutory custodian.", ORANGE),
]
for i,(role, pain, color) in enumerate(problems):
    y = 2.30 + i*0.95
    rounded(s, 0.7, y, 12.0, 0.80, fill=BG, line=color)
    box(s, 0.7, y, 0.20, 0.80, fill=color)
    tb(s, 1.0, y+0.10, 2.2, 0.30, role, size=12, bold=True, color=color)
    tb(s, 3.2, y+0.10, 9.4, 0.60, pain, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 4, TOTAL, "Thesis")

# Slide 5 — the answer
s = slide()
header(s, "Thesis  (continued)", "There is no neutral alternative — until now")
tb(s, 0.5, 1.6, W-1.0, 0.5,
   "APEX delivers what Alexa structurally cannot.",
   size=16, bold=True, color=NAVY)

# Two columns: Alexa | APEX
col_w = 5.9; gap = 0.35; x0 = 0.5
y0 = 2.25
# Header bars
box(s, x0, y0, col_w, 0.45, fill=ORANGE)
tb(s, x0+0.20, y0+0.07, col_w-0.4, 0.32, "ALEXA  (Amazon)", size=14, bold=True, color=WHITE)
box(s, x0+col_w+gap, y0, col_w, 0.45, fill=GREEN)
tb(s, x0+col_w+gap+0.20, y0+0.07, col_w-0.4, 0.32, "APEX  (Deloitte framework)", size=14, bold=True, color=WHITE)

contrast = [
    ("Privacy custodian",       "Amazon (commercial ToS)",                "Telco (CPNI statutory) or Merch (LEDGER)"),
    ("Audit trail",             "None disclosed",                          "14-field WORM hash chain · apex-replay"),
    ("Hard rules",              "Amazon-internal ToS",                     "constitution.yaml — blocks at runtime"),
    ("Multi-merchant",          "Skill-mediated, 30% Amazon cut",          "Pack-mediated, no platform tax"),
    ("In-home privacy",         "'Mute' button (advisory)",                "hard_in_home mode (enforced) · raw frames stay home"),
    ("DPO replay",              "Not available",                           "apex-replay reproduces byte-identical with consent state"),
    ("Cross-cloud",             "Locked to AWS",                           "APEX-M / -G / -A · config change"),
]
for i,(k,a,b) in enumerate(contrast):
    y = y0 + 0.55 + i*0.58
    tb(s, 0.5, y, col_w*2 + gap, 0.20, k, size=10, bold=True, color=GREY)
    rounded(s, x0, y+0.20, col_w, 0.32, fill=BG, line=ORANGE)
    tb(s, x0+0.15, y+0.22, col_w-0.3, 0.28, a, size=10.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    rounded(s, x0+col_w+gap, y+0.20, col_w, 0.32, fill=LITE, line=GREEN)
    tb(s, x0+col_w+gap+0.15, y+0.22, col_w-0.3, 0.28, b, size=10.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 5, TOTAL, "Thesis")

# ============================================================
# Slide 6: What APEX is (and isn't)  (10 min)
# ============================================================
s = slide()
header(s, "00:30 – 00:40  (10 minutes)", "What APEX is — and what it is NOT")

# Two cards
y0 = 1.7
rounded(s, 0.5, y0, 6.1, 5.0, fill=LITE, line=GREEN, line_w=2)
tb(s, 0.7, y0+0.15, 5.7, 0.35, "WHAT APEX IS", size=14, bold=True, color=GREEN)
is_items = [
    "A Deloitte agentic-AI delivery framework",
    "Cloud-neutral: APEX-M / APEX-G / APEX-A (14 interfaces)",
    "10-asset Industry Solution Pack bundle (CFMP, MERML, etc.)",
    "Same 6-Agent Fleet across every Pack",
    "14-field WORM hash-chain LEDGER on every decision",
    "Service Envelope tiers: Lite → Standard → Enterprise → Operate",
    "BVA-funded entry; ISV Marketplace and DCIF funding paths",
    "Microsoft money rides Marketplace, never direct ECIF",
    "Independence-safe by construction",
]
for i, t in enumerate(is_items):
    tb(s, 0.85, y0+0.55 + i*0.45, 5.4, 0.42, f"✓  {t}", size=11.5, color=INK)

rounded(s, 6.85, y0, 6.0, 5.0, fill=BG, line=ORANGE, line_w=2)
tb(s, 7.05, y0+0.15, 5.6, 0.35, "WHAT APEX IS NOT", size=14, bold=True, color=ORANGE)
not_items = [
    "A product license (no SKU)",
    "A Microsoft co-sell (independence-safe by design)",
    "A platform we ask you to host",
    "An Alexa-style commercial provider",
    "A bespoke greenfield build for first-time clients",
    "A black box — Constitution and Pack manifests are inspectable",
    "Owned by Microsoft (the model is one of 14 interfaces)",
    "A one-time deployment (Operate tier is durable run-rate)",
    "Locked to Azure (Pack manifests cloud-portable)",
]
for i, t in enumerate(not_items):
    tb(s, 7.20, y0+0.55 + i*0.45, 5.3, 0.42, f"✗  {t}", size=11.5, color=INK)

footer(s, 6, TOTAL, "Framework")

# ============================================================
# Slide 7: Customer Journey Spine  (10 min)
# ============================================================
s = slide()
header(s, "00:40 – 00:50  (10 minutes)", "The customer-journey spine")
tb(s, 0.5, 1.6, W-1.0, 0.4,
   "Every CFMP / CHC scenario maps to exactly one moment in the customer's day.",
   size=14, color=NAVY)

phases = [
    ("NEED",     "What are we running out of?", "Pantry · fridge · consumption pattern", "home", LITE, GREEN),
    ("CHOOSE",   "What should we get?",          "Recipes · pairings · offers",          "home + store", BG, NAVY),
    ("SELECT",   "Where & when?",                "Wayfinding · OSA · planogram",         "store + home", LITE, GREEN),
    ("BUY",      "Confirm & pay",                "Queue · BOPIS · SCO · cart",           "store + home", BG, NAVY),
    ("SERVICES", "How did it go?",               "Loyalty · complaint · NPS",            "cross-channel", LITE, GREEN),
]
cx = 0.5; cw = (W - 1.0) / 5; gap = 0.10
for i,(name_, q, examples, where, fill, line) in enumerate(phases):
    x = cx + i * cw
    rounded(s, x+gap, 2.30, cw-gap*2, 3.20, fill=fill, line=line, line_w=1.5)
    tb(s, x+gap, 2.40, cw-gap*2, 0.45, name_, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s, x+gap+0.10, 2.85, cw-gap*2-0.20, 0.50, q, size=11, italic_=False, color=GREY, align=PP_ALIGN.CENTER, bold=True)
    tb(s, x+gap+0.10, 3.45, cw-gap*2-0.20, 1.30, examples, size=10, color=INK, align=PP_ALIGN.CENTER)
    tb(s, x+gap+0.10, 4.85, cw-gap*2-0.20, 0.30, where.upper(), size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # arrow
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x+cw-0.20), Inches(3.5), Inches(0.25), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = GREY
        a.line.fill.background()

tb(s, 0.5, 5.7, W-1.0, 0.4,
   "CFMP base = in-store CHOOSE → SELECT → BUY → SERVICES.    CHC overlay = adds NEED (the household runs low).",
   size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 0.5, 6.10, W-1.0, 0.4,
   "Each phase ships its own scenarios, Adaptive Cards, and Virtual Views.",
   size=10, color=GREY, align=PP_ALIGN.CENTER)

footer(s, 7, TOTAL, "Framework")

# ============================================================
# Slide 8: Demo opener  (25 min total across 7 chapters)
# ============================================================
s = slide()
header(s, "00:50 – 01:15  (25 minutes)", "LIVE DEMO — share screen now")

rounded(s, 1.0, 1.7, W-2.0, 5.0, fill=DEEP, line=GREEN, line_w=2)
tb(s, 1.3, 2.0, W-2.6, 0.6, "Open your laptop. Share screen. Open this URL.", size=18, bold=True, color=GREEN)
tb(s, 1.3, 2.7, W-2.6, 0.6,
   "https://ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io",
   size=14, color=WHITE, font="Consolas")

tb(s, 1.3, 3.6, W-2.6, 0.4, "Seven chapters · ~3 minutes each:", size=12, bold=True, color=GREEN)
chapters = [
    "Ch 1 — Your shopper in your store today",
    "Ch 2 — What happens when they linger (Proactive Associate)",
    "Ch 3 — Wayfinding via Azure Maps Creator",
    "Ch 4 — Cart hits $50: the HITL gate trips",
    "Ch 5 — Now imagine this in the kitchen (CHC)",
    "Ch 6 — The privacy contract (Constitution)",
    "Ch 7 — How this becomes an engagement",
]
for i, c in enumerate(chapters):
    tb(s, 1.5, 4.0 + i*0.32, W-3.0, 0.30, c, size=12, color=WHITE)

tb(s, 1.3, 6.30, W-2.6, 0.4,
   "Capture audience reactions: which chapter draws the leaning-in, which draws the eyebrow-raise. That's your follow-on conversation.",
   size=10, italic_=True, color=GREY)

footer(s, 8, TOTAL, "Demo")

# ============================================================
# Slide 9: Break
# ============================================================
s = slide()
box(s, 0, 0, W, H, fill=DEEP)
tb(s, 0.5, 2.5, W-1.0, 1.2, "— BREAK —", size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 0.5, 4.0, W-1.0, 0.5, "15 minutes", size=24, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.5, 5.0, W-1.0, 0.5,
   "Coffee · check in · the second half of the workshop is interactive",
   size=14, color=GREY, align=PP_ALIGN.CENTER)
tb(s, 0.5, H-0.6, W-1.0, 0.3, f"9 / {TOTAL}", size=9, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 10-11: Exercise 1 — Current-State Mapping  (30 min)
# ============================================================
s = slide()
header(s, "01:30 – 02:00  (30 minutes)", "Exercise 1 — Map your current state")
tb(s, 0.5, 1.6, W-1.0, 0.4,
   "On the easel: a 5-column journey-spine grid. Sticky-notes per row.",
   size=12, color=GREY)

# Grid template (printable)
cw = (W - 1.0) / 5
labels = ["NEED", "CHOOSE", "SELECT", "BUY", "SERVICES"]
y_top = 2.20
rows = [
    ("CURRENT EXPERIENCE", "What does your shopper do at this phase today? Be specific."),
    ("CURRENT FRICTION",   "Where do they hesitate, abandon, or substitute against your interest?"),
    ("CURRENT DATA",       "What do you measure? What goes uninstrumented?"),
    ("CURRENT TECH",       "What runs at this phase today (Alexa, app, kiosk, register)?"),
    ("YOUR GOAL",          "What would you most like to move?"),
]

for i,name_ in enumerate(labels):
    x = 0.5 + i * cw
    box(s, x+0.05, y_top, cw-0.10, 0.42, fill=NAVY)
    tb(s, x+0.10, y_top+0.06, cw-0.20, 0.32, name_, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri,(lab, _hint) in enumerate(rows):
    y = y_top + 0.50 + ri*0.85
    box(s, 0.05, y, 0.45, 0.80, fill=GREEN)
    tb(s, 0.07, y+0.05, 0.43, 0.70, lab, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i in range(5):
        x = 0.5 + i * cw
        rounded(s, x+0.05, y, cw-0.10, 0.80, fill=BG, line=GREY, line_w=0.5)

# Facilitator notes
tb(s, 0.5, 6.85, W-1.0, 0.30,
   "FACILITATOR: each table gets 25m to fill the grid; 5m back-share, one phase per table to read out.",
   size=10, italic_=True, color=GREY)
footer(s, 10, TOTAL, "Exercise 1")

# Slide 11: Exercise 1 — Coaching aid
s = slide()
header(s, "Exercise 1 — Coaching aid", "Sample answers (don't reveal — use as facilitator cues only)")

coach = [
    ("NEED",      "Today: nothing. App reorder is manual. Alexa runs the wake-word. Pantry blind spot is real."),
    ("CHOOSE",    "Email · push notifications · paid social. PDP search dies on niche queries. Recipes are static."),
    ("SELECT",    "Aisle directory on paper or in-app static. No live OSA. Planogram updated weekly."),
    ("BUY",       "Self-checkout queues unpredictable. BOPIS pickup counter chronically over/under-staffed."),
    ("SERVICES",  "Loyalty churn measured monthly. Complaint triage routes to a contact center, 18-72hr SLA."),
]
for i,(k,v) in enumerate(coach):
    y = 1.7 + i*0.95
    rounded(s, 0.5, y, 12.3, 0.80, fill=BG, line=NAVY)
    box(s, 0.5, y, 1.6, 0.80, fill=NAVY)
    tb(s, 0.6, y+0.30, 1.4, 0.30, k, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, 2.30, y+0.10, 10.4, 0.65, v, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 11, TOTAL, "Exercise 1")

# ============================================================
# Slide 12: CFMP / CHC catalog overview  (10 min)
# ============================================================
s = slide()
header(s, "02:00 – 02:10  (10 minutes)", "The CFMP + CHC scenario catalog — 30 in total")

phases_data = [
    ("CHOOSE",   5, ["sampling-table-engagement", "in-store-advertising-impact",
                     "personalized-offer-targeting", "review-summary-generation",
                     "product-search-relevance"]),
    ("SELECT",   6, ["wayfinding-walk-to-product", "on-shelf-availability-oos",
                     "shelf-gap-realtime-dispatch", "aisle-engagement-attribution",
                     "end-cap-display-roi", "shelf-price-label-audit"]),
    ("BUY",      4, ["wait-time-prediction-checkout", "cart-dwell-abandonment-rescue",
                     "self-checkout-shrink-detection", "bopis-pickup-counter-load"]),
    ("SERVICES", 3, ["loyalty-churn-winback ⭐", "loyalty-tier-migration",
                     "product-complaint-triage"]),
    ("NEED (CHC)",3,["pantry-inventory-inference", "low-stock-voice-prompt",
                     "replenishment-cart-build"]),
    ("CHC+",     9, ["substitution-acceptance", "pickup-window-optimization",
                     "delivery-vs-pickup-routing", "voice-confirmation-consent",
                     "out-of-home-fallback", "delivery-arrival-handoff",
                     "substitution-NPS-feedback", "loyalty-fusion-attribution ⭐",
                     "consumption-pattern-learning"]),
]
cw = (W - 1.0) / 3
for i, (name_, n, items) in enumerate(phases_data):
    row = i // 3; col = i % 3
    x = 0.5 + col * cw
    y = 1.70 + row * 2.65
    rounded(s, x+0.10, y, cw-0.20, 2.40, fill=BG, line=NAVY)
    box(s, x+0.10, y, cw-0.20, 0.45, fill=NAVY)
    tb(s, x+0.20, y+0.08, cw-0.40, 0.30, f"{name_}  ({n})", size=13, bold=True, color=WHITE)
    for j, it in enumerate(items[:5]):  # cap to 5 visible
        tb(s, x+0.30, y+0.55 + j*0.30, cw-0.50, 0.28, f"• {it}", size=9.5, color=INK)
    if len(items) > 5:
        tb(s, x+0.30, y+2.10, cw-0.50, 0.25, f"+ {len(items)-5} more", size=9, italic_=True, color=GREY)

footer(s, 12, TOTAL, "Catalog")

# ============================================================
# Slide 13-14: Exercise 2 — Scenario Shortlisting  (30 min)
# ============================================================
s = slide()
header(s, "02:10 – 02:40  (30 minutes)", "Exercise 2 — Dot-voting scenario shortlist")
tb(s, 0.5, 1.55, W-1.0, 0.4,
   "Two dots each. One on 'biggest impact'. One on 'fastest to value'. We tally and find Lite candidates.",
   size=12, color=GREY)

# Scoring grid (sample, 12 scenarios)
shortlist_grid = [
    ("rc-sampling-table-engagement",         "CHOOSE",   "Lite"),
    ("rc-cart-dwell-abandonment-rescue",     "BUY",      "Lite"),
    ("cfmp-wayfinding-walk-to-product",      "SELECT",   "Lite"),
    ("rc-on-shelf-availability-oos",         "SELECT",   "Standard"),
    ("rc-shelf-gap-realtime-restock",        "SELECT",   "Standard"),
    ("rc-customer-wait-time-prediction",     "BUY",      "Standard"),
    ("rc-loyalty-churn-winback ⭐",          "SERVICES", "Standard"),
    ("chc-pantry-inventory-inference",       "NEED",     "Lite (CHC)"),
    ("chc-low-stock-voice-prompt",           "NEED",     "Lite (CHC)"),
    ("chc-substitution-acceptance",          "SELECT",   "Standard (CHC)"),
    ("chc-loyalty-fusion-attribution ⭐",    "SERVICES", "Enterprise"),
    ("chc-delivery-arrival-handoff",         "SERVICES", "Standard (CHC)"),
]

# header
hy = 2.10
box(s, 0.5, hy, 5.0, 0.40, fill=NAVY)
box(s, 5.55, hy, 1.5, 0.40, fill=NAVY)
box(s, 7.10, hy, 1.8, 0.40, fill=NAVY)
box(s, 8.95, hy, 1.6, 0.40, fill=NAVY)
box(s, 10.60, hy, 2.2, 0.40, fill=NAVY)
tb(s, 0.60, hy+0.06, 4.9, 0.28, "Scenario", size=11, bold=True, color=WHITE)
tb(s, 5.65, hy+0.06, 1.4, 0.28, "Phase", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 7.20, hy+0.06, 1.7, 0.28, "Best-fit tier", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 9.05, hy+0.06, 1.4, 0.28, "Impact ●", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 10.70, hy+0.06, 2.0, 0.28, "Speed ◆", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i,(name_, ph, tier) in enumerate(shortlist_grid):
    y = hy + 0.45 + i*0.32
    if i % 2 == 0:
        box(s, 0.5, y, 12.30, 0.32, fill=BG)
    tb(s, 0.60, y+0.05, 4.9, 0.25, name_, size=10, color=INK)
    tb(s, 5.65, y+0.05, 1.4, 0.25, ph, size=10, color=GREY, align=PP_ALIGN.CENTER)
    tb(s, 7.20, y+0.05, 1.7, 0.25, tier, size=10, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
    rounded(s, 9.30, y+0.05, 0.9, 0.22, fill=BG, line=NAVY)
    rounded(s, 10.95, y+0.05, 1.4, 0.22, fill=BG, line=NAVY)

tb(s, 0.5, 6.85, W-1.0, 0.30,
   "FACILITATOR: tally dots live on the easel. Lite = top 3-4 impact + speed combined.",
   size=10, italic_=True, color=GREY)
footer(s, 13, TOTAL, "Exercise 2")

# Slide 14: Exercise 2 — Outcome template
s = slide()
header(s, "Exercise 2 — Outcome", "Your shortlist (capture live)")

rounded(s, 0.5, 1.60, 12.30, 1.50, fill=LITE, line=GREEN, line_w=2)
tb(s, 0.65, 1.72, 12.0, 0.32, "PACK LITE — 3 scenarios committed (Wave 1, 4-6 weeks, $150-400K)",
   size=14, bold=True, color=GREEN)
for i in range(3):
    tb(s, 0.85, 2.10 + i*0.30, 11.8, 0.30, f"{i+1}.  [scenario_id]                                       impact: ___   speed: ___",
       size=11, color=INK, font="Consolas")

rounded(s, 0.5, 3.30, 12.30, 1.50, fill=BG, line=NAVY, line_w=1.5)
tb(s, 0.65, 3.42, 12.0, 0.32, "PACK STANDARD candidates (Wave 2, 12-16 weeks)",
   size=12, bold=True, color=NAVY)
for i in range(4):
    tb(s, 0.85, 3.80 + i*0.24, 11.8, 0.24, f"•  [scenario_id]",
       size=10, color=INK, font="Consolas")

rounded(s, 0.5, 5.00, 12.30, 1.20, fill=BG, line=GREY, line_w=1.5)
tb(s, 0.65, 5.12, 12.0, 0.32, "PACK ENTERPRISE / future (Wave 3+)", size=12, bold=True, color=GREY)
for i in range(3):
    tb(s, 0.85, 5.50 + i*0.24, 11.8, 0.24, f"•  [scenario_id]",
       size=10, color=INK, font="Consolas")

footer(s, 14, TOTAL, "Exercise 2")

# ============================================================
# Slide 15: Second Break (10 min)
# ============================================================
s = slide()
box(s, 0, 0, W, H, fill=DEEP)
tb(s, 0.5, 2.5, W-1.0, 1.2, "— BREAK —", size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 0.5, 4.0, W-1.0, 0.5, "10 minutes", size=24, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.5, 5.0, W-1.0, 0.5, "Coming back to ROI and pricing", size=14, color=GREY, align=PP_ALIGN.CENTER)
tb(s, 0.5, H-0.6, W-1.0, 0.3, f"15 / {TOTAL}", size=9, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 16: ROI calculator walkthrough  (20 min)
# ============================================================
s = slide()
header(s, "02:50 – 03:10  (20 minutes)", "ROI calculator walkthrough")
tb(s, 0.5, 1.55, W-1.0, 0.4,
   "Open: BVA-ROI-Calculator.xlsx — five tabs.",
   size=12, color=GREY)

panels = [
    ("README",       "Color conventions · how to use · sources cited from CFMP §9 and CHC §6"),
    ("Inputs",       "Client baseline metrics in BLUE cells; YELLOW = high-leverage assumption (CFO attention)"),
    ("Scenarios",    "21 scenarios with default lift % from the published designs; edit lift % live with client"),
    ("Outputs",      "Wave 1, 2, 3 + cumulative 3-yr TCV; Value : Cost ratio (target >3x for engagement strength)"),
    ("Sensitivity",  "±25% on 4 swing variables; finds the input that most needs CFO sign-off"),
]
for i,(name_, desc) in enumerate(panels):
    y = 2.10 + i*0.95
    rounded(s, 0.5, y, 12.3, 0.80, fill=BG, line=NAVY)
    box(s, 0.5, y, 2.0, 0.80, fill=NAVY)
    tb(s, 0.6, y+0.27, 1.9, 0.30, name_, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, 2.70, y+0.10, 10.0, 0.65, desc, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 16, TOTAL, "ROI")

# ============================================================
# Slide 17: Pack Tier sizing (15 min)
# ============================================================
s = slide()
header(s, "03:10 – 03:25  (15 minutes)", "Pack tier sizing — Lite / Standard / Enterprise")

tiers = [
    ("LITE",       "$150-250K",  "4-6w",  "3 scenarios live\nOne sub-domain\nBVA + DCIF",            LITE,  GREEN),
    ("STANDARD",   "$500K-$1.5M","12-16w","10 scenarios live\nOne sub-domain in prod\nDCIF + T&M",   BG,    NAVY),
    ("ENTERPRISE", "$1.5-3.5M",  "6-9mo", "Full Pack (18)\nMulti-banner rollout\nOperate-ready",     LITE,  GREEN),
]
cw = (W-1.0) / 3
for i,(name_, price, time, scope, fill, line) in enumerate(tiers):
    x = 0.5 + i * cw
    rounded(s, x+0.10, 1.7, cw-0.20, 4.90, fill=fill, line=line, line_w=2)
    box(s, x+0.10, 1.7, cw-0.20, 0.70, fill=line)
    tb(s, x+0.20, 1.83, cw-0.40, 0.50, name_, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x+0.20, 2.55, cw-0.40, 0.50, price, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(s, x+0.20, 3.10, cw-0.40, 0.40, time,  size=14, color=GREY, align=PP_ALIGN.CENTER)
    box(s, x+0.50, 3.65, cw-1.0, 0.04, fill=line)
    tb(s, x+0.30, 3.80, cw-0.60, 2.50, scope, size=12, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.3)

# Below: T1 Foundation + T5 Operate
rounded(s, 0.5, 6.75, 6.0, 0.55, fill=BG, line=GREY)
tb(s, 0.65, 6.82, 5.8, 0.45, "T1 Foundation (one-time per tenant): $400-700K · 6-10w · BVA+DCIF",
   size=11, color=NAVY)
rounded(s, 6.85, 6.75, 6.0, 0.55, fill=BG, line=GREY)
tb(s, 7.00, 6.82, 5.8, 0.45, "T5 Operate (subscription run-rate): $25-110K/mo · continuous · client direct",
   size=11, color=NAVY)

footer(s, 17, TOTAL, "Sizing")

# ============================================================
# Slide 18: Funding paths (10 min)
# ============================================================
s = slide()
header(s, "03:25 – 03:35  (10 minutes)", "Funding paths — independence-safe by construction")

vehicles = [
    ("BVA (client-funded)",       "Workshop fee — sometimes Deloitte-absorbed", "T1, T2",  GREEN),
    ("DCIF (Deloitte co-invest)", "Deloitte commits balance sheet against milestones", "T1, T2, T3", GREEN),
    ("ISV Marketplace burndown",  "Client's Azure commitment burned down via Marketplace", "T2, T3, T4", GREEN),
    ("SI Teaming POC",            "Non-competing SI receives MSFT funding; Deloitte sub-delivers", "T2, T3",  ACCENT),
    ("Client direct",             "Standard Deloitte-on-Deloitte-paper engagement",  "All tiers",  NAVY),
    ("T&M",                       "For custom build (T4) or extensions",             "T2, T3, T4", NAVY),
]
for i,(k, v, tiers_, color) in enumerate(vehicles):
    y = 1.75 + i*0.75
    rounded(s, 0.5, y, 12.3, 0.65, fill=BG, line=color, line_w=1.5)
    box(s, 0.5, y, 0.20, 0.65, fill=color)
    tb(s, 0.85, y+0.10, 3.5, 0.45, k, size=12, bold=True, color=NAVY)
    tb(s, 4.50, y+0.10, 5.8, 0.45, v, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 10.40, y+0.10, 2.3, 0.45, tiers_, size=11, color=GREEN, align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE)

box(s, 0.5, 6.40, 12.3, 0.70, fill=DEEP)
tb(s, 0.7, 6.48, 12.0, 0.30, "INDEPENDENCE RULE", size=11, bold=True, color=GREEN)
tb(s, 0.7, 6.74, 12.0, 0.30,
   "Deloitte is Microsoft's auditor. We never receive ECIF directly. Microsoft money flows via Marketplace or SI Teaming.",
   size=10, color=WHITE)

footer(s, 18, TOTAL, "Funding")

# ============================================================
# Slide 19: Exercise 3 — Pack Lite scoping  (15 min)
# ============================================================
s = slide()
header(s, "03:35 – 03:50  (15 minutes)", "Exercise 3 — Pack Lite scoping commitment")
tb(s, 0.5, 1.6, W-1.0, 0.4,
   "Together: lock the Lite scope. Capture answers live on this slide.",
   size=12, color=GREY)

questions = [
    ("Three scenarios from Exercise 2 shortlist:", "1) _____________   2) _____________   3) _____________"),
    ("Variant — B2B2C, B2C, or in-store only:",    "[A · B · in-store]"),
    ("Pilot scope — banner / store count / household count:", "Banner: ________  Stores: ___  HH: _____"),
    ("Primary source-system adapter (one):",       "[POS · Planogram · Loyalty · Vision · Beacon]"),
    ("Privacy lead from your team:",                "Name: ________________  Org: ______________"),
    ("Cloud architect engagement contact:",         "Name: ________________  Email: _____________"),
    ("Funding vehicle (per slide 18):",             "[BVA · DCIF · ISV Marketplace · SI Teaming · Direct · T&M]"),
    ("Target SOW countersignature date:",           "____ / ____ / 2026"),
]
for i,(q, a) in enumerate(questions):
    y = 2.10 + i*0.62
    tb(s, 0.6, y, 6.0, 0.30, q, size=11, bold=True, color=NAVY)
    rounded(s, 6.70, y, 6.0, 0.40, fill=BG, line=GREEN)
    tb(s, 6.85, y+0.06, 5.8, 0.30, a, size=10.5, color=INK, font="Consolas")

footer(s, 19, TOTAL, "Exercise 3")

# ============================================================
# Slide 20: SOW skeleton review (5 min) + close (5 min)
# ============================================================
s = slide()
header(s, "03:50 – 04:00  (10 minutes)", "Next-step commitments + SOW skeleton")

tb(s, 0.5, 1.6, W-1.0, 0.4, "Three artifacts leave this room today:", size=14, bold=True, color=NAVY)

artifacts = [
    ("1.", "BVA workshop summary",          "One-page memo · scenarios committed · shortlist · sponsor names · target dates",                         GREEN),
    ("2.", "ROI calculator (populated)",    "Wave 1/2/3 dollar impacts · Value:Cost ratio · sensitivity outputs · CFO-ready sign-off page",            NAVY),
    ("3.", "SOW skeleton (Pack Lite)",      "Pre-filled with this session's outcomes · ready for redline · countersignature target [TARGET_DATE]",     GREEN),
]
for i,(n, name_, desc, color) in enumerate(artifacts):
    y = 2.20 + i*1.20
    rounded(s, 0.5, y, 12.3, 1.05, fill=BG, line=color, line_w=2)
    box(s, 0.5, y, 1.0, 1.05, fill=color)
    tb(s, 0.6, y+0.30, 0.8, 0.45, n, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, 1.65, y+0.15, 11.0, 0.35, name_, size=14, bold=True, color=NAVY)
    tb(s, 1.65, y+0.50, 11.0, 0.50, desc, size=11, color=INK, line_spacing=1.2)

box(s, 0.5, 6.40, 12.3, 0.65, fill=DEEP)
tb(s, 0.7, 6.50, 12.0, 0.30, "THE COMMITMENT", size=11, bold=True, color=GREEN)
tb(s, 0.7, 6.75, 12.0, 0.30,
   "By [target date]: SOW countersigned · BVA artifacts archived in Client + Deloitte SharePoint · engagement team named.",
   size=10, color=WHITE)

footer(s, 20, TOTAL, "Close")

# ============================================================
# Slides 21-29: Reference / backup deck (one-liners)
# ============================================================

# 21 — Buyer personas quick-ref
s = slide()
header(s, "Backup", "Talk tracks by buyer persona")
personas = [
    ("CMO / VP-CX",  "Trip conversion · basket size · NPS · Amazon Fresh share — agentic layer that routes the basket to YOU."),
    ("CIO / CDO (Telco)", "ARPU growth · churn · 5G stickiness · CPNI is your moat — Home Concierge tier on fiber/5G."),
    ("Chief Privacy Officer", "Constitution blocks at runtime · consent_hash on every row · apex-replay for DPO + regulator."),
    ("Chief Digital / Head of Loyalty", "Active loyalty in-the-moment · same identity, same trust, now in the kitchen."),
    ("CISO",         "Beelink home-edge privacy floor · Defender for IoT · 14-field forensic ledger on every decision."),
    ("Cloud Architect", "14 interfaces, profile-swappable · pack manifests touch abstract API only · no lock-in."),
    ("Procurement",  "Additive sub-tiers · independence-safe funding paths · BVA filters scope before commit."),
]
for i,(k,v) in enumerate(personas):
    y = 1.7 + i*0.75
    rounded(s, 0.5, y, 12.3, 0.65, fill=BG, line=NAVY)
    tb(s, 0.7, y+0.10, 3.0, 0.45, k, size=11.5, bold=True, color=NAVY)
    tb(s, 3.80, y+0.10, 8.9, 0.45, v, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 21, TOTAL, "Reference")

# 22 — Top objections
s = slide()
header(s, "Backup", "Top 10 objections — one-line handlers")
objs = [
    ("Isn't this just Alexa with a logo?", "Alexa is a commercial provider; we deliver a regulated framework with statutory custodianship."),
    ("Microsoft does this with Foundry.",  "Foundry IS the runtime under APEX. We deliver the engagement layer above it."),
    ("Framework lock-in?",                  "Cloud-neutral. Manifests touch the 14-interface abstraction only."),
    ("How long to production?",            "BVA 4hr → Lite 4-6w live with 3 scenarios → Standard 12-16w."),
    ("Sounds expensive.",                   "BVA filters scope; Lite then converts the workshop to a live engagement."),
    ("Privacy / children?",                 "Constitution hard-blocks at runtime; raw frames never leave home; <13 not enrolled."),
    ("Switch off Azure later?",            "Same pack runs on G or A — config change, not re-platform."),
    ("We already have an Alexa Skill.",    "Skill routes 30% to Amazon AND competes vs Fresh by default. We replace the backend."),
    ("Independence rules?",                "Microsoft money via Marketplace + SI Teaming. Never direct ECIF."),
    ("We don't have the data.",            "Pack ships synthetic demo data + source adapters + 60+ acceptance tests."),
]
for i,(q,a) in enumerate(objs):
    y = 1.65 + i*0.50
    if i % 2 == 0:
        box(s, 0.5, y, 12.3, 0.45, fill=BG)
    tb(s, 0.7, y+0.04, 4.3, 0.20, "Q: " + q, size=10, bold=True, color=NAVY)
    tb(s, 0.7, y+0.22, 12.0, 0.20, "A: " + a, size=10, color=INK)
footer(s, 22, TOTAL, "Reference")

# 23 — Independence language
s = slide()
header(s, "Backup", "Independence language discipline")
tb(s, 0.5, 1.6, 6.0, 0.4, "NEVER SAY", size=14, bold=True, color=ORANGE)
never = ["Microsoft partner / alliance / co-sell", "Microsoft funds Deloitte / pays Deloitte",
         "ECIF (client-facing)", "Joint Microsoft commercial", "Microsoft GTM"]
for i,t in enumerate(never):
    tb(s, 0.7, 2.10 + i*0.42, 5.8, 0.40, "✗  " + t, size=11.5, color=INK)

tb(s, 7.0, 1.6, 6.0, 0.4, "ALWAYS SAY", size=14, bold=True, color=GREEN)
always = ["Deloitte's Microsoft practice / DMTSP", "Microsoft technology / cloud profile",
          "Microsoft money via ISV Marketplace burndown", "Independent of Microsoft (Deloitte's role)",
          "Cloud profile: Microsoft / APEX-M"]
for i,t in enumerate(always):
    tb(s, 7.2, 2.10 + i*0.42, 5.8, 0.40, "✓  " + t, size=11.5, color=INK)

box(s, 0.5, 5.0, 12.3, 1.8, fill=DEEP)
tb(s, 0.7, 5.20, 12.0, 0.40, "THE INDEPENDENCE SENTENCE (always ready)", size=11, bold=True, color=GREEN)
tb(s, 0.7, 5.65, 12.0, 1.0,
   ("\"Deloitte is Microsoft's auditor. SEC rules forbid us receiving ECIF directly from Microsoft. We've "
    "engineered the funding paths around that constraint — Microsoft money flows via Marketplace burndown "
    "or SI Teaming, never directly to Deloitte. Our language discipline reflects the rule.\""),
   size=11, color=WHITE, line_spacing=1.3)
footer(s, 23, TOTAL, "Reference")

# 24 — Constitution hard rules
s = slide()
header(s, "Backup", "Constitution.yaml — six hard rules (in-home mode)")
rules = [
    "1. No raw frames leave the home — embeddings + canonical event JSON only.",
    "2. No off-list inferences — agents may only emit events for pre-enrolled product list.",
    "3. Per-member consent required for voice-print enrollment. Under 13: cannot be enrolled.",
    "4. \"Privacy minute\" voice command pauses all cameras/mics for 60 minutes.",
    "5. Quiet-hours hard floor: no proactive voice prompts 10pm–7am, ever.",
    "6. No cross-household correlation — partition + LEDGER per-tenant filter at query plane.",
]
for i,r in enumerate(rules):
    y = 1.85 + i*0.78
    rounded(s, 0.5, y, 12.3, 0.65, fill=LITE, line=GREEN, line_w=1.5)
    tb(s, 0.7, y+0.12, 12.0, 0.45, r, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

tb(s, 0.5, 6.85, W-1.0, 0.30,
   "These rules block at the RUNTIME layer — not as policy, as code. Auditable. DPO-reviewable. Inspectable.",
   size=10, italic_=True, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 24, TOTAL, "Reference")

# 25 — Demo BOM
s = slide()
header(s, "Backup", "Consumer hardware BOM — shipping today on Amazon Prime")
bom = [
    ("Edge ML + camera", "NVIDIA Jetson Orin Nano Dev Kit",         "1K+/mo · 371 ratings",  "$249", "Overnight"),
    ("",                 "+ Waveshare aluminum case w/ camera",      "100+/mo · 142 ratings", "$24",  "Tomorrow"),
    ("Smart speaker",    "Sonos Era 100 (Alexa-Enabled — replaceable backend)", "500+/mo · 2,566 ratings", "$189 (was $219)", "Tomorrow"),
    ("Kitchen kiosk",    "Android 16 tablet 11\" Octa-Core 28GB RAM","1K+/mo · 1,078 ratings","$110", "Overnight"),
    ("TOTAL household starter kit", "Per-household BOM",              "",                    "~$572","All Prime"),
]
hy = 1.85
box(s, 0.5, hy, 12.3, 0.45, fill=NAVY)
tb(s, 0.65, hy+0.08, 2.0, 0.30, "Role",       size=11, bold=True, color=WHITE)
tb(s, 2.70, hy+0.08, 4.0, 0.30, "Device",     size=11, bold=True, color=WHITE)
tb(s, 6.80, hy+0.08, 2.5, 0.30, "Live signal",size=11, bold=True, color=WHITE)
tb(s, 9.30, hy+0.08, 1.6, 0.30, "Price",      size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 10.95,hy+0.08, 1.8, 0.30, "Delivery",   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i,(role, dev, sig, pr, deliv) in enumerate(bom):
    y = hy + 0.50 + i*0.65
    is_total = (i == len(bom)-1)
    box(s, 0.5, y, 12.3, 0.60, fill=LITE if is_total else BG)
    tb(s, 0.65, y+0.10, 2.0, 0.40, role, size=11, bold=is_total, color=NAVY if is_total else INK, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 2.70, y+0.10, 4.0, 0.40, dev, size=10.5, color=INK, bold=is_total, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 6.80, y+0.10, 2.5, 0.40, sig, size=9.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 9.30, y+0.10, 1.6, 0.40, pr,  size=11, bold=True, color=GREEN if is_total else INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 10.95,y+0.10, 1.8, 0.40, deliv, size=10, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 25, TOTAL, "Reference")

# 26 — Cross-pack chain (Variant A)
s = slide()
header(s, "Backup", "Variant A cross-pack scenario chain — chc-replenishment-cycle")

rounded(s, 0.5, 1.7, 6.0, 4.8, fill=LITE, line=GREEN, line_w=2)
tb(s, 0.7, 1.85, 5.7, 0.4, "TCP — APEX-M PRIMARY (AT&T tenant)", size=13, bold=True, color=GREEN)
tcp = [
    "Agent 1 — Assess pantry need (tcml.pantry_state)",
    "Agent 2 — Classify urgency / substitutions",
    "Agent 3 — Quantify replen qty (consumption pattern 7d)",
    "Agent 4 — Build household card → sonos_voice OR home_kiosk",
    "Persona household — confirm or modify cart",
]
for i,t in enumerate(tcp):
    tb(s, 0.85, 2.30 + i*0.55, 5.3, 0.45, "• " + t, size=10, color=INK)

a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(3.8), Inches(0.4), Inches(0.6))
a.fill.solid(); a.fill.fore_color.rgb = NAVY; a.line.fill.background()
tb(s, 6.5, 4.45, 0.8, 0.30, "imports:", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, 6.5, 4.70, 0.8, 0.30, "Interface #12", size=8, color=GREY, align=PP_ALIGN.CENTER)

rounded(s, 7.2, 1.7, 5.6, 4.8, fill=BG, line=NAVY, line_w=2)
tb(s, 7.4, 1.85, 5.3, 0.4, "CFMP — APEX-A SATELLITE (Kroger tenant)", size=13, bold=True, color=NAVY)
cfmp = [
    "merml.product_with_planogram_location",
    "merml.delivery_slot_inventory",
    "cxml.loyalty_state",
    "cfmp.add_to_cart",
    "cfmp.reserve_delivery_slot",
    "cfmp.confirm_order",
]
for i,t in enumerate(cfmp):
    tb(s, 7.55, 2.30 + i*0.45, 5.0, 0.40, "• " + t, size=10, color=INK)

box(s, 0.5, 6.55, 12.3, 0.50, fill=DEEP)
tb(s, 0.7, 6.65, 12.0, 0.30,
   "Eight LedgerRows land in AT&T's LEDGER. Kroger writes none. One audit chain across two corporate clouds.",
   size=10, color=GREEN, align=PP_ALIGN.CENTER)
footer(s, 26, TOTAL, "Reference")

# 27 — Three-year sequencing
s = slide()
header(s, "Backup", "Three-year sequencing (joint Variant A + B)")
yrs = [
    ("Y1 FY27", "AT&T × Kroger pilot 1 DMA (DFW or Houston)\n~1,000 households",
                "Kroger Plus Premium pilot 2 banners (Smith's + Mariano's)\n~5,000 households"),
    ("Y2 FY28", "AT&T national rollout ~100K households\n+ Verizon × Walmart pilot",
                "Kroger national + Walmart+ launch"),
    ("Y3 FY29", "All Tier-1 carriers × Tier-1 merchandisers\n~10 pairings · multi-million household scale",
                "Every Tier-1 merch has direct variant\nIn-store + in-home + Variant-A federation share schemas"),
]
for i,(y, a, b) in enumerate(yrs):
    yy = 1.75 + i*1.65
    box(s, 0.5, yy, 1.5, 1.5, fill=NAVY)
    tb(s, 0.65, yy+0.55, 1.3, 0.5, y, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rounded(s, 2.10, yy, 5.3, 1.5, fill=LITE, line=GREEN)
    tb(s, 2.20, yy+0.05, 5.1, 0.30, "VARIANT A (B2B2C)", size=10, bold=True, color=GREEN)
    tb(s, 2.30, yy+0.35, 5.0, 1.10, a, size=10.5, color=INK, line_spacing=1.2)
    rounded(s, 7.50, yy, 5.3, 1.5, fill=BG, line=NAVY)
    tb(s, 7.60, yy+0.05, 5.1, 0.30, "VARIANT B (B2C)", size=10, bold=True, color=NAVY)
    tb(s, 7.70, yy+0.35, 5.0, 1.10, b, size=10.5, color=INK, line_spacing=1.2)
footer(s, 27, TOTAL, "Reference")

# 28 — Demo URL big card
s = slide()
box(s, 0, 0, W, H, fill=DEEP)
box(s, 0, H-0.10, W, 0.10, fill=GREEN)
tb(s, 0.5, 1.6, W-1.0, 0.6, "OPEN THE DEMO ANYTIME", size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rounded(s, 1.0, 2.4, W-2.0, 2.0, fill=NAVY, line=GREEN, line_w=2)
tb(s, 1.0, 2.7, W-2.0, 0.5, "Working demo · live in a browser",
   size=14, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 1.0, 3.25, W-2.0, 0.8,
   "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io",
   size=18, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
tb(s, 1.0, 4.6, W-2.0, 0.4, "/architecture page = APEX framework view",
   size=14, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1.0, 5.0, W-2.0, 0.4, "APEX Tour flyout (right side) = guided walkthrough",
   size=14, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.5, 6.8, W-1.0, 0.3, f"28 / {TOTAL}", size=9, color=GREY, align=PP_ALIGN.CENTER)

# 29 — Related artifacts
s = slide()
header(s, "Backup", "Workshop output kit — what leaves the room with you")
arts = [
    "BVA-ROI-Calculator.xlsx — populated with this session's inputs",
    "SOW-APEX-CFMP-Pack-Lite.docx — pre-filled, ready for redline",
    "SOW-APEX-CFMP-Pack-Standard.docx — for Wave-2 follow-on",
    "SOW-APEX-CFMP-Pack-Enterprise.docx — for Wave-3 program design",
    "CFMP-Pack-Lite-Wedge-for-DMTSP-Sellers.md — the engagement one-pager",
    "Connected-Home-Commerce-v0.1.md — CHC design (B2B2C + B2C)",
    "APEX-CFMP-CHC-Sellers-Guide-for-DMTSP.md — field reference",
    "Sellers-Quick-Reference-Card.pdf — print and fold",
    "This deck (BVA-Workshop-Facilitator-Deck.pptx)",
]
for i,a in enumerate(arts):
    y = 1.7 + i*0.50
    rounded(s, 0.5, y, 12.3, 0.42, fill=BG, line=NAVY)
    tb(s, 0.7, y+0.05, 12.0, 0.32, "• " + a, size=11, color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
footer(s, 29, TOTAL, "Reference")

# 30 — Final close slide
s = slide()
box(s, 0, 0, W, H, fill=DEEP)
box(s, 0, H-0.10, W, 0.10, fill=GREEN)
tb(s, 0.5, 2.4, W-1.0, 1.2, "Thank you.", size=64, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tb(s, 0.5, 3.9, W-1.0, 0.6,
   "Three artifacts. One commitment. SOW countersignature by [TARGET_DATE].",
   size=18, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 0.5, 4.7, W-1.0, 0.5,
   "[LCSP_NAME] · Deloitte's Microsoft Technology & Services Practice",
   size=14, color=GREY, align=PP_ALIGN.CENTER)
tb(s, 0.5, 6.8, W-1.0, 0.3, f"30 / {TOTAL}", size=9, color=GREY, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"saved {OUT}")
