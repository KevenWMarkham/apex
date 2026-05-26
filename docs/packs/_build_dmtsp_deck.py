"""DMTSP seller-walkthrough deck — formatted to match APEX-Design-v3.pptx.

Visual spec extracted from the parent deck:

  • Slide:         13.33" × 7.50" widescreen
  • Background:    #FAFAFA
  • Header strip:  #1A1F2E (full-width, top 1.05") + #86BC25 left accent (0.18" × 1.05")
  • Bottom strip:  #1E293B (full-width, 1.20" tall above page footer)
  • Page footer:   #1A1F2E (full-width, 0.32" at very bottom)
  • Card body:     #FFFFFF
  • Card stripe:   #64748B (left/neutral col) or #86BC25 (right/APEX col)
  • Microsoft:     #0078D4 (used on APEX-M card)
  • Accent green:  #86BC25 (Deloitte signature)
  • Muted body:    #475569
  • Strong text:   #1A1F2E
  • Footer text:   #CBD5E1

All type is Calibri. Header title 22pt bold white; subtitle 12pt #CBD5E1;
card label 11pt bold #1A1F2E; card body 9.5pt #475569; bottom-strip
numbered items 01/02/03/04 + label + body.
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\APEX-Walkthrough-Deck-for-DMTSP-Sellers-v2.pptx")

# ─── Palette (exactly the parent's) ───────────────────────────────────────
BG          = RGBColor(0xFA, 0xFA, 0xFA)
HEAD_NAVY   = RGBColor(0x1A, 0x1F, 0x2E)
BOTTOM_NAVY = RGBColor(0x1E, 0x29, 0x3B)
GREEN       = RGBColor(0x86, 0xBC, 0x25)   # Deloitte signature
GRAY_STRIP  = RGBColor(0x64, 0x74, 0x8B)
MS_BLUE     = RGBColor(0x00, 0x78, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
STRONG_INK  = RGBColor(0x1A, 0x1F, 0x2E)
MUTED       = RGBColor(0x47, 0x55, 0x69)
LIGHT_TEXT  = RGBColor(0xCB, 0xD5, 0xE1)
PALE_BLUE   = RGBColor(0xBF, 0xDB, 0xFE)
PALE_GREEN  = RGBColor(0xEA, 0xF4, 0xD4)   # subtle green tint for highlight rows
ROW_HILIGHT = RGBColor(0xF8, 0xFC, 0xEC)
RED_INK     = RGBColor(0xA0, 0x35, 0x35)
PALE_RED    = RGBColor(0xF8, 0xEE, 0xEE)
PALE_GRN_BG = RGBColor(0xEE, 0xF8, 0xEE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    return shp


def set_text(shape, text, *, size=10, bold=False, color=STRONG_INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36_000); tf.margin_right = Emu(36_000)
    tf.margin_top = Emu(18_000);  tf.margin_bottom = Emu(18_000)
    lines = [text] if isinstance(text, str) else list(text)
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, **kw)
    return box


def header(slide, title: str, subtitle: str = "", tagline_label: str = "",
           tagline_body: str = ""):
    """Standard top-of-slide header used on every content slide."""
    # Full-width navy strip
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), fill=HEAD_NAVY)
    # Green left accent
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(1.05), fill=GREEN)
    # Title (22pt bold white)
    add_textbox(slide, Inches(0.45), Inches(0.12), Inches(12.50), Inches(0.50),
                title, size=22, bold=True, color=WHITE)
    # Subtitle (12pt light)
    if subtitle:
        add_textbox(slide, Inches(0.45), Inches(0.62), Inches(12.50), Inches(0.40),
                    subtitle, size=12, color=LIGHT_TEXT)
    # Tagline label (green) + body (muted slate)
    if tagline_label:
        add_textbox(slide, Inches(0.45), Inches(1.18), Inches(1.20), Inches(0.30),
                    tagline_label, size=10, bold=True, color=GREEN)
        add_textbox(slide, Inches(1.65), Inches(1.18), Inches(11.20), Inches(0.55),
                    tagline_body or "", size=10.5, color=MUTED)


def page_footer(slide, slide_no: int, total: int):
    """Bottom 0.32" navy strip with attribution + page number."""
    add_rect(slide, Inches(0), Inches(7.18), SLIDE_W, Inches(0.32), fill=HEAD_NAVY)
    add_textbox(slide, Inches(0.45), Inches(7.22), Inches(8.0), Inches(0.26),
                "Internal  •  Deloitte Microsoft Technology & Services Practice  •  DMTSP walkthrough",
                size=8.5, color=LIGHT_TEXT)
    add_textbox(slide, Inches(11.00), Inches(7.22), Inches(2.00), Inches(0.26),
                f"{slide_no} / {total}", size=8.5, bold=True, color=GREEN,
                align=PP_ALIGN.RIGHT)


def bottom_strip(slide, label_text: str, summary_text: str, items: list[tuple[str, str, str]]):
    """The 1.20" 'BOTTOM LINE' strip pattern from the parent.

    items: list of (number, name, description). Typically 4 items.
    """
    add_rect(slide, Inches(0), Inches(5.95), SLIDE_W, Inches(1.20), fill=BOTTOM_NAVY)
    # Left label + summary
    add_textbox(slide, Inches(0.45), Inches(6.05), Inches(3.60), Inches(0.30),
                label_text, size=10, bold=True, color=GREEN)
    add_textbox(slide, Inches(0.45), Inches(6.31), Inches(3.60), Inches(0.80),
                summary_text, size=9.5, color=LIGHT_TEXT)
    # 4 columns
    x_starts = [Inches(4.30), Inches(6.53), Inches(8.77), Inches(11.00)]
    col_w = Inches(2.13)
    for (num, name, desc), x in zip(items, x_starts):
        add_textbox(slide, x, Inches(6.05), col_w, Inches(0.26),
                    num, size=9, bold=True, color=GREEN)
        add_textbox(slide, x, Inches(6.29), col_w, Inches(0.30),
                    name, size=10, bold=True, color=WHITE)
        add_textbox(slide, x, Inches(6.57), col_w, Inches(0.60),
                    desc, size=8.5, color=LIGHT_TEXT)


def card(slide, left, top, width, height, *, stripe_color, stripe_text, body_fn=None):
    """White card with a colored top stripe + label + optional body."""
    # White card
    add_rect(slide, left, top, width, height, fill=WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    # Top stripe
    add_rect(slide, left, top, width, Inches(0.32), fill=stripe_color)
    # Stripe label (bold white centered)
    add_textbox(slide, left + Inches(0.10), top + Inches(0.03),
                width - Inches(0.20), Inches(0.28),
                stripe_text, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if body_fn:
        body_fn(left, top + Inches(0.40), width, height - Inches(0.40))


def kv_item(slide, x, y, w, *, label, desc):
    """Card body item: bold label + muted description below."""
    add_textbox(slide, x, y, w, Inches(0.26), label,
                size=11, bold=True, color=STRONG_INK)
    add_textbox(slide, x, y + Inches(0.26), w, Inches(0.50), desc,
                size=9.5, color=MUTED)


# ─── Build the deck ──────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 15

    # ── Slide 1: Title ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.05), fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(1.05), fill=GREEN)
    add_textbox(s, Inches(0.45), Inches(0.12), Inches(12.50), Inches(0.50),
                "APEX  —  Walkthrough for DMTSP",
                size=22, bold=True, color=WHITE)
    add_textbox(s, Inches(0.45), Inches(0.62), Inches(12.50), Inches(0.40),
                "Deloitte's agentic-AI delivery framework — what it is · how to demo it · how to make it real",
                size=12, color=LIGHT_TEXT)

    # Hero
    add_rect(s, Inches(0.95), Inches(2.00), Inches(11.45), Inches(3.30),
             fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(s, Inches(0.95), Inches(2.00), Inches(0.18), Inches(3.30), fill=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.20), Inches(0.50), Inches(0.6),
                "“", size=72, bold=True, color=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.85), Inches(10.50), Inches(2.30),
                [
                    "APEX is Deloitte's agentic-AI delivery framework — the engagement layer cloud vendors don't ship.",
                    "",
                    "Foundry, Vertex, and Bedrock are agent runtimes. They are not engagement systems. APEX is what",
                    "turns an agent runtime into a billable, audit-defensible, repeatable Deloitte delivery — sold as",
                    "agentic services on top of the framework, scoped per Industry Pack.",
                ],
                size=18, color=STRONG_INK)

    # Bottom strip
    bottom_strip(s, "WHAT YOU'LL LEAVE WITH",
                 "Three concrete outcomes from this 30-minute walkthrough.",
                 [
                     ("01", "Vocabulary",     "Pack · VV · profile · LEDGER · envelope tier."),
                     ("02", "Demo path",      "Click sequence you can show any client in 10 minutes."),
                     ("03", "Sales motion",   "BVA → Pack Lite → Standard → Operate — every engagement."),
                     ("04", "Positioning",    "Why APEX is NOT a product — and what we actually sell."),
                 ])
    page_footer(s, 1, TOTAL)

    # ── Slide 2: APEX in one sentence ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "APEX  —  In One Sentence",
           "The simplest framing — start every client conversation here",
           "Memorize:",
           "“APEX is the engagement layer cloud vendors don't ship.”")

    # Hero quote card
    add_rect(s, Inches(0.95), Inches(2.00), Inches(11.45), Inches(3.30),
             fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(s, Inches(0.95), Inches(2.00), Inches(0.18), Inches(3.30), fill=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.20), Inches(0.50), Inches(0.6),
                "“", size=72, bold=True, color=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.85), Inches(10.50), Inches(2.30),
                [
                    "Foundry, Vertex, and Bedrock are excellent agent runtimes.",
                    "They are not engagement systems.",
                    "",
                    "APEX is what turns an agent runtime into a billable,",
                    "audit-defensible, repeatable Deloitte delivery.",
                ],
                size=20, color=STRONG_INK)

    bottom_strip(s, "BOTTOM LINE",
                 "Without APEX, an agent project competes with what a client can build alone from cloud-vendor docs. With APEX, Deloitte sells the engagement layer the cloud vendors do not.",
                 [
                     ("01", "Billable",     "Defined methodology + SOW templates."),
                     ("02", "Auditable",    "WORM LEDGER + hash-chain replay."),
                     ("03", "Repeatable",   "Industry packs ship the canonical playbook."),
                     ("04", "Sellable",     "Envelopes give sellers a defined motion."),
                 ])
    page_footer(s, 2, TOTAL)

    # ── Slide 3: The Problem — without vs with APEX (two-card pattern) ─
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "The Problem  —  An Agent Runtime Is Not An Engagement",
           "Cloud vendors ship the runtime. We sell what surrounds it.",
           "Why APEX:",
           "Foundry, Vertex, and Bedrock are excellent agent runtimes. They are not engagement systems. APEX is the layer that turns an agent runtime into a billable, audit-defensible, repeatable Deloitte delivery.")

    # Left card — WITHOUT (gray stripe)
    def left_body(x, y, w, h):
        items = [
            ("Agent runtime",           "Foundry · Vertex Agent Builder · Bedrock Agents — orchestration, sessions, tools"),
            ("LLM hosting + selection", "Azure OpenAI · Vertex Model Garden · Bedrock — model catalog and inference"),
            ("Storage primitives",      "OneLake · GCS · S3 — lakehouse foundations"),
            ("Identity primitives",     "Entra · Cloud Identity · IAM — directory, OBO, federation"),
            ("Catalog primitives",      "Purview · Dataplex · DataZone — lineage, classification"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.07 + i * 0.65), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(0.45), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GRAY_STRIP, stripe_text="WHAT CLOUD VENDORS GIVE YOU",
         body_fn=left_body)

    # Right card — WITH (green stripe)
    def right_body(x, y, w, h):
        items = [
            ("Methodology",            "Chain pattern · Wave 1/2/3 commercial arc · six discovery openers"),
            ("Canonical schemas",      "35+ industry schemas across MFGML · FINML · RISKML · MERML · CXML · OPSML · ESGML"),
            ("LEDGER audit substrate", "14-field WORM hash chain · cryptographically replayable decision attribution"),
            ("Industry verticalization", "Seven Industry Solution Packs · industry-tuned thresholds and resolution routes"),
            ("Commercial framework",   "5-tier Service Envelopes · BVA → DCIF → ISV → T&M → Operate funding paths"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.07 + i * 0.65), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(6.85), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GREEN, stripe_text="WHAT APEX ADDS  (the Margin Spine)",
         body_fn=right_body)

    bottom_strip(s, "BOTTOM LINE",
                 "Without APEX, an agent project looks like staff augmentation. With APEX, Deloitte sells the engagement layer cloud vendors do not — and the deal shape changes accordingly.",
                 [
                     ("01", "Billable",   "Defined methodology + SOW templates."),
                     ("02", "Auditable",  "WORM ledger + hash-chain replay."),
                     ("03", "Repeatable", "Industry packs ship the canonical playbook."),
                     ("04", "Sellable",   "Envelopes give sellers a defined motion."),
                 ])
    page_footer(s, 3, TOTAL)

    # ── Slide 4: The Five Pieces ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "The Five Pieces of APEX",
           "Core  •  Profiles  •  Packs  •  Envelopes  •  LEDGER  —  every engagement assembles from these",
           "Position:",
           "APEX is Deloitte's agentic-AI delivery framework: a Cloud-Neutral Core running on top of cloud-native agent runtimes (Foundry · Vertex · Bedrock) plus a Margin Spine — Virtual Views, Industry Solution Packs, Service Envelopes, and a WORM LEDGER — that makes agent work sellable, auditable, and repeatable.")

    # Quote card (mid)
    add_rect(s, Inches(0.95), Inches(2.00), Inches(11.45), Inches(3.30),
             fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(s, Inches(0.95), Inches(2.00), Inches(0.18), Inches(3.30), fill=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.20), Inches(0.50), Inches(0.6),
                "“", size=72, bold=True, color=GREEN)
    add_textbox(s, Inches(1.45), Inches(2.85), Inches(10.50), Inches(2.30),
                [
                    "APEX is Deloitte's agentic-AI delivery framework: a Cloud-Neutral Core that runs on top of",
                    "cloud-native agent runtimes (Foundry · Vertex · Bedrock) and adds a Margin Spine — Virtual",
                    "Views, Industry Solution Packs, Service Envelopes, and a WORM LEDGER — that makes agent",
                    "work sellable, auditable, and repeatable.",
                ],
                size=17, color=STRONG_INK)

    bottom_strip(s, "THE FIVE PIECES",
                 "Each piece is a section of this deck.",
                 [
                     ("01", "Cloud-Neutral Core", "VV runtime · LEDGER · 6-Agent Fleet · MCP contracts."),
                     ("02", "Cloud Profiles",     "APEX-M · APEX-G · APEX-A · 14-interface contract."),
                     ("03", "Industry Packs",     "10-asset bundle · Manufacturing v1 · 6 more in catalog."),
                     ("04", "Service Envelopes",  "5 tiers Deloitte sells · 3 pack sub-tiers."),
                 ])
    page_footer(s, 4, TOTAL)

    # ── Slide 5: The APEX family — M / G / A ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "The APEX Family  —  M  /  G  /  A",
           "Three cloud-native variants  •  one shared Cloud-Neutral Core  •  one shared pack catalog",
           "Position:",
           "APEX-M, APEX-G, APEX-A are not three products. They are three implementations of the same Cloud Profile contract. The Core (VV runtime, LEDGER, agent fleet) is cloud-neutral. Industry Packs are cloud-agnostic. Only the 14-interface profile changes per cloud.")

    families = [
        ("APEX-M", MS_BLUE, "Microsoft  /  Azure", "GA today",
         [("Foundry + Azure OpenAI",  ""),
          ("OneLake (Delta) + Fabric SQL", ""),
          ("Entra ID · Key Vault · Purview", ""),
          ("Teams Adaptive Cards (HITL)",  ""),
          ("Azure Maps Creator (CFMP iface #15)", "")]),
        ("APEX-G", RGBColor(0x0F, 0x9D, 0x58), "Google  /  GCP", "Beta Q3 FY27",
         [("Vertex AI + Gemini", ""),
          ("BigQuery + BigLake", ""),
          ("Cloud Identity · Secret Manager · Dataplex", ""),
          ("Google Chat cards / GWS", ""),
          ("Google Maps Geospatial Creator", "")]),
        ("APEX-A", RGBColor(0xFF, 0x99, 0x00), "Amazon  /  AWS", "Beta Q4 FY27",
         [("Bedrock (Claude / Nova)", ""),
          ("S3 + Lake Formation", ""),
          ("IAM Identity Center · Secrets Manager", ""),
          ("Slack Connect / Workmail", ""),
          ("Amazon Location Service indoor maps", "")]),
    ]
    x0 = Inches(0.45); gap = Inches(0.10); cw = Inches((13.33 - 0.45 * 2 - 0.10 * 2) / 3)
    for i, (tag, color, name, status, items) in enumerate(families):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        def make_body(items_local, name_local, status_local):
            def body(x, y, w, h):
                add_textbox(s, x + Inches(0.10), y + Inches(0.05),
                            w - Inches(0.20), Inches(0.32),
                            name_local, size=11, color=RGBColor(0x64, 0x74, 0x8B))
                add_rect(s, x + Inches(0.10), y + Inches(0.42),
                         w - Inches(0.20), Inches(0.30), fill=color)
                add_textbox(s, x + Inches(0.10), y + Inches(0.42),
                            w - Inches(0.20), Inches(0.30),
                            status_local, size=9.5, bold=True, color=WHITE,
                            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                for i_, (lbl, desc) in enumerate(items_local):
                    add_textbox(s, x + Inches(0.10), y + Inches(0.85 + i_ * 0.50),
                                w - Inches(0.20), Inches(0.42),
                                "•  " + lbl, size=10.5, bold=True, color=STRONG_INK)
            return body
        card(s, cx, Inches(2.05), cw, Inches(3.85),
             stripe_color=color, stripe_text=tag,
             body_fn=make_body(items, name, status))

    bottom_strip(s, "ONE FAMILY",
                 "Same Core. Same packs. Same envelopes. Different cloud profile. Replatform = change `profile:` in cloud-profile.yaml.",
                 [
                     ("01", "Shared Core",      "VV runtime · LEDGER · agent fleet · MCP."),
                     ("02", "Shared Packs",     "CFMP & every other Pack runs identically on M/G/A."),
                     ("03", "Cloud-specific",   "14 interfaces · storage · LLM · HITL · identity · maps."),
                     ("04", "Marketplace",      "Azure · GCP · AWS marketplaces for burndown."),
                 ])
    page_footer(s, 5, TOTAL)

    # ── Slide 6: Industry Packs catalog ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "Industry Solution Packs  —  The Catalog (7 today)",
           "Same 10-asset structure every Pack — that's why new Packs land predictably",
           "Catalog:",
           "Each Pack ships the same 10 assets (VVs · scenarios · adapters · cards · personas · demo-data · BVA · SOW · tests · runbook). Only the contents change per industry. CFMP is the 7th and the first organized around a customer-moment spine.")

    packs = [
        ("01", "Manufacturing v1",          "Plant Manager · Ops Director",       "MFGML",                                "GA"),
        ("02", "Retail Merchandising v1",   "Chief Merchant · DMM · Planning Lead", "MERML",                              "v1 Q2 FY27"),
        ("03", "Finance / Controllership",  "Controller · CFO",                    "FINML",                              "v1 Q4 FY27"),
        ("04", "Risk & Compliance",         "SOX PMO Lead · CRO",                  "RISKML",                             "v1 Q4 FY27"),
        ("05", "Customer Experience",       "Care Ops Lead",                       "CXML",                               "v1 Q4 FY27"),
        ("06", "ESG",                       "Sustainability Officer",              "ESGML",                              "v1 Q3 FY28"),
        ("07", "CFMP · Customer Focused Merchandise", "CMO · CX-VP · Loyalty Director",
                                                                                   "CXML + MERML + new cfmp namespace",  "Phase 1 LIVE"),
    ]
    # Table positioned at 0.45,1.85 (below tagline)
    y_top = Inches(1.95)
    rowh = Inches(0.42)
    col_x = [Inches(0.45), Inches(1.05), Inches(5.20), Inches(8.90), Inches(11.40)]
    col_w = [Inches(0.55), Inches(4.10), Inches(3.65), Inches(2.45), Inches(1.50)]

    # Header row (dark navy with white text)
    for x, w, txt in zip(col_x, col_w, ["#", "Pack", "Buyer persona", "Schema", "Status"]):
        add_rect(s, x, y_top, w, rowh, fill=HEAD_NAVY)
        add_textbox(s, x + Inches(0.08), y_top, w - Inches(0.16), rowh,
                    txt, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    for i, (num, name, persona, schema, status) in enumerate(packs):
        ry = y_top + Inches((i + 1) * 0.42)
        is_cfmp = num == "07"
        fill = ROW_HILIGHT if is_cfmp else WHITE
        line = GREEN if is_cfmp else RGBColor(0xE5, 0xE7, 0xEB)
        # Row background
        for x, w in zip(col_x, col_w):
            add_rect(s, x, ry, w, rowh, fill=fill, line=line)
        # Cells
        add_textbox(s, col_x[0] + Inches(0.08), ry, col_w[0] - Inches(0.16), rowh,
                    num, size=10, bold=is_cfmp, color=STRONG_INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, col_x[1] + Inches(0.08), ry, col_w[1] - Inches(0.16), rowh,
                    name, size=10.5, bold=is_cfmp, color=STRONG_INK, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, col_x[2] + Inches(0.08), ry, col_w[2] - Inches(0.16), rowh,
                    persona, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, col_x[3] + Inches(0.08), ry, col_w[3] - Inches(0.16), rowh,
                    schema, size=9.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        status_color = GREEN if status in ("GA", "Phase 1 LIVE") else GRAY_STRIP
        add_rect(s, col_x[4] + Inches(0.10), ry + Inches(0.08),
                 col_w[4] - Inches(0.20), Inches(0.26), fill=status_color)
        add_textbox(s, col_x[4] + Inches(0.10), ry + Inches(0.08),
                    col_w[4] - Inches(0.20), Inches(0.26),
                    status, size=9, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bottom_strip(s, "WHY THIS COMPOUNDS",
                 "Each new Pack is identical engineering work — populate the same 10 slots. Pack catalog growth doesn't scale the platform team.",
                 [
                     ("01", "Repeatable",   "10-slot bundle every time."),
                     ("02", "Predictable",  "Engineering time per Pack is bounded."),
                     ("03", "Composable",   "Cross-pack scenarios via imports."),
                     ("04", "Sellable",     "Each Pack maps to industry-led conversations."),
                 ])
    page_footer(s, 6, TOTAL)

    # ── Slide 7: The 10-Asset Bundle ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "Industry Solution Packs  —  10-Asset Bundle",
           "Virtual Views + Scenarios bundled per industry  •  identical structure across every Pack",
           "What you buy:",
           "Industry Packs are the productization unit. Every Pack ships the same 10 assets — only the contents change. Same engineering structure means new Packs land in predictable time.")

    # 5×2 grid of mini cards
    assets = [
        ("1", "VV Manifests",       "views/",              "Declarative federated views with thresholds + resolution routes."),
        ("2", "Scenario Manifests", "scenarios/",          "End-to-end business outcomes composed from VVs."),
        ("3", "Source Adapters",    "adapters/",           "Connectors to common source systems for the industry."),
        ("4", "Adaptive Cards",     "cards/",              "JSON cards per persona role and breach band."),
        ("5", "Persona Map",        "personas.yaml",       "Role-to-identity binding + SLAs + escalation chains."),
        ("6", "Demo Data",          "demo-data/",          "Synthetic data so the Pack runs on a laptop."),
        ("7", "BVA Worksheet",      "bva/<ind>-roi.xlsx",  "Industry-tuned ROI calculator for BVA workshops."),
        ("8", "Sample SOW",         "envelopes/",          "Lite / Standard / Enterprise SOW templates."),
        ("9", "Acceptance Tests",   "tests/",              "60+ tests proving Pack health in CI."),
        ("10", "Runbook + Training", "runbooks/",          "Operate procedures + analyst training materials."),
    ]
    x0 = Inches(0.45); y0 = Inches(2.00); gap = Inches(0.10)
    cw = Inches(2.43); ch = Inches(1.85)
    for i, (n, title, path, desc) in enumerate(assets):
        col = i % 5; row = i // 5
        cx = x0 + Inches(col * (cw / Inches(1) + gap / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, cy, cw, ch, fill=WHITE,
                 line=RGBColor(0xE5, 0xE7, 0xEB))
        # Number stripe
        add_rect(s, cx, cy, cw, Inches(0.30), fill=GREEN)
        add_textbox(s, cx + Inches(0.10), cy + Inches(0.02),
                    cw - Inches(0.20), Inches(0.26),
                    f"{n}   {title}", size=10.5, bold=True, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Path
        add_textbox(s, cx + Inches(0.10), cy + Inches(0.36),
                    cw - Inches(0.20), Inches(0.26), path,
                    size=9.5, bold=True, color=MS_BLUE)
        # Desc
        add_textbox(s, cx + Inches(0.10), cy + Inches(0.65),
                    cw - Inches(0.20), ch - Inches(0.75), desc,
                    size=9, color=MUTED)

    bottom_strip(s, "BOTTOM LINE",
                 "Same engineering structure every Pack means new Packs land in predictable time — Pack catalog growth doesn't scale the platform team.",
                 [
                     ("01", "Repeatable",   "10-slot bundle every time."),
                     ("02", "Predictable",  "Engineering time per Pack is bounded."),
                     ("03", "Composable",   "Cross-pack scenarios via imports."),
                     ("04", "Sellable",     "Each Pack maps to industry-led conversations."),
                 ])
    page_footer(s, 7, TOTAL)

    # ── Slide 8: Spotlight CFMP ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "Spotlight  —  CFMP, the Live Worked Example",
           "Customer Focused Merchandise Pack v0.2 · Phase 1 LIVE on APEX-M dev substrate",
           "Why CFMP:",
           "First Industry Pack organized around a customer-moment spine — sibling to (not a replacement for) the operator-side Retail Merchandising Pack v1. Buyer: CMO / CX-VP / Loyalty Director.")

    # 4 phase columns
    phases = [
        ("CHOOSE",   "What should I buy?",
         ["Recipes · pairings",
          "Dietary filters",
          "Personalized offers",
          "In-store ad attribution",
          "Sampling-table engagement"], MS_BLUE),
        ("SELECT",   "Where is it?",
         ["Wayfinding (iface #15)",
          "OSA / shelf-gap dispatch",
          "Planogram compliance",
          "Aisle engagement attribution",
          "End-cap ROI tracking"], GREEN),
        ("BUY",      "How do I pay?",
         ["Scan-and-go",
          "Self-checkout assist",
          "Queue prediction",
          "BOPIS pickup queue",
          "Cart-abandon recovery"], RGBColor(0xFF, 0x99, 0x00)),
        ("SERVICES", "What's next from you?",
         ["Loyalty winback ⭐",
          "Loyalty tier migration",
          "Complaint triage",
          "Review summary",
          "Sentiment + NPS"], RGBColor(0xB4, 0x42, 0xCC)),
    ]
    cw2 = Inches(3.05); gap2 = Inches(0.10); x0 = Inches(0.45)
    for i, (label, question, scenarios, color) in enumerate(phases):
        cx = x0 + Inches(i * (cw2 / Inches(1) + gap2 / Inches(1)))
        # Card
        add_rect(s, cx, Inches(2.05), cw2, Inches(3.85), fill=WHITE,
                 line=RGBColor(0xE5, 0xE7, 0xEB))
        # Phase stripe
        add_rect(s, cx, Inches(2.05), cw2, Inches(0.32), fill=color)
        add_textbox(s, cx + Inches(0.10), Inches(2.08), cw2 - Inches(0.20), Inches(0.28),
                    label, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # Question
        add_textbox(s, cx + Inches(0.15), Inches(2.50), cw2 - Inches(0.30), Inches(0.34),
                    f"“{question}”", size=11.5, italic=True, color=MUTED)
        # Scenarios list
        for j, sc in enumerate(scenarios):
            add_textbox(s, cx + Inches(0.20), Inches(2.95 + j * 0.45),
                        cw2 - Inches(0.40), Inches(0.40),
                        "•  " + sc, size=10.5, color=STRONG_INK)

    bottom_strip(s, "WHAT'S LIVE TODAY",
                 "Phase 1 of CFMP is live on APEX-M dev substrate. Demo URL: ca-visionkit-portal…/architecture.",
                 [
                     ("01", "18 scenarios",        "Mapped to the 4 customer-moment phases."),
                     ("02", "6 new VVs",           "merml.* + cxml.* + cfmp.storemap_view."),
                     ("03", "Iface #15 (CFMP)",   "Azure Maps Creator for indoor wayfinding."),
                     ("04", "Lite engagement",    "$150K–$250K · 4–6 weeks · BVA + DCIF."),
                 ])
    page_footer(s, 8, TOTAL)

    # ── Slide 9: Demo walk-through (URL + click path) ──────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "How to Drive the Live Demo",
           "Open in front of the client — talk through each surface · 10 minutes end-to-end",
           "URL:",
           "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture  ←  open this for the architecture overview · root URL for the customer-facing CFMP demo")

    steps_left = [
        ("1", "Start at /",                  "“This is what the customer sees at the sampling table or kiosk.” Note the 4 phase chips: CHOOSE · SELECT · BUY · SERVICES."),
        ("2", "Type: Where can I find Coca-Cola?", "Watch the Wayfinding route card render — zone, aisle, distance, turn-by-turn. This rides proposed Interface #15 (Azure Maps Creator)."),
        ("3", "Toggle a dietary chip",       "(e.g. dairy-free) Then hold a dairy product up. Proactive associate fires a red Dietary alert + suggests a compliant swap."),
        ("4", "Add to cart, click a cart line", "Cart is client-authoritative; agent emits cart_actions in its reply. Click a cart item → agent describes it in chat."),
    ]
    steps_right = [
        ("5", "Open /architecture",          "12-section page. Hit “expand all” on the dependency tree (Section 3) to walk customer→portal→orchestrator→Azure→devices."),
        ("6", "Click Details ⓘ on any Azure row", "Modal opens with host env, configuration, cost estimate, consumption metrics. Proves it's real Azure, not a mockup."),
        ("7", "Show Section 11 (APEX accelerator)", "12,649 lines of APEX framework code · 341 lines of CFMP glue · 509 passing tests. That's why this shipped in weeks."),
        ("8", "Toggle 🎓 APEX tour (flyout)", "Right-side flyout cycles through 10 educational topics about Packs, VVs, CFMP — built-in conversation starter."),
    ]

    def render_steps(steps, x0):
        for i, (n, title, body) in enumerate(steps):
            ry = Inches(1.95 + i * 0.97)
            # Circle
            circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x0, ry, Inches(0.42), Inches(0.42))
            circ.fill.solid(); circ.fill.fore_color.rgb = GREEN
            circ.line.fill.background()
            tfc = circ.text_frame; tfc.clear()
            tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfc.margin_left = Emu(0); tfc.margin_right = Emu(0)
            pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
            rc = pc.add_run(); rc.text = n
            rc.font.name = "Calibri"; rc.font.size = Pt(13); rc.font.bold = True
            rc.font.color.rgb = WHITE
            # Title
            add_textbox(s, x0 + Inches(0.55), ry - Inches(0.02),
                        Inches(5.8), Inches(0.35),
                        title, size=11.5, bold=True, color=STRONG_INK)
            add_textbox(s, x0 + Inches(0.55), ry + Inches(0.30),
                        Inches(5.8), Inches(0.70),
                        body, size=9.5, color=MUTED)
    render_steps(steps_left,  Inches(0.45))
    render_steps(steps_right, Inches(6.85))

    bottom_strip(s, "PACING",
                 "Eight steps × about 60 seconds each = a full walkthrough in 8–10 minutes. Keep the architecture page open in a second tab from the start.",
                 [
                     ("01", "Steps 1-4",   "Customer chat surface · ~4 minutes."),
                     ("02", "Steps 5-7",   "/architecture deep-dive · ~4 minutes."),
                     ("03", "Step 8",      "APEX tour flyout · ~2 minutes."),
                     ("04", "Then close",  "Recap with the BVA → Pack Lite path (slide 12)."),
                 ])
    page_footer(s, 9, TOTAL)

    # ── Slide 10: Engagement ladder ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "The Service Envelope Stack  —  5 Tiers Deloitte Sells",
           "Stackable engagement tiers with independence-safe funding paths",
           "Position:",
           "Every APEX engagement assembles from five stackable tiers. Foundation is sold once per tenant. Industry Packs are sold per industry in three sub-tiers. Scenarios and Custom builds layer on top. Operate is the run-rate.")

    tiers = [
        ("T1", "FOUNDATION", "APEX Platform Foundation",
         "VV runtime · federator · MCP auto-reg · LEDGER · Entra/Purview wiring",
         "Sold once per client tenant",
         "Fixed fee · 6–10 wks · $400–700K · BVA + DCIF"),
        ("T2", "INDUSTRY PACK", "Pack Lite / Standard / Enterprise",
         "10-asset Pack bundle — manifests, adapters, cards, BVA, SOW",
         "Sold per industry · 3 sub-tiers",
         "Tiered · 4 wks–9 mo · $100K–$3.5M · DCIF · T&M · ISV"),
        ("T3", "SCENARIO", "Single Scenario Chain Activation",
         "One named outcome (e.g. SOX PMO testing) live on client's sources",
         "Sold per business outcome",
         "Fixed fee · 4–8 wks · $150–400K · DCIF · T&M · ISV"),
        ("T4", "CUSTOM BUILD", "New VV / new Scenario  (bespoke)",
         "Client-specific manifest added to their Pack",
         "Sold as advisory + dev",
         "T&M · Sprints · Variable"),
        ("T5", "OPERATE", "APEX Operate Managed Service",
         "24×7 monitoring · threshold tuning · Pack version uptake",
         "Run-rate / steady state",
         "Subscription · Continuous · $/mo · Client direct"),
    ]
    y_top = Inches(1.95)
    rowh = Inches(0.78)
    for i, (tag, kind, name, scope, sale, price) in enumerate(tiers):
        ry = y_top + Inches(i * 0.78)
        # Tag chip
        add_rect(s, Inches(0.45), ry + Inches(0.04),
                 Inches(0.70), Inches(0.70), fill=HEAD_NAVY)
        add_textbox(s, Inches(0.45), ry + Inches(0.04),
                    Inches(0.70), Inches(0.70),
                    tag, size=16, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Kind chip (under tag)
        # Name + scope
        add_textbox(s, Inches(1.25), ry + Inches(0.02),
                    Inches(5.8), Inches(0.30),
                    name, size=12, bold=True, color=STRONG_INK)
        add_textbox(s, Inches(1.25), ry + Inches(0.30),
                    Inches(5.8), Inches(0.50),
                    scope, size=9.5, color=MUTED)
        # Sale model chip
        add_rect(s, Inches(7.30), ry + Inches(0.16), Inches(2.45), Inches(0.50),
                 fill=WHITE, line=RGBColor(0xCB, 0xD5, 0xE1))
        add_textbox(s, Inches(7.30), ry + Inches(0.16),
                    Inches(2.45), Inches(0.50),
                    sale, size=9.5, color=STRONG_INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Price chip
        add_rect(s, Inches(9.95), ry + Inches(0.16), Inches(2.90), Inches(0.50),
                 fill=GREEN)
        add_textbox(s, Inches(9.95), ry + Inches(0.16),
                    Inches(2.90), Inches(0.50),
                    price, size=9.5, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bottom_strip(s, "WHY THIS SELLS",
                 "Five compartmentalized envelopes with defined SOW templates, BVA worksheets, and acceptance criteria let an LCSP price a new conversation in minutes.",
                 [
                     ("01", "Enter anywhere",   "BVA-led Pack Lite is the standard wedge."),
                     ("02", "Grow downward",    "Foundation underpins all higher tiers."),
                     ("03", "Grow outward",     "Add packs and scenarios per industry."),
                     ("04", "Run-rate",         "Operate tier compounds into managed-service."),
                 ])
    page_footer(s, 10, TOTAL)

    # ── Slide 11: Independence (DO / DO NOT two-card pattern) ─────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "Funding & Independence",
           "Independence-safe paths for every envelope  •  Deloitte never receives ECIF directly",
           "Rule:",
           "Deloitte is Microsoft's auditor. Per SEC rules, Deloitte cannot receive ECIF directly. Workarounds via ISV Marketplace + SI Teaming keep Microsoft money flowing without compromising independence. APEX-G / APEX-A have no such constraint.")

    def left_do(x, y, w, h):
        items = [
            ("BVA",                "Client-funded discovery workshop"),
            ("DCIF",               "Deloitte co-invest fund"),
            ("ISV Marketplace",    "Microsoft money flows via the Marketplace burndown"),
            ("SI Teaming POC",     "Non-competing SI runs the POC"),
            ("T&M",                "Custom Build (T4) and Operate (T5)"),
            ("Client direct",      "Operate run-rate billing"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.10 + i * 0.55), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(0.45), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GREEN, stripe_text="DO  (independence-safe funding)",
         body_fn=left_do)

    def right_dont(x, y, w, h):
        items = [
            ("Route ECIF through Deloitte",  "Direct ECIF receipt breaks independence"),
            ("Imply joint commercial entity",  "“Partnership” / “alliance” language is OUT"),
            ("Describe APEX as a Microsoft product", "APEX is Deloitte IP, not Microsoft's"),
            ("Take Microsoft co-sell credit",     "Engagement is Deloitte's, not joint"),
            ("Skip the language discipline",   "Always: “Deloitte's Microsoft practice” / DMTSP"),
            ("Sell APEX itself",                  "We sell agentic delivery services on the framework"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.10 + i * 0.55), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(6.85), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GRAY_STRIP, stripe_text="DO NOT  (independence violations)",
         body_fn=right_dont)

    bottom_strip(s, "CROSS-CLOUD FUNDING",
                 "APEX-G uses Google Industry Accelerator (GIA). APEX-A uses AWS Migration Acceleration Program (MAP). No independence wrinkle on G/A.",
                 [
                     ("01", "APEX-M",  "ISV/SI burndown · BVA · DCIF."),
                     ("02", "APEX-G",  "GIA + GCP Marketplace burndown."),
                     ("03", "APEX-A",  "AWS MAP + AWS Marketplace burndown."),
                     ("04", "All",     "BVA-led discovery is the universal first step."),
                 ])
    page_footer(s, 11, TOTAL)

    # ── Slide 12: What we DO and DO NOT sell ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "What We DO and DO NOT Sell",
           "The single most important framing — repeat it every meeting",
           "Rule:",
           "APEX is Deloitte IP. Industry Packs (including CFMP) are scoped delivery engagements sold per industry as Deloitte agentic services. The framework itself is never licensed out.")

    # NOT card
    def not_body(x, y, w, h):
        items = [
            ("Not a SKU",          "Not a license, not a subscription a client can buy"),
            ("Not co-branded",     "Not Microsoft / Google / Amazon co-sell"),
            ("Not on Marketplace", "Pack-specific ISV listings exist for burndown — those are different from APEX itself"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.15 + i * 0.85), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(0.45), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GRAY_STRIP, stripe_text="❌  APEX is NOT a product",
         body_fn=not_body)

    # IS card
    def is_body(x, y, w, h):
        items = [
            ("Delivery framework",  "An engagement layer that turns cloud-vendor runtimes into Deloitte deliveries"),
            ("Agentic services",    "What we sell is Deloitte agentic-AI delivery services on the APEX framework"),
            ("Scoped per Pack",     "Engagement signs an SOW for a specific Pack at a specific sub-tier (e.g. CFMP Lite)"),
            ("Deloitte IP",         "Framework code (Core, Profiles, Packs) is Deloitte IP — never licensed out standalone"),
        ]
        for i, (lbl, desc) in enumerate(items):
            kv_item(s, x + Inches(0.20), y + Inches(0.10 + i * 0.70), w - Inches(0.40),
                    label=lbl, desc=desc)
    card(s, Inches(6.85), Inches(2.05), Inches(6.20), Inches(3.85),
         stripe_color=GREEN, stripe_text="✅  APEX IS Deloitte's delivery framework",
         body_fn=is_body)

    bottom_strip(s, "THE LINE",
                 "When clients ask “Can I have APEX?” the answer is: “You can engage Deloitte to deliver an APEX-based solution. Here's what that looks like…”",
                 [
                     ("01", "Buyer",    "CMO / CX-VP (CFMP) · Ops VP (Manufacturing) · CRO (Risk) · etc."),
                     ("02", "Vehicle",  "SOW under the right Service Envelope tier."),
                     ("03", "Funding",  "BVA · DCIF · ISV burndown · T&M · Client direct."),
                     ("04", "Outcome",  "Pack Lite live in 4–6 wks · Standard in 12–16 · Operate run-rate."),
                 ])
    page_footer(s, 12, TOTAL)

    # ── Slide 13: BVA → Pack Lite wedge ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "The Standard Wedge  —  BVA → Pack Lite",
           "How every CFMP-style engagement actually starts",
           "Motion:",
           "BVA earns the pilot. The pilot earns the rollout. The rollout earns Operate. Each step is additive — no client ever pays to redo earlier work.")

    waves = [
        ("BVA",         "Discovery workshop",        "4 hours",     "Client-funded",
         "Scenario shortlist · stakeholder + persona map · pilot SOW skeleton ready"),
        ("Pack Lite",   "Pilot — 3 scenarios + wayfinding", "4–6 wks", "$150–250K · BVA + DCIF",
         "1 sub-domain of VVs live · LEDGER + replay proven on client data"),
        ("Pack Standard", "First production rollout",      "12–16 wks", "$500K–1.5M · DCIF + ISV",
         "All 6 agents wired · Operate-readiness assessment · production cycle"),
        ("Pack Enterprise", "Full chain + Operate-ready",   "6–9 mo",   "$1.5–3.5M · Client direct",
         "Complete VV catalog · all scenarios · DataOps embedded"),
        ("Operate",     "Run-rate managed service",  "Continuous",  "$/mo · Client direct",
         "24×7 monitoring · threshold tuning · Pack version uptake"),
    ]
    yC = Inches(1.95); cw3 = Inches(2.43); gap3 = Inches(0.10); x0 = Inches(0.45)
    for i, (label, role, time, fund, outcome) in enumerate(waves):
        cx = x0 + Inches(i * (cw3 / Inches(1) + gap3 / Inches(1)))
        # Step circle above the card
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx + Inches(0.95), yC - Inches(0.05),
                                   Inches(0.50), Inches(0.50))
        circ.fill.solid(); circ.fill.fore_color.rgb = HEAD_NAVY
        circ.line.fill.background()
        tfc = circ.text_frame; tfc.clear()
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc.margin_left = Emu(0); tfc.margin_right = Emu(0)
        pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run(); rc.text = str(i + 1)
        rc.font.name = "Calibri"; rc.font.size = Pt(15); rc.font.bold = True
        rc.font.color.rgb = GREEN
        # Card
        add_rect(s, cx, yC + Inches(0.55), cw3, Inches(3.40), fill=WHITE,
                 line=RGBColor(0xE5, 0xE7, 0xEB))
        # Phase stripe
        add_rect(s, cx, yC + Inches(0.55), cw3, Inches(0.32), fill=GREEN)
        add_textbox(s, cx + Inches(0.10), yC + Inches(0.58),
                    cw3 - Inches(0.20), Inches(0.28),
                    label, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # Role
        add_textbox(s, cx + Inches(0.15), yC + Inches(0.95),
                    cw3 - Inches(0.30), Inches(0.35),
                    role, size=11, bold=True, color=STRONG_INK)
        # Time
        add_textbox(s, cx + Inches(0.15), yC + Inches(1.30),
                    cw3 - Inches(0.30), Inches(0.28),
                    f"⏱  {time}", size=9.5, color=MUTED)
        # Funding chip
        add_rect(s, cx + Inches(0.15), yC + Inches(1.60),
                 cw3 - Inches(0.30), Inches(0.40), fill=HEAD_NAVY)
        add_textbox(s, cx + Inches(0.15), yC + Inches(1.60),
                    cw3 - Inches(0.30), Inches(0.40),
                    fund, size=9, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Outcome
        add_textbox(s, cx + Inches(0.15), yC + Inches(2.10),
                    cw3 - Inches(0.30), Inches(1.80),
                    outcome, size=9, color=MUTED)

    bottom_strip(s, "ADDITIVE BY DESIGN",
                 "Lite → Standard → Enterprise → Operate. No client ever pays to redo earlier work. The ladder IS the natural progression.",
                 [
                     ("01", "Lite as wedge",      "BVA earns the Lite engagement."),
                     ("02", "Standard as proof",  "First production cycle on client data."),
                     ("03", "Enterprise as run",  "All views, all scenarios, ready for Operate."),
                     ("04", "No rework",          "Each tier extends, never replaces, prior."),
                 ])
    page_footer(s, 13, TOTAL)

    # ── Slide 14: Q&A ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "Client Questions You Will Hear  —  Answers Ready",
           "Memorize these for the first 10 minutes of every conversation",
           "Pattern:",
           "These six questions cover ~90% of the opening objections. Answer crisp, then redirect to the demo URL.")

    qa = [
        ("Is this Microsoft's product?",
         "No — APEX is Deloitte's delivery framework. We use Foundry as the underlying runtime, but the engagement layer on top is Deloitte IP."),
        ("Can we buy APEX?",
         "You engage Deloitte for an APEX-based delivery. We scope an Industry Pack at the right sub-tier (Lite / Standard / Enterprise) and sign an SOW. The framework itself isn't licensed."),
        ("How is this different from what I'd build on Foundry?",
         "Foundry ships the agent runtime. We ship the engagement layer: 10-asset Industry Packs, WORM LEDGER, reasoning-trace + replay, persona/HITL plumbing, three cloud profiles. 12,649 lines of vetted Deloitte IP."),
        ("What if we're on AWS / Google?",
         "Same Packs, same Core, different profile. APEX-G uses Vertex; APEX-A uses Bedrock. Replatforming a Pack is a config change — manifests, agent code, LEDGER schema all unchanged."),
        ("How fast can we see something working?",
         "BVA workshop in week one. Pack Lite live in 4–6 weeks with one sub-domain of Virtual Views on your data. The demo runs that exact path."),
        ("What does the audit story actually look like?",
         "Every decision an agent makes lands a 14-field hash-chained LedgerRow. Every step has a reasoning-trace. Any decision can be replayed against the original Bronze snapshot — byte-identical. No cloud vendor ships this."),
    ]
    y_top = Inches(1.95)
    rowh = Inches(0.65)
    for i, (q, a) in enumerate(qa):
        ry = y_top + Inches(i * 0.65)
        # Row
        add_rect(s, Inches(0.45), ry, Inches(12.45), rowh, fill=WHITE,
                 line=RGBColor(0xE5, 0xE7, 0xEB))
        # Q
        add_textbox(s, Inches(0.60), ry + Inches(0.06),
                    Inches(5.6), Inches(0.55),
                    "Q   " + q, size=11, bold=True, color=STRONG_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        # A
        add_textbox(s, Inches(6.30), ry + Inches(0.06),
                    Inches(6.45), Inches(0.55),
                    "A   " + a, size=9.5, color=MUTED,
                    anchor=MSO_ANCHOR.MIDDLE)

    bottom_strip(s, "WHEN STUCK",
                 "If you can't answer in 30 seconds, redirect: “Let me show you in the live demo — it's faster.”",
                 [
                     ("01", "Lead with story",  "Show the wayfinding card. Concrete > abstract."),
                     ("02", "Then show audit",  "Click Details ⓘ + APEX panel — proves it's real."),
                     ("03", "Then talk price",  "BVA is free for the client. Pack Lite starts $150K."),
                     ("04", "Close on time",    "Pack Lite in 4–6 wks · Standard in 12–16 wks."),
                 ])
    page_footer(s, 14, TOTAL)

    # ── Slide 15: What to do this week ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header(s, "What to Do This Week",
           "Three concrete next steps for every DMTSP seller",
           "Call to action:",
           "The demo is live and ready for client conversations. The BVA template is in docs/packs/. The Pack Lite SOW skeleton ships with each Pack.")

    steps = [
        ("01", "Open the live demo with one client this week",
         "Use the architecture URL above. Walk the 8-step path on slide 9.",
         "60 minutes · low effort · one meeting"),
        ("02", "Stand up a BVA workshop",
         "Pick the Pack (CFMP for retail customer-side, Manufacturing for OEE, Risk for SOX, etc.) and book 4 hours. BVA is client-funded.",
         "Week 1–2 · zero Deloitte cost"),
        ("03", "Convert into a Pack Lite SOW",
         "BVA → scenario shortlist → 1 demo VV on client data → Pack Lite ($150K-$250K · 4–6 wks · BVA + DCIF). That's the wedge.",
         "Week 3–4 · DCIF approval · SOW signed"),
    ]
    yS = Inches(2.05)
    for i, (n, t, body, foot) in enumerate(steps):
        ry = yS + Inches(i * 1.25)
        # Card
        add_rect(s, Inches(0.45), ry, Inches(12.45), Inches(1.15), fill=WHITE,
                 line=RGBColor(0xE5, 0xE7, 0xEB))
        # Number
        add_rect(s, Inches(0.45), ry, Inches(1.05), Inches(1.15), fill=HEAD_NAVY)
        add_textbox(s, Inches(0.45), ry, Inches(1.05), Inches(1.15),
                    n, size=32, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_textbox(s, Inches(1.65), ry + Inches(0.10),
                    Inches(8.4), Inches(0.40),
                    t, size=14, bold=True, color=STRONG_INK)
        add_textbox(s, Inches(1.65), ry + Inches(0.50),
                    Inches(8.4), Inches(0.60),
                    body, size=10, color=MUTED)
        # Foot chip
        add_rect(s, Inches(10.20), ry + Inches(0.30), Inches(2.55), Inches(0.55),
                 fill=GREEN)
        add_textbox(s, Inches(10.20), ry + Inches(0.30),
                    Inches(2.55), Inches(0.55),
                    foot, size=9.5, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bottom_strip(s, "DEMO URL",
                 "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture — the architecture overview · root URL for the customer-facing demo.",
                 [
                     ("01", "Today",       "Open the demo in front of one teammate · 30 min."),
                     ("02", "This week",   "Pick a client + a Pack · email + book the BVA · 4 hr."),
                     ("03", "Week 2",      "Run the BVA · capture scenario shortlist."),
                     ("04", "Week 3–4",    "Pack Lite SOW under DCIF · sign · kick off pilot."),
                 ])
    page_footer(s, 15, TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"size: {OUT.stat().st_size:,} bytes · {TOTAL} slides")


if __name__ == "__main__":
    main()
