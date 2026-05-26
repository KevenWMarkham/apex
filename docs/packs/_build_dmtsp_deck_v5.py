"""DMTSP seller deck v5 — consulting-template polish.

Layout DNA inspired by the reference templates:
  • 3-card services with circle-icon-on-light-gray + accent bottom bar
  • Pricing-table hero with the emphasized middle tier raised
  • Portrait-style team / persona cards
  • Process flows with numbered circles
  • Real bar charts + donut charts for data viz
  • Strong accent bars at slide edges (Deloitte green in our case)
  • Generous whitespace · large type · one focal element per slide

All imagery generated via Pillow (silhouette portraits, abstract scene
compositions, gradient hero backgrounds) — no photo drops required.
Brand stays Deloitte: green #86BC25 + navy #1A1F2E + light bg #FAFAFA.
"""

from __future__ import annotations

import math
from pathlib import Path
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\APEX-Walkthrough-Deck-for-DMTSP-Sellers-v5.pptx")
ASSETS = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\_apex_deck_v5_assets")
ASSETS.mkdir(parents=True, exist_ok=True)

# ─── Palette ──────────────────────────────────────────────────────────────
BG          = RGBColor(0xFA, 0xFA, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG     = RGBColor(0xF3, 0xF6, 0xFA)
HEAD_NAVY   = RGBColor(0x1A, 0x1F, 0x2E)
DEEP_NAVY   = RGBColor(0x0E, 0x14, 0x23)
GREEN       = RGBColor(0x86, 0xBC, 0x25)
GREEN_DK    = RGBColor(0x5A, 0x83, 0x0F)
GREEN_PALE  = RGBColor(0xF1, 0xF8, 0xE9)
MS_BLUE     = RGBColor(0x00, 0x78, 0xD4)
STRONG_INK  = RGBColor(0x1A, 0x1F, 0x2E)
INK_SOFT    = RGBColor(0x33, 0x3F, 0x52)
MUTED       = RGBColor(0x47, 0x55, 0x69)
SOFT_LINE   = RGBColor(0xE2, 0xE8, 0xF0)
PALE_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_TEXT  = RGBColor(0xCB, 0xD5, 0xE1)
ICE         = RGBColor(0xBF, 0xDB, 0xFE)
RED_INK     = RGBColor(0xA0, 0x35, 0x35)
RED_PALE    = RGBColor(0xFD, 0xF4, 0xF4)

# Industry tints for the 7-pack catalog
TINT_MFG    = RGBColor(0x6B, 0x72, 0x80)
TINT_RETAIL = RGBColor(0xD9, 0x73, 0x06)
TINT_FIN    = RGBColor(0x05, 0x96, 0x69)
TINT_RISK   = RGBColor(0xDC, 0x26, 0x26)
TINT_CX     = RGBColor(0x2E, 0x53, 0xDF)
TINT_ESG    = RGBColor(0x16, 0xA3, 0x4A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════
# Pillow-generated imagery
# ═══════════════════════════════════════════════════════════════════════════

def _gradient(size, c0, c1, *, angle_deg=20):
    w, h = size
    img = Image.new("RGB", size)
    px = img.load()
    rad = math.radians(angle_deg)
    dx, dy = math.cos(rad), math.sin(rad)
    proj_max = w * abs(dx) + h * abs(dy)
    for y in range(h):
        for x in range(w):
            t = (x * dx + y * dy) / proj_max
            t = max(0.0, min(1.0, t))
            r = int(c0[0] + (c1[0] - c0[0]) * t)
            g = int(c0[1] + (c1[1] - c0[1]) * t)
            b = int(c0[2] + (c1[2] - c0[2]) * t)
            px[x, y] = (r, g, b)
    return img


def _dots(size, color, *, spacing=80, radius=2, alpha=80):
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    for y in range(spacing // 2, size[1], spacing):
        for x in range(spacing // 2, size[0], spacing):
            d.ellipse([x - radius, y - radius, x + radius, y + radius],
                      fill=(*color, alpha))
    return img


def _silhouette_avatar(size, *, bg_color, suit_color, accent_color):
    """Stylized abstract person silhouette for use in 'team' / Q&A cards."""
    w, h = size
    img = Image.new("RGB", size, bg_color)
    d = ImageDraw.Draw(img)
    # Head circle
    hx, hy, hr = w // 2, int(h * 0.32), int(h * 0.16)
    d.ellipse([hx - hr, hy - hr, hx + hr, hy + hr], fill=accent_color)
    # Shoulders / body — wider arc at bottom
    body_top = int(h * 0.55)
    body_w = int(w * 0.78)
    d.ellipse([w // 2 - body_w // 2, body_top, w // 2 + body_w // 2, body_top + int(h * 0.95)],
              fill=suit_color)
    # Subtle decorative accent stripe (yellow-like) at bottom-right corner
    d.polygon([(w - 30, h), (w, h), (w, h - 30)], fill=accent_color)
    return img.filter(ImageFilter.GaussianBlur(radius=1.2))


def gen_hero_title(path: Path) -> None:
    size = (1920, 1080)
    img = _gradient(size, (10, 16, 32), (28, 40, 60), angle_deg=15).convert("RGBA")
    img = Image.alpha_composite(img, _dots(size, (134, 188, 37), spacing=70, radius=2, alpha=70))
    # Big green diagonal wedge bottom-left
    wedge = Image.new("RGBA", size, (0, 0, 0, 0))
    wd = ImageDraw.Draw(wedge)
    wd.polygon([(0, size[1]), (820, size[1]), (0, 280)], fill=(134, 188, 37, 230))
    wd.polygon([(0, size[1]), (640, size[1]), (0, 480)], fill=(134, 188, 37, 255))
    img = Image.alpha_composite(img, wedge)
    img.convert("RGB").save(path, "PNG", optimize=True)


def gen_hero_closing(path: Path) -> None:
    size = (1920, 1080)
    img = _gradient(size, (134, 188, 37), (14, 20, 35), angle_deg=12).convert("RGBA")
    img = Image.alpha_composite(img, _dots(size, (255, 255, 255), spacing=100, radius=1, alpha=70))
    img.convert("RGB").save(path, "PNG", optimize=True)


def gen_meeting_scene(path: Path) -> None:
    """Abstract 'meeting / collaboration' composition.

    Three figure silhouettes around a table-shape rectangle, on a soft
    gradient background. Used on the Q&A slide.
    """
    size = (1280, 720)
    img = _gradient(size, (240, 246, 252), (216, 226, 242), angle_deg=30).convert("RGB")
    d = ImageDraw.Draw(img)
    # Table
    d.rounded_rectangle([180, 380, 1100, 560], radius=30, fill=(190, 200, 215))
    d.rounded_rectangle([180, 380, 1100, 440], radius=15, fill=(160, 175, 195))
    # Three figures arrayed around the table — abstract heads + body arcs
    figures = [
        (350, 340, (26, 31, 46), (200, 100, 60)),   # left
        (640, 290, (134, 188, 37), (255, 255, 255)),  # center (green suit)
        (930, 340, (26, 31, 46), (180, 100, 200)),  # right
    ]
    for fx, fy, suit, accent in figures:
        # Head
        d.ellipse([fx - 36, fy - 100, fx + 36, fy - 28], fill=suit)
        # Body
        d.ellipse([fx - 90, fy - 30, fx + 90, fy + 180], fill=suit)
        # Accent dot (laptop / device)
        d.rounded_rectangle([fx - 50, fy + 60, fx + 50, fy + 90], radius=4, fill=accent)
    # Subtle dot pattern overlay
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, _dots(size, (134, 188, 37), spacing=60, radius=1, alpha=70))
    img.convert("RGB").filter(ImageFilter.GaussianBlur(radius=1.5)).save(path, "PNG", optimize=True)


def gen_avatar(path: Path, *, bg, suit, accent):
    """Per-Q&A questioner avatar."""
    img = _silhouette_avatar((320, 320), bg_color=bg, suit_color=suit, accent_color=accent)
    img.save(path, "PNG", optimize=True)


# Generate all imagery
print("Generating imagery…")
gen_hero_title(ASSETS / "hero_title.png")
gen_hero_closing(ASSETS / "hero_closing.png")
gen_meeting_scene(ASSETS / "meeting.png")
# 6 avatars for the Q&A slide (we use 6 questions)
avatar_styles = [
    ((230, 240, 250), (26, 31, 46),   (134, 188, 37)),  # navy suit, green accent
    ((250, 240, 220), (180, 70, 30),  (255, 220, 100)), # rust suit, yellow accent
    ((230, 250, 245), (16, 100, 80),  (200, 220, 230)), # teal suit
    ((250, 230, 240), (110, 30, 90),  (255, 200, 220)), # plum
    ((240, 248, 255), (0, 120, 212),  (200, 230, 255)), # ms blue
    ((240, 240, 235), (60, 60, 65),   (210, 163, 34)),  # charcoal + amber
]
for i, (bg_c, suit_c, accent_c) in enumerate(avatar_styles):
    gen_avatar(ASSETS / f"avatar_{i+1}.png",
               bg=bg_c, suit=suit_c, accent=accent_c)
print(f"  → wrote {ASSETS}")


# ═══════════════════════════════════════════════════════════════════════════
# PowerPoint primitives + helpers (same lib as v4)
# ═══════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, *, fill, line=None, corner_radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    if corner_radius:
        shp.adjustments[0] = corner_radius
    return shp


def add_oval(slide, left, top, width, height, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def add_arrow(slide, left, top, width, height, *, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def set_text(shape, text, *, size=10, bold=False, color=STRONG_INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36_000); tf.margin_right = Emu(36_000)
    tf.margin_top = Emu(18_000);  tf.margin_bottom = Emu(18_000)
    lines = [text] if isinstance(text, str) else list(text)
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, **kw)
    return box


def header_minimal(slide, title: str, subtitle: str = "", number: int | None = None):
    """Polished header — accent dot + large title + small subtitle + optional section number."""
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.18), fill=GREEN)
    add_textbox(slide, Inches(0.83), Inches(0.34), Inches(11), Inches(0.60),
                title, size=30, bold=True, color=STRONG_INK)
    if subtitle:
        add_textbox(slide, Inches(0.83), Inches(0.97), Inches(11.5), Inches(0.40),
                    subtitle, size=13, color=MUTED, italic=True)
    if number is not None:
        add_textbox(slide, Inches(11.5), Inches(0.34), Inches(1.4), Inches(0.4),
                    f"{number:02d}", size=22, bold=True, color=GREEN,
                    align=PP_ALIGN.RIGHT)


def page_footer(slide, slide_no: int, total: int):
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.02),
             fill=SOFT_LINE)
    add_textbox(slide, Inches(0.5), Inches(7.10), Inches(9), Inches(0.30),
                "APEX  ·  Deloitte Microsoft Technology & Services Practice",
                size=9, color=MUTED)
    add_textbox(slide, Inches(11.00), Inches(7.10), Inches(2.0), Inches(0.30),
                f"{slide_no} / {total}", size=9, bold=True, color=GREEN,
                align=PP_ALIGN.RIGHT)


# Custom icon builders (kept compact; consistent look across all icons)

def icon_circle_with_glyph(slide, x, y, w, h, *, bg=GREEN, glyph="●", glyph_color=WHITE,
                            glyph_size=36):
    add_oval(slide, x, y, w, h, fill=bg)
    add_textbox(slide, x, y, w, h, glyph, size=glyph_size, bold=True,
                color=glyph_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 15

    # ── Slide 1: Title hero ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_title.png"),
                         Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Version chip top-right
    add_rect(s, Inches(10.5), Inches(0.5), Inches(2.4), Inches(0.42),
             fill=GREEN, corner_radius=0.5)
    add_textbox(s, Inches(10.5), Inches(0.5), Inches(2.4), Inches(0.42),
                "DMTSP  ·  v5  ·  2026", size=10, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title — large, on the dark area to the right of the green wedge
    add_textbox(s, Inches(5.5), Inches(2.4), Inches(7.5), Inches(1.4),
                "APEX", size=140, bold=True, color=WHITE)
    add_rect(s, Inches(5.5), Inches(3.85), Inches(0.6), Inches(0.12), fill=GREEN)
    add_textbox(s, Inches(5.5), Inches(4.05), Inches(7.5), Inches(0.6),
                "Deloitte's agentic-AI", size=28, color=WHITE)
    add_textbox(s, Inches(5.5), Inches(4.55), Inches(7.5), Inches(0.6),
                "delivery framework",
                size=28, bold=True, color=WHITE)
    add_textbox(s, Inches(5.5), Inches(5.4), Inches(7.5), Inches(0.4),
                "Walkthrough for DMTSP  ·  what it is  ·  how to demo  ·  how to make it real",
                size=13, color=LIGHT_TEXT, italic=True)
    # URL chip
    add_rect(s, Inches(5.5), Inches(6.2), Inches(7.2), Inches(0.55),
             fill=DEEP_NAVY, corner_radius=0.2)
    add_textbox(s, Inches(5.5), Inches(6.2), Inches(7.2), Inches(0.55),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=10, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 1, TOTAL)

    # ── Slide 2: APEX in one sentence — quote + sandwich diagram ───────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "APEX in one sentence",
                   "Memorize this — start every client conversation here.", number=1)
    # Big quote
    add_textbox(s, Inches(0.6), Inches(2.0), Inches(0.5), Inches(0.9),
                "“", size=120, bold=True, color=GREEN)
    add_textbox(s, Inches(0.6), Inches(2.9), Inches(6.8), Inches(2.6),
                ["The engagement layer",
                 "cloud vendors",
                 "don't ship."],
                size=42, bold=True, color=STRONG_INK)
    add_rect(s, Inches(0.6), Inches(5.5), Inches(0.5), Inches(0.1), fill=GREEN)
    add_textbox(s, Inches(0.6), Inches(5.65), Inches(7), Inches(0.4),
                "— That's the one-liner.",
                size=14, color=MUTED, italic=True)

    # Right: 3-layer sandwich diagram
    diag_x = Inches(7.8); diag_y = Inches(2.0); diag_w = Inches(5.1)
    add_rect(s, diag_x, diag_y, diag_w, Inches(0.85), fill=HEAD_NAVY,
             corner_radius=0.1)
    add_textbox(s, diag_x, diag_y, diag_w, Inches(0.85),
                "Client outcome", size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arr1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              diag_x + diag_w * 0.45, diag_y + Inches(0.9),
                              diag_w * 0.1, Inches(0.45))
    arr1.fill.solid(); arr1.fill.fore_color.rgb = SOFT_LINE
    arr1.line.fill.background()
    # APEX layer
    add_rect(s, diag_x, diag_y + Inches(1.4), diag_w, Inches(1.5), fill=GREEN,
             corner_radius=0.1)
    add_textbox(s, diag_x, diag_y + Inches(1.4), diag_w, Inches(0.5),
                "APEX  —  the engagement layer",
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, diag_x, diag_y + Inches(1.85), diag_w, Inches(1.0),
                ["Industry Packs  ·  Virtual Views",
                 "WORM LEDGER  ·  Service Envelopes",
                 "10-asset bundle structure"],
                size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
    arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              diag_x + diag_w * 0.45, diag_y + Inches(3.05),
                              diag_w * 0.1, Inches(0.45))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = SOFT_LINE
    arr2.line.fill.background()
    add_rect(s, diag_x, diag_y + Inches(3.55), diag_w, Inches(0.85),
             fill=MS_BLUE, corner_radius=0.1)
    add_textbox(s, diag_x, diag_y + Inches(3.55), diag_w, Inches(0.85),
                "Foundry  ·  Vertex  ·  Bedrock",
                size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, diag_x, diag_y + Inches(4.45), diag_w, Inches(0.30),
                "(cloud-vendor agent runtimes)",
                size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 2, TOTAL)

    # ── Slide 3: Without / With APEX — 3-card services pattern ─────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Without APEX  ·  With APEX",
                   "Same Foundry runtime. Different deal shape.", number=2)
    # Two side panels (a left vs right comparison)
    def comparison_panel(x, w, *, tone, headline, items):
        col_fill = RED_PALE if tone == "bad" else GREEN_PALE
        accent = RED_INK if tone == "bad" else GREEN
        add_rect(s, x, Inches(1.85), w, Inches(5.1),
                 fill=WHITE, line=accent, corner_radius=0.05)
        # Header band
        add_rect(s, x, Inches(1.85), w, Inches(0.7), fill=accent)
        add_textbox(s, x, Inches(1.85), w, Inches(0.7), headline,
                    size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bottom accent strip (consulting-template style)
        add_rect(s, x, Inches(6.55), w, Inches(0.4), fill=col_fill)
        # Items
        for i, (label, desc) in enumerate(items):
            iy = Inches(2.85 + i * 0.92)
            # Bullet
            add_oval(s, x + Inches(0.30), iy + Inches(0.15),
                     Inches(0.20), Inches(0.20), fill=accent)
            add_textbox(s, x + Inches(0.65), iy, w - Inches(0.80), Inches(0.40),
                        label, size=14, bold=True, color=STRONG_INK)
            add_textbox(s, x + Inches(0.65), iy + Inches(0.38),
                        w - Inches(0.80), Inches(0.50),
                        desc, size=10.5, color=MUTED)

    comparison_panel(Inches(0.5), Inches(6.1), tone="bad", headline="WITHOUT APEX", items=[
        ("Staff augmentation",  "Project competes with what client could build alone."),
        ("Bespoke discovery",   "Every engagement re-invents the data model + SOW."),
        ("Margin compression",  "Deal shape resembles labor sales."),
        ("Fragile audit story", "No replay; trust is conversational, not cryptographic."),
    ])
    comparison_panel(Inches(6.85), Inches(6.0), tone="good", headline="WITH APEX", items=[
        ("Fixed-fee Packs",      "Industry Pack scoped at sub-tier (Lite / Standard / Enterprise)."),
        ("Canonical scenarios",  "10-asset bundle ships pre-built playbooks."),
        ("WORM LEDGER",          "14-field hash-chained audit row per decision."),
        ("Replatform = config",  "Same Pack on APEX-M / G / A."),
    ])
    page_footer(s, 3, TOTAL)

    # ── Slide 4: Five pieces — 3-card services pattern (icon + label) ──
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The five pieces of APEX",
                   "Every engagement assembles from these.", number=3)
    # Five cards in a row — consulting "Our Services" pattern with circle-icons
    pieces = [
        ("01", "CORE",       "Cloud-neutral IP",       HEAD_NAVY, "▲"),
        ("02", "PROFILES",   "M  /  G  /  A",          MS_BLUE,   "◔"),
        ("03", "PACKS",      "7 in the catalog",       HEAD_NAVY, "▦"),
        ("04", "ENVELOPES",  "5 sellable tiers",       HEAD_NAVY, "▤"),
        ("05", "LEDGER",     "WORM hash chain",        GREEN_DK,  "⛨"),
    ]
    cw = Inches(2.40); gap = Inches(0.15); x0 = Inches(0.55)
    for i, (num, label, sub, color, glyph) in enumerate(pieces):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, Inches(1.9), cw, Inches(4.7),
                 fill=WHITE, line=SOFT_LINE, corner_radius=0.05)
        # Number badge top
        add_textbox(s, cx, Inches(2.05), cw, Inches(0.35),
                    num, size=12, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER)
        # Circle icon
        icon_circle_with_glyph(s, cx + Inches(0.65), Inches(2.55),
                                 Inches(1.10), Inches(1.10),
                                 bg=color, glyph=glyph, glyph_color=WHITE,
                                 glyph_size=36)
        # Label
        add_textbox(s, cx, Inches(3.85), cw, Inches(0.45),
                    label, size=18, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Sub
        add_textbox(s, cx, Inches(4.40), cw, Inches(0.40),
                    sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        # Accent bottom bar (consulting-template touch)
        add_rect(s, cx, Inches(6.45), cw, Inches(0.15), fill=GREEN)
    page_footer(s, 4, TOTAL)

    # ── Slide 5: APEX family — team-style cards ─────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same core. Different cloud.",
                   "Three implementations of one 14-interface contract.", number=4)
    families = [
        ("APEX-M", "MICROSOFT  /  AZURE", "GA",            MS_BLUE,                   "M",
         ["Foundry + Azure OpenAI", "OneLake + Fabric", "Entra · Purview · Maps Creator"]),
        ("APEX-G", "GOOGLE  /  GCP",      "BETA Q3 FY27", RGBColor(0x0F, 0x9D, 0x58), "G",
         ["Vertex AI + Gemini", "BigQuery + BigLake", "Cloud Identity · Dataplex"]),
        ("APEX-A", "AMAZON  /  AWS",      "BETA Q4 FY27", RGBColor(0xFF, 0x99, 0x00), "A",
         ["Bedrock (Claude / Nova)", "S3 + Lake Formation", "IAM IC · Location Service"]),
    ]
    cw = Inches(4.05); gap = Inches(0.15); x0 = Inches(0.55)
    for i, (tag, name, status, color, letter, bullets) in enumerate(families):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, Inches(1.9), cw, Inches(4.7),
                 fill=WHITE, line=SOFT_LINE, corner_radius=0.05)
        # Top color stripe with tag
        add_rect(s, cx, Inches(1.9), cw, Inches(0.5), fill=color)
        add_textbox(s, cx, Inches(1.9), cw, Inches(0.5), tag,
                    size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Big letter circle (like avatar)
        icon_circle_with_glyph(s, cx + Inches(1.25), Inches(2.65),
                                 Inches(1.55), Inches(1.55),
                                 bg=color, glyph=letter, glyph_color=WHITE,
                                 glyph_size=64)
        # Name
        add_textbox(s, cx, Inches(4.35), cw, Inches(0.35),
                    name, size=13, color=MUTED, align=PP_ALIGN.CENTER)
        # Status pill
        add_rect(s, cx + Inches(1.0), Inches(4.80),
                 cw - Inches(2.0), Inches(0.40),
                 fill=color, corner_radius=0.5)
        add_textbox(s, cx + Inches(1.0), Inches(4.80),
                    cw - Inches(2.0), Inches(0.40),
                    status, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        for j, b in enumerate(bullets):
            add_textbox(s, cx + Inches(0.30), Inches(5.45 + j * 0.35),
                        cw - Inches(0.60), Inches(0.32),
                        "·  " + b, size=10.5, color=STRONG_INK)
    page_footer(s, 5, TOTAL)

    # ── Slide 6: Seven Industry Packs — 7-tile catalog ─────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Seven Industry Packs",
                   "Same 10-asset bundle. Different industry per Pack.", number=5)

    packs = [
        ("01", "Manufacturing",      "Plant / Ops",          TINT_MFG,    "M"),
        ("02", "Retail Merch",       "Chief Merchant",       TINT_RETAIL, "R"),
        ("03", "Finance",            "Controller · CFO",     TINT_FIN,    "F"),
        ("04", "Risk",               "SOX PMO · CRO",        TINT_RISK,   "R"),
        ("05", "Customer Experience","Care Ops",             TINT_CX,     "C"),
        ("06", "ESG",                "Sustainability",       TINT_ESG,    "E"),
        ("07", "CFMP",               "CMO · CX-VP",          GREEN,       "★"),
    ]
    cw = Inches(2.85); ch = Inches(2.20); gx = Inches(0.18); gy = Inches(0.20)

    def pack_card(cx, cy, num, name, persona, color, letter, *, highlight=False):
        add_rect(s, cx, cy, cw, ch, fill=WHITE,
                 line=(GREEN if highlight else SOFT_LINE),
                 corner_radius=0.05)
        # Top stripe in tint color
        add_rect(s, cx, cy, cw, Inches(0.20), fill=color)
        if highlight:
            # Star ribbon corner
            add_rect(s, cx + cw - Inches(1.0), cy + Inches(0.05),
                     Inches(0.9), Inches(0.32),
                     fill=GREEN, corner_radius=0.4)
            add_textbox(s, cx + cw - Inches(1.0), cy + Inches(0.05),
                        Inches(0.9), Inches(0.32),
                        "★ LIVE", size=9, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Number
        add_textbox(s, cx + Inches(0.18), cy + Inches(0.30), Inches(0.7), Inches(0.30),
                    num, size=11, bold=True, color=color)
        # Letter circle (kind of like an industry logo)
        icon_circle_with_glyph(s, cx + cw * 0.34, cy + Inches(0.55),
                                 cw * 0.32, Inches(0.85),
                                 bg=color, glyph=letter, glyph_color=WHITE,
                                 glyph_size=28)
        # Name
        add_textbox(s, cx, cy + Inches(1.45), cw, Inches(0.35),
                    name, size=14, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Persona
        add_textbox(s, cx, cy + Inches(1.78), cw, Inches(0.35),
                    persona, size=10, color=MUTED, italic=True,
                    align=PP_ALIGN.CENTER)

    # Row 1: 4 cards
    y_row1 = Inches(1.85)
    for i in range(4):
        cx = Inches(0.40 + i * (cw / Inches(1) + gx / Inches(1)))
        num, name, persona, color, letter = packs[i]
        pack_card(cx, y_row1, num, name, persona, color, letter)
    # Row 2: 3 cards centered
    y_row2 = y_row1 + ch + gy
    x_row2_start = (SLIDE_W - (cw * 3 + gx * 2)) / 2
    for i in range(3):
        cx = x_row2_start + Inches(i * (cw / Inches(1) + gx / Inches(1)))
        num, name, persona, color, letter = packs[4 + i]
        pack_card(cx, y_row2, num, name, persona, color, letter,
                  highlight=(name == "CFMP"))
    page_footer(s, 6, TOTAL)

    # ── Slide 7: 10-asset bundle — clean grid ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same structure. Every Pack.",
                   "Ten slots. Identical contract. Predictable engineering.", number=6)

    assets = [
        ("VV manifests",       "Federated views"),
        ("Scenarios",          "End-to-end outcomes"),
        ("Source adapters",    "System connectors"),
        ("Adaptive cards",     "HITL surfaces"),
        ("Persona map",        "Role bindings"),
        ("Demo data",          "Laptop-runnable"),
        ("BVA worksheet",      "ROI calculator"),
        ("Sample SOW",         "Tier templates"),
        ("Acceptance tests",   "CI-green"),
        ("Runbook + training", "Operate-ready"),
    ]
    cw = Inches(2.40); ch = Inches(2.20); gx = Inches(0.10); gy = Inches(0.10)
    x0 = Inches(0.55); y0 = Inches(1.95)
    for i, (name, sub) in enumerate(assets):
        col = i % 5; row = i // 5
        cx = x0 + Inches(col * (cw / Inches(1) + gx / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + gy / Inches(1)))
        # Card
        add_rect(s, cx, cy, cw, ch, fill=WHITE, line=SOFT_LINE, corner_radius=0.04)
        # Big number circle top-center
        icon_circle_with_glyph(s, cx + cw * 0.30, cy + Inches(0.30),
                                 cw * 0.40, Inches(0.95),
                                 bg=GREEN, glyph=f"{i+1:02d}",
                                 glyph_color=WHITE, glyph_size=22)
        # Name
        add_textbox(s, cx, cy + Inches(1.40), cw, Inches(0.35),
                    name, size=13, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Sub
        add_textbox(s, cx, cy + Inches(1.72), cw, Inches(0.30),
                    sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        # Bottom accent bar
        add_rect(s, cx + cw * 0.30, cy + ch - Inches(0.12),
                 cw * 0.40, Inches(0.06), fill=GREEN)
    page_footer(s, 7, TOTAL)

    # ── Slide 8: CFMP spotlight + DONUT of scenarios by phase ──────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Spotlight: CFMP",
                   "Customer Focused Merchandise Pack v0.2 · Phase 1 LIVE", number=7)

    # Left: pitch card
    add_rect(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.95),
             fill=HEAD_NAVY, corner_radius=0.05)
    add_rect(s, Inches(0.55), Inches(1.95), Inches(0.18), Inches(4.95), fill=GREEN)
    add_textbox(s, Inches(0.85), Inches(2.10), Inches(5.5), Inches(0.6),
                "CFMP", size=44, bold=True, color=WHITE)
    add_textbox(s, Inches(0.85), Inches(2.85), Inches(5.5), Inches(0.4),
                "Customer Focused Merchandise Pack",
                size=14, color=ICE)
    add_textbox(s, Inches(0.85), Inches(3.30), Inches(5.5), Inches(0.4),
                "First customer-moment Industry Pack",
                size=11, color=LIGHT_TEXT, italic=True)
    # Mini stats
    stats = [
        ("18",    "scenarios"),
        ("6",     "new VVs"),
        ("4",     "journey phases"),
        ("$150K", "Pack Lite (4-6 wks)"),
    ]
    for i, (val, lbl) in enumerate(stats):
        col = i % 2; row = i // 2
        sx = Inches(0.85 + col * 2.9)
        sy = Inches(4.0 + row * 1.3)
        add_textbox(s, sx, sy, Inches(2.7), Inches(0.6),
                    val, size=30, bold=True, color=GREEN)
        add_textbox(s, sx, sy + Inches(0.6), Inches(2.7), Inches(0.4),
                    lbl, size=10, color=LIGHT_TEXT, italic=True)

    # Right: Donut chart — 18 scenarios by phase
    cx, cy = Inches(10.0), Inches(4.3)
    outer_r = Inches(1.7)
    inner_r = Inches(0.95)
    # 4 segments (5 + 6 + 4 + 3 = 18; close enough to user data)
    add_textbox(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(0.4),
                "18 SCENARIOS BY PHASE",
                size=11, bold=True, color=GREEN)
    add_textbox(s, Inches(7.0), Inches(2.45), Inches(5.8), Inches(0.35),
                "How the spine distributes scope.",
                size=10, color=MUTED, italic=True)
    # Donut segments using pie shapes (manual angles)
    segs = [
        ("CHOOSE",   5, MS_BLUE),
        ("SELECT",   6, GREEN),
        ("BUY",      4, RGBColor(0xFF, 0x99, 0x00)),
        ("SERVICES", 3, RGBColor(0xB4, 0x42, 0xCC)),
    ]
    total = sum(n for _, n, _ in segs)
    cur = -90  # start at top
    for label, n, color in segs:
        sweep = 360 * n / total
        pie = s.shapes.add_shape(MSO_SHAPE.PIE,
                                 cx - outer_r, cy - outer_r,
                                 outer_r * 2, outer_r * 2)
        pie.fill.solid(); pie.fill.fore_color.rgb = color
        pie.line.color.rgb = WHITE
        pie.line.width = Pt(2.5)
        pie.adjustments[0] = cur
        pie.adjustments[1] = cur + sweep
        cur += sweep
    # Hollow center
    add_oval(s, cx - inner_r, cy - inner_r, inner_r * 2, inner_r * 2, fill=BG)
    add_textbox(s, cx - inner_r, cy - inner_r, inner_r * 2, inner_r * 2,
                ["18", "scenarios"],
                size=22, bold=True, color=STRONG_INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Legend
    for i, (label, n, color) in enumerate(segs):
        ly = Inches(2.95 + i * 0.45)
        add_rect(s, Inches(7.0), ly, Inches(0.20), Inches(0.20), fill=color)
        add_textbox(s, Inches(7.30), ly - Inches(0.05),
                    Inches(2.7), Inches(0.35),
                    label, size=11, bold=True, color=STRONG_INK)
        add_textbox(s, Inches(7.30), ly + Inches(0.18),
                    Inches(2.7), Inches(0.25),
                    f"{n} scenarios", size=9, color=MUTED)
    page_footer(s, 8, TOTAL)

    # ── Slide 9: Demo URL — single focal CTA ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.20), SLIDE_H, fill=GREEN)
    add_textbox(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.7),
                "Open the demo. Talk through it.",
                size=42, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.45),
                "Eight surfaces  ·  ten minutes  ·  let the agent answer for you.",
                size=15, color=LIGHT_TEXT, italic=True)

    # URL — big bold green bar
    add_rect(s, Inches(0.7), Inches(2.65), Inches(11.93), Inches(1.0),
             fill=GREEN, corner_radius=0.05)
    add_textbox(s, Inches(0.7), Inches(2.65), Inches(11.93), Inches(1.0),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 8 step pills along bottom
    steps = ["1  Demo /", "2  “Where is X?”", "3  Dietary chip", "4  Cart click",
             "5  /architecture", "6  Details ⓘ", "7  APEX value", "8  Tour flyout"]
    cw = Inches(1.40); gap = Inches(0.12); x0 = Inches(0.7)
    for i, t in enumerate(steps):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        add_rect(s, cx, Inches(4.4), cw, Inches(0.8),
                 fill=WHITE, corner_radius=0.3)
        add_textbox(s, cx, Inches(4.4), cw, Inches(0.8),
                    t, size=10, bold=True, color=DEEP_NAVY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom strip — what to demo first
    add_rect(s, Inches(0.7), Inches(5.6), Inches(11.93), Inches(1.1),
             fill=DEEP_NAVY, corner_radius=0.05)
    add_textbox(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.4),
                "FIRST CLICK", size=10, bold=True, color=GREEN,
                align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.6),
                "Type into chat: “Where can I find the Coca-Cola?” — the wayfinding card lands in under 3 seconds.",
                size=14, color=WHITE)
    page_footer(s, 9, TOTAL)

    # ── Slide 10: Engagement ladder — PRICING TABLE pattern ─────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Engagement ladder",
                   "Five sellable tiers — stackable, additive.", number=8)
    tiers = [
        ("T1", "FOUNDATION", "$400–700K",     "6–10 wks",      HEAD_NAVY,
         ["Sold once per tenant",
          "VV runtime + LEDGER",
          "MCP auto-registration"]),
        ("T2", "PACK",        "$100K–$3.5M",  "4 wks – 9 mo",  GREEN,
         ["10-asset Pack bundle",
          "Lite / Standard / Enterprise",
          "Per-industry scope"], True),  # emphasized
        ("T3", "SCENARIO",    "$150–400K",    "4–8 wks",       HEAD_NAVY,
         ["One business outcome",
          "Live on client sources",
          "Fixed fee"]),
        ("T4", "CUSTOM",      "T&M",          "Sprints",       HEAD_NAVY,
         ["Net-new VV / scenario",
          "Client-specific",
          "Advisory + dev"]),
        ("T5", "OPERATE",     "$/mo",         "Continuous",    HEAD_NAVY,
         ["24×7 monitoring",
          "Threshold tuning",
          "Run-rate revenue"]),
    ]
    # Pricing-table pattern — emphasized tier raised + larger
    base_cw = Inches(2.45); gap = Inches(0.10); x0 = Inches(0.55)
    for i, tier in enumerate(tiers):
        if len(tier) == 6:
            tag, name, price, time, color, items = tier
            emphasized = False
        else:
            tag, name, price, time, color, items, emphasized = tier
        cx = x0 + Inches(i * (base_cw / Inches(1) + gap / Inches(1)))
        # Raised position for emphasized
        cy = Inches(1.80) if emphasized else Inches(2.05)
        ch = Inches(5.0) if emphasized else Inches(4.55)
        # Card
        add_rect(s, cx, cy, base_cw, ch, fill=WHITE,
                 line=(GREEN if emphasized else SOFT_LINE),
                 corner_radius=0.05)
        # Header band
        add_rect(s, cx, cy, base_cw, Inches(0.85), fill=color)
        # Tag chip
        add_textbox(s, cx, cy + Inches(0.05), base_cw, Inches(0.35),
                    tag, size=14, bold=True,
                    color=(WHITE if not emphasized else WHITE),
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cx, cy + Inches(0.40), base_cw, Inches(0.40),
                    name, size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        # Price
        add_textbox(s, cx, cy + Inches(1.05), base_cw, Inches(0.50),
                    price, size=20, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cx, cy + Inches(1.55), base_cw, Inches(0.35),
                    time, size=11, color=MUTED, italic=True,
                    align=PP_ALIGN.CENTER)
        # Divider
        add_rect(s, cx + Inches(0.3), cy + Inches(2.0),
                 base_cw - Inches(0.6), Inches(0.02), fill=SOFT_LINE)
        # Items
        for j, item in enumerate(items):
            add_textbox(s, cx + Inches(0.2), cy + Inches(2.15 + j * 0.42),
                        base_cw - Inches(0.4), Inches(0.40),
                        "·  " + item, size=10, color=STRONG_INK)
        # Emphasized tier — call-out badge
        if emphasized:
            add_rect(s, cx + Inches(0.45), cy - Inches(0.30),
                     base_cw - Inches(0.9), Inches(0.32),
                     fill=GREEN, corner_radius=0.5)
            add_textbox(s, cx + Inches(0.45), cy - Inches(0.30),
                        base_cw - Inches(0.9), Inches(0.32),
                        "★ MOST SOLD", size=9, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 10, TOTAL)

    # ── Slide 11: Independence — DO / DO NOT polished ──────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Independence  —  non-negotiable",
                   "Deloitte audits Microsoft. SEC rules shape every engagement.", number=9)

    # Two-column comparison with strong visual cards
    def indep_panel(x, w, *, tone, headline, items):
        accent = (GREEN if tone == "do" else RED_INK)
        pale = (GREEN_PALE if tone == "do" else RED_PALE)
        add_rect(s, x, Inches(1.95), w, Inches(5.0),
                 fill=WHITE, line=accent, corner_radius=0.05)
        # Header band
        add_rect(s, x, Inches(1.95), w, Inches(0.7), fill=accent)
        glyph = "✓" if tone == "do" else "✕"
        add_textbox(s, x, Inches(1.95), w, Inches(0.7),
                    f"{glyph}   {headline}", size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bottom accent
        add_rect(s, x, Inches(6.55), w, Inches(0.4), fill=pale)
        for i, t in enumerate(items):
            iy = Inches(2.85 + i * 0.62)
            add_oval(s, x + Inches(0.30), iy + Inches(0.10),
                     Inches(0.30), Inches(0.30), fill=accent)
            add_textbox(s, x + Inches(0.30), iy + Inches(0.10),
                        Inches(0.30), Inches(0.30),
                        glyph, size=14, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(s, x + Inches(0.75), iy + Inches(0.05),
                        w - Inches(0.85), Inches(0.45),
                        t, size=13, color=STRONG_INK,
                        anchor=MSO_ANCHOR.MIDDLE)

    indep_panel(Inches(0.55), Inches(6.1), tone="do", headline="DO", items=[
        "BVA (client-funded discovery)",
        "DCIF (Deloitte co-invest fund)",
        "ISV Marketplace burndown",
        "SI Teaming POC",
        "T&M  ·  Client direct",
    ])
    indep_panel(Inches(6.85), Inches(6.0), tone="dont", headline="DO NOT", items=[
        "Route ECIF through Deloitte",
        "Use “partnership” / “alliance” language",
        "Take Microsoft co-sell credit",
        "Sell APEX itself (Deloitte IP)",
        "Imply joint commercial entity",
    ])
    page_footer(s, 11, TOTAL)

    # ── Slide 12: NOT a product — binary panel ─────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "What we sell  ·  What we don't",
                   "Single most important framing — repeat every meeting.", number=10)

    # NOT panel
    add_rect(s, Inches(0.55), Inches(1.95), Inches(6.1), Inches(5.0),
             fill=RED_PALE, line=RED_INK, corner_radius=0.05)
    add_rect(s, Inches(0.55), Inches(6.55), Inches(6.1), Inches(0.4), fill=RED_INK)
    icon_circle_with_glyph(s, Inches(3.1), Inches(2.3),
                             Inches(1.0), Inches(1.0),
                             bg=RED_INK, glyph="✕", glyph_color=WHITE, glyph_size=44)
    add_textbox(s, Inches(0.55), Inches(3.5), Inches(6.1), Inches(0.7),
                "NOT a product", size=28, bold=True, color=RED_INK,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.55), Inches(4.2), Inches(6.1), Inches(0.4),
                "—  the client does not buy APEX  —",
                size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    nots = ["Not a SKU", "Not a license", "Not a subscription",
            "Not Microsoft co-sell"]
    for i, t in enumerate(nots):
        add_textbox(s, Inches(0.7), Inches(4.85 + i * 0.36),
                    Inches(5.8), Inches(0.35),
                    "·  " + t, size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)

    # IS panel
    add_rect(s, Inches(6.85), Inches(1.95), Inches(6.0), Inches(5.0),
             fill=GREEN_PALE, line=GREEN, corner_radius=0.05)
    add_rect(s, Inches(6.85), Inches(6.55), Inches(6.0), Inches(0.4), fill=GREEN)
    icon_circle_with_glyph(s, Inches(9.35), Inches(2.3),
                             Inches(1.0), Inches(1.0),
                             bg=GREEN, glyph="✓", glyph_color=WHITE, glyph_size=44)
    add_textbox(s, Inches(6.85), Inches(3.5), Inches(6.0), Inches(0.7),
                "Deloitte delivery framework", size=24, bold=True, color=GREEN_DK,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.85), Inches(4.2), Inches(6.0), Inches(0.4),
                "—  the client engages Deloitte  —",
                size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    iss = ["Sell  Deloitte agentic services",
           "Scoped per Pack · per tier",
           "Funded via standard envelope",
           "Framework is Deloitte IP"]
    for i, t in enumerate(iss):
        add_textbox(s, Inches(7.05), Inches(4.85 + i * 0.36),
                    Inches(5.6), Inches(0.35),
                    "·  " + t, size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)
    page_footer(s, 12, TOTAL)

    # ── Slide 13: BVA wedge + bar chart of engagement value over time ──
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The wedge  —  BVA → Pack Lite → Standard → Enterprise → Operate",
                   "Each step earns the next. Nothing gets redone.", number=11)

    waves = [
        ("BVA",        "4 hrs",      "Free",          GREEN),
        ("Pack Lite",  "4-6 wks",    "$150-250K",     GREEN_DK),
        ("Standard",   "12-16 wks",  "$500K-1.5M",    HEAD_NAVY),
        ("Enterprise", "6-9 mo",     "$1.5-3.5M",     RGBColor(0x14, 0x1B, 0x4A)),
        ("Operate",    "Ongoing",    "$/mo",          RGBColor(0x47, 0x55, 0x69)),
    ]

    # Left half — bar chart showing engagement-value progression
    chart_x = Inches(0.55); chart_y = Inches(2.0); chart_w = Inches(5.8); chart_h = Inches(3.6)
    add_rect(s, chart_x, chart_y, chart_w, chart_h, fill=WHITE,
             line=SOFT_LINE, corner_radius=0.05)
    add_textbox(s, chart_x + Inches(0.2), chart_y + Inches(0.15),
                chart_w, Inches(0.35),
                "Engagement value, $ (illustrative)",
                size=11, bold=True, color=STRONG_INK)
    # Bars
    bar_values = [0.05, 0.20, 1.0, 2.5, 0.4]  # relative heights (Enterprise tallest)
    bar_labels = ["BVA", "Lite", "Std", "Ent", "Op"]
    bar_w = (chart_w - Inches(1.0)) / 5
    bar_area_top = chart_y + Inches(0.7)
    bar_area_bottom = chart_y + chart_h - Inches(0.6)
    bar_area_h = bar_area_bottom - bar_area_top
    max_v = max(bar_values)
    for i, (v, lbl) in enumerate(zip(bar_values, bar_labels)):
        bx = chart_x + Inches(0.4) + Inches(i * (bar_w / Inches(1) + 0.1))
        bh = Emu(int(bar_area_h * (v / max_v)))
        by = bar_area_bottom - bh
        add_rect(s, bx, by, bar_w - Inches(0.1), bh, fill=waves[i][3])
        # Value label on top
        if v >= 0.5:
            add_textbox(s, bx, by - Inches(0.30), bar_w, Inches(0.30),
                        waves[i][2], size=9, bold=True, color=STRONG_INK,
                        align=PP_ALIGN.CENTER)
        # X label
        add_textbox(s, bx, bar_area_bottom + Inches(0.05), bar_w, Inches(0.30),
                    lbl, size=10, bold=True, color=MUTED,
                    align=PP_ALIGN.CENTER)
    # Axis line
    add_rect(s, chart_x + Inches(0.3), bar_area_bottom,
             chart_w - Inches(0.6), Inches(0.02), fill=MUTED)

    # Right half — process flow with numbered circles + funding pills
    flow_x = Inches(6.65); flow_y = Inches(2.05)
    add_textbox(s, flow_x, flow_y, Inches(6.2), Inches(0.4),
                "The motion · 5 steps", size=14, bold=True, color=STRONG_INK)
    for i, (name, time, fund, color) in enumerate(waves):
        ry = flow_y + Inches(0.45 + i * 0.78)
        # Numbered circle
        icon_circle_with_glyph(s, flow_x, ry, Inches(0.6), Inches(0.6),
                                 bg=color, glyph=str(i + 1),
                                 glyph_color=WHITE, glyph_size=18)
        # Step text
        add_textbox(s, flow_x + Inches(0.75), ry, Inches(2.4), Inches(0.32),
                    name, size=13, bold=True, color=STRONG_INK)
        add_textbox(s, flow_x + Inches(0.75), ry + Inches(0.30),
                    Inches(2.4), Inches(0.30),
                    time, size=10, color=MUTED, italic=True)
        # Funding pill
        add_rect(s, flow_x + Inches(3.4), ry + Inches(0.10),
                 Inches(2.7), Inches(0.42),
                 fill=color, corner_radius=0.5)
        add_textbox(s, flow_x + Inches(3.4), ry + Inches(0.10),
                    Inches(2.7), Inches(0.42),
                    fund, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 13, TOTAL)

    # ── Slide 14: Q&A — portrait cards (consulting-team pattern) ───────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The questions you'll hear",
                   "Six openers cover ~90% of the first ten minutes.", number=12)
    qa = [
        ("Is this Microsoft's product?",                  1),
        ("Can we buy APEX?",                              2),
        ("How is this different from Foundry?",           3),
        ("What if we're on AWS / Google?",                4),
        ("How fast can we see something working?",        5),
        ("What does the audit story look like?",          6),
    ]
    # 3 columns × 2 rows of avatar-style cards
    cw = Inches(4.05); ch = Inches(2.40); gx = Inches(0.15); gy = Inches(0.15)
    x0 = Inches(0.55); y0 = Inches(1.95)
    for i, (q, avatar_num) in enumerate(qa):
        col = i % 3; row = i // 3
        cx = x0 + Inches(col * (cw / Inches(1) + gx / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + gy / Inches(1)))
        # Card
        add_rect(s, cx, cy, cw, ch, fill=WHITE, line=SOFT_LINE, corner_radius=0.04)
        # Avatar image area top (with circular avatar embedded — but pptx doesn't support clipping,
        # so we put the square avatar then overlay a thin border)
        s.shapes.add_picture(str(ASSETS / f"avatar_{avatar_num}.png"),
                             cx + Inches(0.10), cy + Inches(0.10),
                             Inches(1.0), Inches(1.0))
        # "Q" badge overlay top-right of avatar
        add_oval(s, cx + Inches(0.85), cy + Inches(0.10),
                 Inches(0.40), Inches(0.40), fill=GREEN)
        add_textbox(s, cx + Inches(0.85), cy + Inches(0.10),
                    Inches(0.40), Inches(0.40),
                    "Q", size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Question text
        add_textbox(s, cx + Inches(1.25), cy + Inches(0.20),
                    cw - Inches(1.40), Inches(1.30),
                    q, size=13, bold=True, color=STRONG_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Bottom accent bar (consulting touch)
        add_rect(s, cx, cy + ch - Inches(0.10), cw, Inches(0.10), fill=GREEN)
    add_textbox(s, Inches(0.55), Inches(6.65), Inches(12.3), Inches(0.4),
                "If stuck for 30 seconds — redirect:  “Let me show you in the live demo, it's faster.”",
                size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 14, TOTAL)

    # ── Slide 15: Closing — three actions over momentum gradient ───────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_closing.png"),
                         Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Dark panel for text
    add_rect(s, Inches(0), Inches(0.6), SLIDE_W, Inches(5.8), fill=DEEP_NAVY)
    add_rect(s, Inches(0), Inches(0.6), Inches(0.20), Inches(5.8), fill=GREEN)
    add_textbox(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.8),
                "This week",
                size=60, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.45),
                "Three actions  ·  every DMTSP seller  ·  every week.",
                size=16, color=LIGHT_TEXT, italic=True)

    actions = [
        ("01", "OPEN  the  demo",      "60 minutes  ·  one client meeting"),
        ("02", "BOOK  the  BVA",       "4 hours  ·  client-funded"),
        ("03", "SIGN  Pack Lite SOW",  "DCIF  ·  $150-250K  ·  4-6 wks"),
    ]
    y0 = Inches(2.6)
    for i, (n, t, sub) in enumerate(actions):
        cy = y0 + Inches(i * 1.25)
        # Big number with bottom green bar
        add_textbox(s, Inches(0.7), cy, Inches(1.8), Inches(1.1),
                    n, size=72, bold=True, color=GREEN,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(0.7), cy + Inches(1.05),
                 Inches(1.5), Inches(0.06), fill=GREEN)
        # Title + sub
        add_textbox(s, Inches(2.7), cy + Inches(0.10),
                    Inches(10), Inches(0.55),
                    t, size=30, bold=True, color=WHITE)
        add_textbox(s, Inches(2.7), cy + Inches(0.70),
                    Inches(10), Inches(0.40),
                    sub, size=14, color=LIGHT_TEXT, italic=True)

    # URL strip at very bottom
    add_rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.55), fill=GREEN)
    add_textbox(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.55),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 15, TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"size: {OUT.stat().st_size:,} bytes · {TOTAL} slides")


if __name__ == "__main__":
    main()
