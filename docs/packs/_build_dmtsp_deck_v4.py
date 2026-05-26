"""DMTSP seller deck v4 — diagrams + generated imagery. Self-contained.

Design principles:
  • Consulting-deck polish (McKinsey/BCG style); non-technical audience.
  • Every slide has ONE dominant visual: hero gradient image OR rich
    vector diagram OR custom shape icons composed from primitives.
  • Pillow generates 5 gradient hero backgrounds (saved next to the
    pptx); python-pptx inserts them and composes everything else.
  • No external photo drops needed. Self-contained on first open.
  • Total text ≤ 35 words per content slide; speaker carries narrative.
"""

from __future__ import annotations

from pathlib import Path
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\APEX-Walkthrough-Deck-for-DMTSP-Sellers-v4.pptx")
ASSETS = Path(r"C:\Stage\Clients\Industries\APEX\docs\reference\_apex_deck_v4_assets")
ASSETS.mkdir(parents=True, exist_ok=True)

# ─── Palette ──────────────────────────────────────────────────────────────
BG          = RGBColor(0xFA, 0xFA, 0xFA)
HEAD_NAVY   = RGBColor(0x1A, 0x1F, 0x2E)
DEEP_NAVY   = RGBColor(0x0E, 0x14, 0x23)
GREEN       = RGBColor(0x86, 0xBC, 0x25)
GREEN_DK    = RGBColor(0x5A, 0x83, 0x0F)
MS_BLUE     = RGBColor(0x00, 0x78, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
STRONG_INK  = RGBColor(0x1A, 0x1F, 0x2E)
MUTED       = RGBColor(0x47, 0x55, 0x69)
SOFT_LINE   = RGBColor(0xE2, 0xE8, 0xF0)
PALE_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_TEXT  = RGBColor(0xCB, 0xD5, 0xE1)
ICE         = RGBColor(0xBF, 0xDB, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════
# Pillow image generation — hero backgrounds
# ═══════════════════════════════════════════════════════════════════════════

def _diagonal_gradient(size, c0, c1, *, angle_deg=30):
    """Linear gradient from c0 (top-left-ish) to c1 (bottom-right-ish)."""
    import math
    w, h = size
    img = Image.new("RGB", size)
    px = img.load()
    rad = math.radians(angle_deg)
    dx, dy = math.cos(rad), math.sin(rad)
    # Project each pixel onto the gradient axis, normalize 0..1
    proj_min = 0
    proj_max = w * abs(dx) + h * abs(dy)
    for y in range(h):
        for x in range(w):
            t = (x * dx + y * dy) / proj_max
            t = max(0.0, min(1.0, t))
            r = int(c0[0] + (c1[0] - c0[0]) * t)
            g = int(c0[1] + (c1[1] - c0[1]) * t)
            b = int(c0[2] + (c1[2] - c0[2]) * t)
            px[x, y] = (r, g, b)
    return img


def _radial_dots_overlay(size, *, dot_color=(134, 188, 37), spacing=80, radius=2, alpha=70):
    """Subtle dot grid overlay (Deloitte green accent)."""
    w, h = size
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    for y in range(spacing // 2, h, spacing):
        for x in range(spacing // 2, w, spacing):
            draw.ellipse([x - radius, y - radius, x + radius, y + radius],
                         fill=(*dot_color, alpha))
    return img


def _diagonal_lines_overlay(size, *, line_color=(255, 255, 255), spacing=120,
                             alpha=20, line_width=1):
    """Faint diagonal pinstripe overlay for tech feel."""
    w, h = size
    img = Image.new("RGBA", size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # 45-degree lines
    for offset in range(-h, w + h, spacing):
        draw.line([(offset, 0), (offset + h, h)],
                  fill=(*line_color, alpha), width=line_width)
    return img


def gen_hero_title(path: Path) -> None:
    """Slide-1 hero: navy gradient + green dot pattern + diagonal accent."""
    size = (1920, 1080)
    img = _diagonal_gradient(size, (14, 20, 35), (30, 41, 59), angle_deg=20)
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, _diagonal_lines_overlay(size, spacing=140, alpha=18))
    img = Image.alpha_composite(img, _radial_dots_overlay(size, spacing=80, radius=2, alpha=80))
    # Green diagonal accent in the corner
    accent = Image.new("RGBA", size, (0, 0, 0, 0))
    a_draw = ImageDraw.Draw(accent)
    a_draw.polygon([(0, 0), (380, 0), (0, 380)], fill=(134, 188, 37, 200))
    a_draw.polygon([(0, 0), (260, 0), (0, 260)], fill=(134, 188, 37, 255))
    img = Image.alpha_composite(img, accent)
    img.convert("RGB").save(path, "PNG", optimize=True)


def gen_hero_warm(path: Path) -> None:
    """CFMP spotlight: warm gradient evoking retail / customer experience."""
    size = (1920, 1080)
    img = _diagonal_gradient(size, (255, 218, 165), (255, 153, 102), angle_deg=35)
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, _radial_dots_overlay(size, dot_color=(180, 80, 30),
                                                          spacing=100, radius=3, alpha=60))
    # Soft blur for depth
    img.convert("RGB").filter(ImageFilter.GaussianBlur(radius=2)).save(path, "PNG", optimize=True)


def gen_hero_demo(path: Path) -> None:
    """Demo URL slide: tech-blue gradient + circuit-like pattern."""
    size = (1920, 1080)
    img = _diagonal_gradient(size, (10, 25, 60), (0, 120, 212), angle_deg=20)
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, _diagonal_lines_overlay(size, spacing=80,
                                                              line_color=(255, 255, 255),
                                                              alpha=15))
    img = Image.alpha_composite(img, _radial_dots_overlay(size, dot_color=(255, 255, 255),
                                                          spacing=60, radius=1, alpha=110))
    img.convert("RGB").save(path, "PNG", optimize=True)


def gen_hero_closing(path: Path) -> None:
    """Closing slide: green-to-navy diagonal momentum gradient."""
    size = (1920, 1080)
    img = _diagonal_gradient(size, (134, 188, 37), (14, 20, 35), angle_deg=15)
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, _diagonal_lines_overlay(size, spacing=100, alpha=25))
    img.convert("RGB").save(path, "PNG", optimize=True)


def gen_subtle_blueprint(path: Path) -> None:
    """Generic subtle blueprint pattern for half-bleed use."""
    size = (1280, 1080)
    img = Image.new("RGB", size, (240, 246, 252))
    img = img.convert("RGBA")
    # Light navy dot grid
    img = Image.alpha_composite(img, _radial_dots_overlay(size,
                                                          dot_color=(60, 80, 130),
                                                          spacing=40, radius=1, alpha=50))
    # Sparse green dots
    img = Image.alpha_composite(img, _radial_dots_overlay(size,
                                                          dot_color=(134, 188, 37),
                                                          spacing=120, radius=2, alpha=120))
    img.convert("RGB").save(path, "PNG", optimize=True)


# Generate all five
print("Generating hero backgrounds…")
gen_hero_title(ASSETS / "hero_title.png")
gen_hero_warm(ASSETS / "hero_warm.png")
gen_hero_demo(ASSETS / "hero_demo.png")
gen_hero_closing(ASSETS / "hero_closing.png")
gen_subtle_blueprint(ASSETS / "blueprint.png")
print(f"  → wrote 5 assets in {ASSETS}")


# ═══════════════════════════════════════════════════════════════════════════
# PowerPoint primitives
# ═══════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, *, fill, line=None, corner_radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    if corner_radius:
        shp.adjustments[0] = corner_radius
    return shp


def add_oval(slide, left, top, width, height, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def add_hexagon(slide, left, top, width, height, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    return shp


def add_arrow(slide, left, top, width, height, *, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def set_text(shape, text, *, size=10, bold=False, color=STRONG_INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36_000); tf.margin_right = Emu(36_000)
    tf.margin_top = Emu(18_000);  tf.margin_bottom = Emu(18_000)
    lines = [text] if isinstance(text, str) else list(text)
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, **kw)
    return box


def header_minimal(slide, title: str, subtitle: str = ""):
    """Lean exec-deck header — small green accent + large title."""
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.20), Inches(0.20),
             fill=GREEN)
    add_textbox(slide, Inches(0.85), Inches(0.30), Inches(11), Inches(0.55),
                title, size=28, bold=True, color=STRONG_INK)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(0.90), Inches(11.5), Inches(0.40),
                    subtitle, size=13, color=MUTED, italic=True)


def page_footer(slide, slide_no: int, total: int):
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.02),
             fill=SOFT_LINE)
    add_textbox(slide, Inches(0.5), Inches(7.10), Inches(9), Inches(0.30),
                "APEX  ·  Deloitte Microsoft Technology & Services Practice",
                size=9, color=MUTED)
    add_textbox(slide, Inches(11.00), Inches(7.10), Inches(2.0), Inches(0.30),
                f"{slide_no} / {total}", size=9, bold=True, color=GREEN,
                align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# Custom icon builders — composed from primitives, look professional
# ═══════════════════════════════════════════════════════════════════════════

def icon_factory(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Factory icon — stacked rectangles + chimney triangle."""
    # Main building
    add_rect(slide, x, y + h * 0.45, w, h * 0.55, fill=color)
    # Roof zigzag (3 triangles)
    for i in range(3):
        tri_x = x + w * (i / 3)
        tri_w = w / 3
        add_rect(slide, tri_x, y + h * 0.30, tri_w * 0.95, h * 0.20, fill=color)
    # Chimney
    add_rect(slide, x + w * 0.15, y + h * 0.05, w * 0.12, h * 0.30, fill=color)
    # Smoke (lighter circle)
    add_oval(slide, x + w * 0.08, y, w * 0.20, h * 0.20,
             fill=RGBColor(0xCB, 0xD5, 0xE1))


def icon_shopping(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Shopping bag icon."""
    # Handles
    add_oval(slide, x + w * 0.18, y, w * 0.20, h * 0.40, fill=BG)
    add_oval(slide, x + w * 0.18, y + h * 0.05, w * 0.20, h * 0.35, fill=color)
    add_oval(slide, x + w * 0.62, y, w * 0.20, h * 0.40, fill=BG)
    add_oval(slide, x + w * 0.62, y + h * 0.05, w * 0.20, h * 0.35, fill=color)
    # Main body (trapezoid-like; using rectangle)
    add_rect(slide, x + w * 0.05, y + h * 0.30, w * 0.90, h * 0.70,
             fill=color, corner_radius=0.15)


def icon_bars(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Bar chart icon (Finance) — 3 ascending bars."""
    heights = [0.4, 0.7, 1.0]
    for i, ratio in enumerate(heights):
        bar_x = x + w * (0.1 + i * 0.30)
        bar_h = h * 0.85 * ratio
        add_rect(slide, bar_x, y + h - bar_h, w * 0.22, bar_h, fill=color)


def icon_shield(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Shield icon (Risk)."""
    # Top rectangle + bottom rounded
    add_rect(slide, x + w * 0.10, y, w * 0.80, h * 0.55, fill=color)
    add_oval(slide, x + w * 0.10, y + h * 0.20, w * 0.80, h * 0.80, fill=color)
    # Inner mark (checkmark via diamond)
    add_rect(slide, x + w * 0.40, y + h * 0.30, w * 0.20, h * 0.20,
             fill=BG, corner_radius=0.3)


def icon_headset(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Headset icon (CX)."""
    # Arc band — approximated by a thick oval with cutout via overlay
    add_oval(slide, x + w * 0.10, y + h * 0.05, w * 0.80, h * 0.80, fill=color)
    # Cut interior
    add_oval(slide, x + w * 0.22, y + h * 0.17, w * 0.56, h * 0.66, fill=BG)
    # Bottom mask to keep U-shape only
    add_rect(slide, x, y + h * 0.55, w, h * 0.50, fill=BG)
    # Earpieces
    add_oval(slide, x + w * 0.05, y + h * 0.45, w * 0.20, h * 0.30, fill=color)
    add_oval(slide, x + w * 0.75, y + h * 0.45, w * 0.20, h * 0.30, fill=color)
    # Mic dot
    add_oval(slide, x + w * 0.55, y + h * 0.72, w * 0.10, h * 0.15, fill=color)


def icon_leaf(slide, x, y, w, h, *, color=GREEN):
    """Leaf icon (ESG)."""
    # Teardrop approximated by oval rotated visually + small base
    add_oval(slide, x + w * 0.15, y + h * 0.05, w * 0.70, h * 0.85, fill=color)
    # Stem
    add_rect(slide, x + w * 0.46, y + h * 0.80, w * 0.08, h * 0.20, fill=GREEN_DK)


def icon_pin(slide, x, y, w, h, *, color=GREEN):
    """Map pin icon (CFMP / wayfinding)."""
    # Pinhead
    add_oval(slide, x + w * 0.15, y, w * 0.70, h * 0.65, fill=color)
    # Inner dot
    add_oval(slide, x + w * 0.36, y + h * 0.18, w * 0.28, h * 0.28, fill=BG)
    # Pin tail (triangle approximated by narrow rotated rect)
    # Use a regular triangle-ish via three small rects? Simplest: a small pointer rect.
    add_rect(slide, x + w * 0.42, y + h * 0.55, w * 0.16, h * 0.45, fill=color)


def icon_pillar(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Pillar / column icon (CORE)."""
    # Top cap
    add_rect(slide, x + w * 0.10, y, w * 0.80, h * 0.10, fill=color)
    # Shaft
    add_rect(slide, x + w * 0.25, y + h * 0.10, w * 0.50, h * 0.80, fill=color)
    # Base
    add_rect(slide, x + w * 0.05, y + h * 0.90, w * 0.90, h * 0.10, fill=color)


def icon_cloud(slide, x, y, w, h, *, color=MS_BLUE):
    """Cloud icon (PROFILES) — three connected ovals."""
    add_oval(slide, x + w * 0.05, y + h * 0.40, w * 0.40, h * 0.55, fill=color)
    add_oval(slide, x + w * 0.25, y + h * 0.20, w * 0.50, h * 0.65, fill=color)
    add_oval(slide, x + w * 0.50, y + h * 0.40, w * 0.45, h * 0.50, fill=color)
    # Bottom mask to flatten the cloud's base
    add_rect(slide, x, y + h * 0.78, w, h * 0.22, fill=BG)


def icon_box_stack(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Box stack icon (PACKS)."""
    # 3 stacked boxes
    add_rect(slide, x + w * 0.10, y + h * 0.55, w * 0.55, h * 0.42, fill=color)
    add_rect(slide, x + w * 0.30, y + h * 0.25, w * 0.55, h * 0.32, fill=GREEN)
    add_rect(slide, x + w * 0.05, y + h * 0.05, w * 0.50, h * 0.22, fill=color)


def icon_envelope(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Envelope icon (ENVELOPES)."""
    add_rect(slide, x + w * 0.05, y + h * 0.20, w * 0.90, h * 0.60, fill=color)
    # Top flap (triangle approximated by a thin rect rotated — using a triangle shape via MSO_SHAPE.ISOCELES_TRIANGLE)
    flap = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x + w * 0.05, y + h * 0.18, w * 0.90, h * 0.40)
    flap.fill.solid(); flap.fill.fore_color.rgb = GREEN
    flap.line.fill.background()
    flap.rotation = 180


def icon_chain(slide, x, y, w, h, *, color=HEAD_NAVY):
    """Chain links icon (LEDGER) — 3 interlocked ovals."""
    add_oval(slide, x + w * 0.05, y + h * 0.30, w * 0.40, h * 0.40,
             fill=BG, line=color)
    add_oval(slide, x + w * 0.30, y + h * 0.30, w * 0.40, h * 0.40,
             fill=BG, line=color)
    add_oval(slide, x + w * 0.55, y + h * 0.30, w * 0.40, h * 0.40,
             fill=BG, line=color)
    # Note: line + fill=BG gives an outlined ring appearance


def icon_check(slide, x, y, w, h, *, color=GREEN):
    add_oval(slide, x, y, w, h, fill=color)
    add_textbox(slide, x, y, w, h, "✓", size=int(h / 914400 * 72 * 0.7),
                bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def icon_x(slide, x, y, w, h, *, color=RGBColor(0xA0, 0x35, 0x35)):
    add_oval(slide, x, y, w, h, fill=color)
    add_textbox(slide, x, y, w, h, "✕", size=int(h / 914400 * 72 * 0.6),
                bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# Build the deck
# ═══════════════════════════════════════════════════════════════════════════

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 15

    # ── Slide 1: Title hero — full-bleed gradient ─────────────────────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_title.png"),
                         Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Top-right version chip
    add_rect(s, Inches(10.6), Inches(0.5), Inches(2.3), Inches(0.4),
             fill=GREEN, corner_radius=0.5)
    add_textbox(s, Inches(10.6), Inches(0.5), Inches(2.3), Inches(0.4),
                "DMTSP  ·  v4  ·  2026", size=10, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title block (left-aligned, lower-third)
    add_textbox(s, Inches(0.7), Inches(3.2), Inches(11), Inches(1.2),
                "APEX", size=108, bold=True, color=WHITE, font="Calibri")
    # Green accent bar
    add_rect(s, Inches(0.7), Inches(4.6), Inches(0.6), Inches(0.10), fill=GREEN)
    add_textbox(s, Inches(0.7), Inches(4.8), Inches(11), Inches(0.6),
                "Deloitte's agentic-AI delivery framework",
                size=28, color=WHITE, font="Calibri")
    add_textbox(s, Inches(0.7), Inches(5.45), Inches(11), Inches(0.4),
                "Walkthrough for DMTSP  ·  what it is  ·  how to demo  ·  how to make it real",
                size=14, color=LIGHT_TEXT, italic=True)
    # URL chip near bottom
    add_rect(s, Inches(0.7), Inches(6.3), Inches(10.2), Inches(0.55),
             fill=RGBColor(0x00, 0x00, 0x00), corner_radius=0.2)
    add_textbox(s, Inches(0.7), Inches(6.3), Inches(10.2), Inches(0.55),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=12, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 1, TOTAL)

    # ── Slide 2: APEX in one sentence — visual sandwich diagram ─────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "APEX in one sentence",
                   "The engagement layer cloud vendors don't ship.")
    # Big quote, centered, left
    add_textbox(s, Inches(0.6), Inches(2.0), Inches(0.5), Inches(0.9),
                "“", size=120, bold=True, color=GREEN)
    add_textbox(s, Inches(0.6), Inches(2.9), Inches(6.8), Inches(2.2),
                ["The engagement layer",
                 "cloud vendors", "don't ship."],
                size=42, bold=True, color=STRONG_INK)

    # Sandwich diagram on the right
    diag_x = Inches(7.7); diag_y = Inches(1.9); diag_w = Inches(5.2)
    # Top layer — Client outcome
    add_rect(s, diag_x, diag_y, diag_w, Inches(0.85), fill=HEAD_NAVY,
             corner_radius=0.15)
    add_textbox(s, diag_x, diag_y, diag_w, Inches(0.85),
                "Client outcome", size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Arrow down
    arr1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              diag_x + diag_w * 0.45, diag_y + Inches(0.9),
                              diag_w * 0.1, Inches(0.45))
    arr1.fill.solid(); arr1.fill.fore_color.rgb = SOFT_LINE
    arr1.line.fill.background()
    # APEX layer
    add_rect(s, diag_x, diag_y + Inches(1.4), diag_w, Inches(1.5), fill=GREEN,
             corner_radius=0.15)
    add_textbox(s, diag_x, diag_y + Inches(1.4), diag_w, Inches(0.5),
                "APEX  —  Deloitte's engagement layer",
                size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, diag_x, diag_y + Inches(1.85), diag_w, Inches(1.0),
                ["Industry Packs  ·  Virtual Views",
                 "WORM LEDGER  ·  Service Envelopes",
                 "10-asset bundle structure"],
                size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow down 2
    arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              diag_x + diag_w * 0.45, diag_y + Inches(3.05),
                              diag_w * 0.1, Inches(0.45))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = SOFT_LINE
    arr2.line.fill.background()
    # Cloud runtime
    add_rect(s, diag_x, diag_y + Inches(3.55), diag_w, Inches(0.85),
             fill=MS_BLUE, corner_radius=0.15)
    add_textbox(s, diag_x, diag_y + Inches(3.55), diag_w, Inches(0.85),
                "Foundry  ·  Vertex  ·  Bedrock",
                size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, diag_x, diag_y + Inches(4.45), diag_w, Inches(0.35),
                "(cloud-vendor agent runtimes)",
                size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 2, TOTAL)

    # ── Slide 3: Without / With APEX — split-screen ────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Without APEX  ·  With APEX",
                   "Same Foundry runtime. Different deal shape.")

    # Two panels
    def panel(x, w, tone, headline, items):
        col_fill = RGBColor(0xFD, 0xF4, 0xF4) if tone == "bad" else RGBColor(0xF1, 0xF8, 0xE9)
        accent = RGBColor(0xA0, 0x35, 0x35) if tone == "bad" else GREEN
        add_rect(s, x, Inches(1.85), w, Inches(5.05), fill=col_fill,
                 line=accent, corner_radius=0.04)
        # Header band
        add_rect(s, x, Inches(1.85), w, Inches(0.6), fill=accent)
        add_textbox(s, x, Inches(1.85), w, Inches(0.6), headline,
                    size=15, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Items
        for i, (label, desc) in enumerate(items):
            iy = Inches(2.65 + i * 1.05)
            # Bullet dot
            add_oval(s, x + Inches(0.30), iy + Inches(0.15),
                     Inches(0.16), Inches(0.16), fill=accent)
            # Label + body
            add_textbox(s, x + Inches(0.60), iy, w - Inches(0.70), Inches(0.40),
                        label, size=14, bold=True, color=STRONG_INK)
            add_textbox(s, x + Inches(0.60), iy + Inches(0.40),
                        w - Inches(0.70), Inches(0.50),
                        desc, size=10.5, color=MUTED)

    panel(Inches(0.5), Inches(6.1), "bad", "WITHOUT", [
        ("Staff augmentation", "Project competes with what client could build alone."),
        ("Bespoke discovery", "Every engagement re-invents the data model + SOW."),
        ("Margin compression", "Deal shape resembles labor sales."),
        ("Fragile audit story", "No replay; trust is conversational, not cryptographic."),
    ])
    panel(Inches(6.85), Inches(6.0), "good", "WITH APEX", [
        ("Fixed-fee Packs", "Industry Solution Pack scoped at sub-tier."),
        ("Canonical scenarios", "10-asset bundle ships pre-built playbooks."),
        ("WORM LEDGER", "14-field hash-chained audit row per decision."),
        ("Replatform = config", "Same Pack on APEX-M / G / A."),
    ])
    page_footer(s, 3, TOTAL)

    # ── Slide 4: Five Pieces — radial diagram ──────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Five pieces. Every engagement.",
                   "Core  ·  Profiles  ·  Packs  ·  Envelopes  ·  LEDGER")

    pieces = [
        (icon_pillar,    "CORE",       "Cloud-neutral IP",   HEAD_NAVY),
        (icon_cloud,     "PROFILES",   "M  /  G  /  A",      MS_BLUE),
        (icon_box_stack, "PACKS",      "7 in the catalog",   HEAD_NAVY),
        (icon_envelope,  "ENVELOPES",  "5 sellable tiers",   HEAD_NAVY),
        (icon_chain,     "LEDGER",     "WORM hash chain",    GREEN_DK),
    ]
    cw = Inches(2.40); gap = Inches(0.15); x0 = Inches(0.55)
    for i, (icon_fn, label, sub, color) in enumerate(pieces):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, Inches(1.9), cw, Inches(4.7),
                 fill=WHITE, line=SOFT_LINE, corner_radius=0.05)
        # Green top band
        add_rect(s, cx, Inches(1.9), cw, Inches(0.10), fill=GREEN)
        # Icon area (centered)
        ix = cx + Inches(0.6); iy_ = Inches(2.30)
        iw = Inches(1.20); ih = Inches(1.20)
        icon_fn(s, ix, iy_, iw, ih, color=color)
        # Number
        add_textbox(s, cx, Inches(3.70), cw, Inches(0.35),
                    f"0{i+1}", size=12, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER)
        # Label
        add_textbox(s, cx, Inches(4.10), cw, Inches(0.50),
                    label, size=22, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Sub
        add_textbox(s, cx, Inches(4.75), cw, Inches(1.5),
                    sub, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    page_footer(s, 4, TOTAL)

    # ── Slide 5: APEX Family M/G/A — three clouds ──────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same core. Different cloud.",
                   "Three implementations of one 14-interface contract.")

    families = [
        ("APEX-M", "MICROSOFT  /  AZURE", "GA",            MS_BLUE),
        ("APEX-G", "GOOGLE  /  GCP",      "BETA Q3 FY27", RGBColor(0x0F, 0x9D, 0x58)),
        ("APEX-A", "AMAZON  /  AWS",      "BETA Q4 FY27", RGBColor(0xFF, 0x99, 0x00)),
    ]
    cw = Inches(4.0); gap = Inches(0.2); x0 = Inches(0.55)
    for i, (tag, name, status, color) in enumerate(families):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Card
        add_rect(s, cx, Inches(1.9), cw, Inches(4.7),
                 fill=WHITE, line=SOFT_LINE, corner_radius=0.05)
        # Top color stripe
        add_rect(s, cx, Inches(1.9), cw, Inches(0.5), fill=color)
        add_textbox(s, cx, Inches(1.9), cw, Inches(0.5), tag,
                    size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Cloud illustration
        icon_cloud(s, cx + Inches(0.95), Inches(2.7),
                   Inches(2.10), Inches(1.5), color=color)
        # Name
        add_textbox(s, cx, Inches(4.45), cw, Inches(0.4),
                    name, size=13, color=MUTED, align=PP_ALIGN.CENTER)
        # Status pill
        add_rect(s, cx + Inches(0.9), Inches(4.95),
                 cw - Inches(1.8), Inches(0.45),
                 fill=color, corner_radius=0.5)
        add_textbox(s, cx + Inches(0.9), Inches(4.95),
                    cw - Inches(1.8), Inches(0.45),
                    status, size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        bullets = {
            "APEX-M": ["Foundry + Azure OpenAI", "OneLake + Fabric", "Entra · Purview · Maps Creator"],
            "APEX-G": ["Vertex AI + Gemini", "BigQuery + BigLake", "Cloud Identity · Dataplex"],
            "APEX-A": ["Bedrock (Claude / Nova)", "S3 + Lake Formation", "IAM IC · Location Service"],
        }[tag]
        for j, b in enumerate(bullets):
            add_textbox(s, cx + Inches(0.25), Inches(5.55 + j * 0.35),
                        cw - Inches(0.50), Inches(0.32),
                        "·  " + b, size=11, color=STRONG_INK)
    page_footer(s, 5, TOTAL)

    # ── Slide 6: Seven Industry Packs — hex catalog ────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Seven Industry Packs",
                   "Same 10-asset bundle. Different industry per Pack.")

    packs = [
        (icon_factory,  "Manufacturing",        HEAD_NAVY, "Plant / Ops"),
        (icon_shopping, "Retail Merch",          HEAD_NAVY, "Chief Merchant"),
        (icon_bars,     "Finance",               HEAD_NAVY, "Controller · CFO"),
        (icon_shield,   "Risk",                  HEAD_NAVY, "SOX PMO · CRO"),
        (icon_headset,  "Customer Experience",   HEAD_NAVY, "Care Ops"),
        (icon_leaf,     "ESG",                   GREEN,     "Sustainability"),
        (icon_pin,      "CFMP",                  GREEN,     "CMO · CX-VP"),
    ]
    # Lay out as 4 + 3 (centered second row)
    cw = Inches(2.8); ch = Inches(2.30); gap_x = Inches(0.30); gap_y = Inches(0.25)

    def pack_card(cx, cy, icon_fn, name, color, persona, *, highlight=False):
        # Card body
        add_rect(s, cx, cy, cw, ch,
                 fill=WHITE, line=(GREEN if highlight else SOFT_LINE),
                 corner_radius=0.05)
        if highlight:
            # Star ribbon corner
            add_rect(s, cx + cw - Inches(0.9), cy, Inches(0.9), Inches(0.32),
                     fill=GREEN)
            add_textbox(s, cx + cw - Inches(0.9), cy, Inches(0.9), Inches(0.32),
                        "★ LIVE", size=9, bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Icon
        icon_fn(s, cx + cw * 0.35, cy + Inches(0.20), cw * 0.30, Inches(0.90),
                color=color)
        # Name
        add_textbox(s, cx, cy + Inches(1.18), cw, Inches(0.40),
                    name, size=15, bold=True, color=STRONG_INK,
                    align=PP_ALIGN.CENTER)
        # Persona
        add_textbox(s, cx, cy + Inches(1.62), cw, Inches(0.4),
                    persona, size=11, color=MUTED, italic=True,
                    align=PP_ALIGN.CENTER)

    # Row 1 (4 cards)
    y_row1 = Inches(1.85)
    x_row1 = Inches(0.30)
    for i in range(4):
        cx = x_row1 + Inches(i * (cw / Inches(1) + gap_x / Inches(1)))
        icon_fn, name, color, persona = packs[i]
        pack_card(cx, y_row1, icon_fn, name, color, persona)
    # Row 2 (3 cards centered)
    y_row2 = y_row1 + ch + gap_y
    # Center 3 cards across 13.33 width:  3 * cw + 2 * gap = 3 * 2.8 + 2 * 0.3 = 9.0 → start at (13.33-9)/2 = 2.165
    x_row2 = Inches(2.165)
    for i in range(3):
        cx = x_row2 + Inches(i * (cw / Inches(1) + gap_x / Inches(1)))
        icon_fn, name, color, persona = packs[4 + i]
        is_cfmp = name == "CFMP"
        pack_card(cx, y_row2, icon_fn, name, color, persona, highlight=is_cfmp)
    page_footer(s, 6, TOTAL)

    # ── Slide 7: 10-Asset Bundle — neat grid ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Same engineering structure. Every Pack.",
                   "Ten slots. Identical contract. New Packs land predictably.")

    assets = [
        ("01", "VV manifests",       "Federated views"),
        ("02", "Scenarios",          "End-to-end outcomes"),
        ("03", "Source adapters",    "System connectors"),
        ("04", "Adaptive cards",     "HITL surfaces"),
        ("05", "Persona map",        "Role bindings"),
        ("06", "Demo data",          "Laptop-runnable"),
        ("07", "BVA worksheet",      "ROI calculator"),
        ("08", "Sample SOW",         "Tier templates"),
        ("09", "Acceptance tests",   "CI-green"),
        ("10", "Runbook + training", "Operate-ready"),
    ]
    # 5 cols × 2 rows
    cw = Inches(2.40); ch = Inches(2.10); gx = Inches(0.10); gy = Inches(0.10)
    x0 = Inches(0.55); y0 = Inches(1.95)
    for i, (num, name, sub) in enumerate(assets):
        col = i % 5; row = i // 5
        cx = x0 + Inches(col * (cw / Inches(1) + gx / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + gy / Inches(1)))
        # Card
        add_rect(s, cx, cy, cw, ch,
                 fill=WHITE, line=SOFT_LINE, corner_radius=0.05)
        # Number circle
        add_oval(s, cx + Inches(0.20), cy + Inches(0.18),
                 Inches(0.65), Inches(0.65), fill=GREEN)
        add_textbox(s, cx + Inches(0.20), cy + Inches(0.18),
                    Inches(0.65), Inches(0.65),
                    num, size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name
        add_textbox(s, cx + Inches(0.95), cy + Inches(0.25),
                    cw - Inches(1.05), Inches(0.40),
                    name, size=13, bold=True, color=STRONG_INK)
        # Sub
        add_textbox(s, cx + Inches(0.95), cy + Inches(0.65),
                    cw - Inches(1.05), Inches(0.35),
                    sub, size=10, color=MUTED)
        # Decorative bar
        add_rect(s, cx + Inches(0.20), cy + Inches(1.35),
                 cw - Inches(0.40), Inches(0.05), fill=PALE_GRAY)
    page_footer(s, 7, TOTAL)

    # ── Slide 8: CFMP spotlight — warm hero + journey arrows ────────────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_warm.png"),
                         Inches(0), Inches(0), SLIDE_W, Inches(3.4))
    # Dark navy block top-left for title legibility
    add_rect(s, Inches(0), Inches(0), Inches(6.5), Inches(3.4), fill=HEAD_NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(3.4), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(0.6), Inches(5.5), Inches(0.7),
                "Spotlight: CFMP",
                size=42, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                "Customer Focused Merchandise Pack",
                size=18, color=ICE)
    add_textbox(s, Inches(0.5), Inches(2.05), Inches(5.5), Inches(0.4),
                "v0.2  ·  Phase 1 LIVE  ·  the live worked example",
                size=12, color=LIGHT_TEXT, italic=True)
    add_rect(s, Inches(0.5), Inches(2.65), Inches(2.8), Inches(0.45),
             fill=GREEN, corner_radius=0.5)
    add_textbox(s, Inches(0.5), Inches(2.65), Inches(2.8), Inches(0.45),
                "Pack Lite · $150–250K · 4–6 wks",
                size=10, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom area — journey ribbon
    add_rect(s, Inches(0), Inches(3.4), SLIDE_W, Inches(4.1), fill=BG)
    add_textbox(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.45),
                "The customer-journey spine",
                size=18, bold=True, color=STRONG_INK)
    add_textbox(s, Inches(0.5), Inches(4.10), Inches(12.3), Inches(0.35),
                "Four moments. Four phases. Eighteen scenarios mapped.",
                size=12, color=MUTED, italic=True)
    # 4 phase arrows
    phases = [
        ("CHOOSE",   "What should I buy?",   MS_BLUE),
        ("SELECT",   "Where is it?",         GREEN),
        ("BUY",      "How do I pay?",        RGBColor(0xFF, 0x99, 0x00)),
        ("SERVICES", "What's next?",         RGBColor(0xB4, 0x42, 0xCC)),
    ]
    aw = Inches(3.0); gap = Inches(0.10); x0 = Inches(0.5)
    for i, (label, sub, color) in enumerate(phases):
        cx = x0 + Inches(i * (aw / Inches(1) + gap / Inches(1)))
        if i == 3:
            # Last one — rounded rectangle (no arrow tail beyond slide)
            add_rect(s, cx, Inches(4.7), aw, Inches(1.6), fill=color,
                     corner_radius=0.12)
        else:
            add_arrow(s, cx, Inches(4.7), aw, Inches(1.6), fill=color)
        # Label
        add_textbox(s, cx, Inches(4.85), aw - Inches(0.4), Inches(0.55),
                    label, size=22, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, cx, Inches(5.45), aw - Inches(0.4), Inches(0.4),
                    sub, size=12, color=WHITE, italic=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Below ribbon — concrete capability counts
    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                "18 scenarios  ·  6 new Virtual Views  ·  proposed Interface #15 (Maps & Wayfinding)  ·  new `customer` persona type  ·  $50 cart consent gate",
                size=10.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 8, TOTAL)

    # ── Slide 9: Demo URL — bold hero ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_demo.png"),
                         Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Dark overlay for text legibility (full coverage with 50% transparency via blend - we'll use solid dark)
    add_rect(s, Inches(0), Inches(2.0), SLIDE_W, Inches(3.6), fill=DEEP_NAVY)
    add_rect(s, Inches(0), Inches(2.0), Inches(0.20), Inches(3.6), fill=GREEN)
    add_textbox(s, Inches(0.5), Inches(2.2), Inches(12.5), Inches(0.7),
                "Open the demo. Talk through it.",
                size=40, bold=True, color=WHITE)
    add_textbox(s, Inches(0.5), Inches(2.95), Inches(12.5), Inches(0.45),
                "Eight surfaces  ·  ten minutes  ·  let the agent answer for you.",
                size=15, color=LIGHT_TEXT, italic=True)
    # URL on green bar
    add_rect(s, Inches(0.5), Inches(3.85), Inches(12.33), Inches(0.95),
             fill=GREEN, corner_radius=0.05)
    add_textbox(s, Inches(0.5), Inches(3.85), Inches(12.33), Inches(0.95),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 8 step pills along bottom
    steps = ["1  Demo /", "2  “Where is X?”", "3  Dietary chip", "4  Cart click",
             "5  /architecture", "6  Details ⓘ", "7  APEX value", "8  Tour flyout"]
    cw = Inches(1.50); gap = Inches(0.05); x0 = Inches(0.55)
    for i, t in enumerate(steps):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        add_rect(s, cx, Inches(6.0), cw, Inches(0.6),
                 fill=WHITE, corner_radius=0.4)
        add_textbox(s, cx, Inches(6.0), cw, Inches(0.6),
                    t, size=10, bold=True, color=DEEP_NAVY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 9, TOTAL)

    # ── Slide 10: Engagement Ladder — ascending staircase diagram ──────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Engagement ladder",
                   "Five tiers. Stackable. Each tier earns the next.")
    tiers = [
        ("T1", "FOUNDATION",  "$400–700K",       "6–10 wks"),
        ("T2", "PACK",        "$100K–$3.5M",     "4 wks – 9 mo"),
        ("T3", "SCENARIO",    "$150–400K",       "4–8 wks"),
        ("T4", "CUSTOM",      "T&M",             "Sprints"),
        ("T5", "OPERATE",     "$/mo",            "Continuous"),
    ]
    # Ascending stair: each step gets taller + further right
    base_x = Inches(0.7); base_y = Inches(6.4); step_w = Inches(2.45); step_h_base = Inches(0.55)
    for i, (tag, name, price, time) in enumerate(tiers):
        h = step_h_base + Inches(i * 0.20)
        x = base_x + Inches(i * (step_w / Inches(1) - 0.05))
        y = base_y - h
        # Step block
        add_rect(s, x, y, step_w, h, fill=HEAD_NAVY,
                 corner_radius=0.05)
        # Green leading-edge stripe
        add_rect(s, x, y, Inches(0.10), h, fill=GREEN)
        # Tag
        add_textbox(s, x + Inches(0.2), y + Inches(0.05),
                    Inches(0.6), Inches(0.4),
                    tag, size=14, bold=True, color=GREEN,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Name
        add_textbox(s, x + Inches(0.85), y + Inches(0.05),
                    step_w - Inches(1.0), Inches(0.4),
                    name, size=13, bold=True, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Price + time below name (only if step is tall enough)
        if i > 0:
            add_textbox(s, x + Inches(0.2), y + Inches(0.50),
                        step_w - Inches(0.3), Inches(0.3),
                        price, size=10, color=ICE)
            add_textbox(s, x + Inches(0.2), y + Inches(0.78),
                        step_w - Inches(0.3), Inches(0.3),
                        time, size=9, color=LIGHT_TEXT, italic=True)
    # Up-arrow showing growth direction (left of staircase)
    arr = s.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                             Inches(0.2), Inches(2.0),
                             Inches(0.4), Inches(4.3))
    arr.fill.solid(); arr.fill.fore_color.rgb = GREEN
    arr.line.fill.background()
    add_textbox(s, Inches(0.0), Inches(2.0), Inches(0.7), Inches(0.6),
                "GROW", size=10, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Right-side callouts (specific Pack Lite emphasis)
    add_rect(s, Inches(9.0), Inches(1.95), Inches(4.0), Inches(2.6),
             fill=GREEN, corner_radius=0.05)
    add_textbox(s, Inches(9.2), Inches(2.10), Inches(3.6), Inches(0.4),
                "The standard wedge",
                size=14, bold=True, color=WHITE)
    add_textbox(s, Inches(9.2), Inches(2.55), Inches(3.6), Inches(0.6),
                "Pack Lite  ·  T2",
                size=26, bold=True, color=WHITE)
    add_textbox(s, Inches(9.2), Inches(3.25), Inches(3.6), Inches(0.5),
                "$150–250K  ·  4–6 weeks",
                size=14, color=WHITE)
    add_textbox(s, Inches(9.2), Inches(3.85), Inches(3.6), Inches(0.6),
                "BVA + DCIF funding · client-funded discovery → fixed-fee pilot",
                size=10, color=ICE, italic=True)
    page_footer(s, 10, TOTAL)

    # ── Slide 11: Independence — Shield + DO / DO NOT ──────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "Independence  —  non-negotiable",
                   "Deloitte audits Microsoft. SEC rules shape funding paths.")

    # Big shield icon on left
    icon_shield(s, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.5),
                color=HEAD_NAVY)
    # Add ✓ inside shield
    add_textbox(s, Inches(0.8), Inches(3.2), Inches(3.2), Inches(2.0),
                "✓", size=180, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # DO column
    do_items = ["BVA", "DCIF", "ISV Marketplace burndown", "SI Teaming POC", "T&M  ·  Client direct"]
    add_textbox(s, Inches(4.6), Inches(1.95), Inches(4), Inches(0.5),
                "✓  DO", size=22, bold=True, color=GREEN)
    for i, t in enumerate(do_items):
        iy = Inches(2.55 + i * 0.65)
        add_rect(s, Inches(4.6), iy, Inches(4), Inches(0.55),
                 fill=RGBColor(0xF1, 0xF8, 0xE9),
                 line=GREEN, corner_radius=0.1)
        add_textbox(s, Inches(4.6), iy, Inches(4), Inches(0.55),
                    t, size=13, bold=True, color=GREEN_DK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # DO NOT column
    dont_items = ["Direct ECIF flow", "“Partnership”/“Alliance”", "Joint co-sell credit",
                  "Sell APEX itself"]
    bad_red = RGBColor(0xA0, 0x35, 0x35)
    add_textbox(s, Inches(8.85), Inches(1.95), Inches(4), Inches(0.5),
                "✕  DO NOT", size=22, bold=True, color=bad_red)
    for i, t in enumerate(dont_items):
        iy = Inches(2.55 + i * 0.65)
        add_rect(s, Inches(8.85), iy, Inches(4), Inches(0.55),
                 fill=RGBColor(0xFD, 0xF4, 0xF4),
                 line=bad_red, corner_radius=0.1)
        add_textbox(s, Inches(8.85), iy, Inches(4), Inches(0.55),
                    t, size=13, bold=True, color=bad_red,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 11, TOTAL)

    # ── Slide 12: NOT a product (binary) ───────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "What we sell  ·  What we don't",
                   "The single most important framing. Repeat it every meeting.")

    # NOT panel
    bad_red = RGBColor(0xA0, 0x35, 0x35)
    add_rect(s, Inches(0.5), Inches(1.95), Inches(6.1), Inches(4.8),
             fill=RGBColor(0xFD, 0xF4, 0xF4),
             line=bad_red, corner_radius=0.05)
    # Big X icon
    icon_x(s, Inches(3.0), Inches(2.30), Inches(1.1), Inches(1.1), color=bad_red)
    add_textbox(s, Inches(0.5), Inches(3.65), Inches(6.1), Inches(0.7),
                "NOT a product",
                size=28, bold=True, color=bad_red, align=PP_ALIGN.CENTER)
    nots = ["Not a SKU", "Not a license", "Not a subscription a client can buy",
            "Not Microsoft co-sell"]
    for i, t in enumerate(nots):
        add_textbox(s, Inches(1.0), Inches(4.6 + i * 0.45),
                    Inches(5.1), Inches(0.4),
                    "·  " + t, size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)

    # IS panel
    add_rect(s, Inches(6.85), Inches(1.95), Inches(6.0), Inches(4.8),
             fill=RGBColor(0xF1, 0xF8, 0xE9),
             line=GREEN, corner_radius=0.05)
    icon_check(s, Inches(9.30), Inches(2.30), Inches(1.1), Inches(1.1), color=GREEN)
    add_textbox(s, Inches(6.85), Inches(3.65), Inches(6.0), Inches(0.7),
                "Deloitte delivery framework",
                size=24, bold=True, color=GREEN_DK, align=PP_ALIGN.CENTER)
    iss = ["We sell  agentic-AI delivery services",
           "Scoped per Pack  ·  per tier  ·  per SOW",
           "Funded via the standard envelope mix",
           "Framework is Deloitte IP. Always."]
    for i, t in enumerate(iss):
        add_textbox(s, Inches(7.05), Inches(4.6 + i * 0.45),
                    Inches(5.6), Inches(0.4),
                    "·  " + t, size=14, color=STRONG_INK, align=PP_ALIGN.CENTER)
    page_footer(s, 12, TOTAL)

    # ── Slide 13: BVA Wedge — funnel diagram ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The wedge  —  BVA → Pack Lite → Standard → Enterprise → Operate",
                   "Each step earns the next. Nothing gets redone.")

    waves = [
        ("BVA",        "4 hrs",      "Free / client-funded", GREEN),
        ("Pack Lite",  "4-6 wks",    "$150-250K",            GREEN_DK),
        ("Standard",   "12-16 wks",  "$500K-1.5M",           HEAD_NAVY),
        ("Enterprise", "6-9 mo",     "$1.5-3.5M",            RGBColor(0x14, 0x1B, 0x4A)),
        ("Operate",    "Ongoing",    "$/mo subscription",    RGBColor(0x47, 0x55, 0x69)),
    ]
    # Horizontal arrow flow
    cw = Inches(2.45); gap = Inches(0.10); x0 = Inches(0.55)
    yc = Inches(3.0)
    for i, (name, time, fund, color) in enumerate(waves):
        cx = x0 + Inches(i * (cw / Inches(1) + gap / Inches(1)))
        # Arrow shape (last one is a regular rectangle)
        if i < len(waves) - 1:
            add_arrow(s, cx, yc, cw, Inches(1.6), fill=color)
        else:
            add_rect(s, cx, yc, cw, Inches(1.6), fill=color, corner_radius=0.08)
        # Number circle on top-left of each
        add_oval(s, cx + Inches(0.15), yc - Inches(0.30),
                 Inches(0.50), Inches(0.50), fill=GREEN)
        add_textbox(s, cx + Inches(0.15), yc - Inches(0.30),
                    Inches(0.50), Inches(0.50),
                    str(i + 1), size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name
        add_textbox(s, cx + Inches(0.10), yc + Inches(0.10),
                    cw - Inches(0.50), Inches(0.55),
                    name, size=18, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Time
        add_textbox(s, cx + Inches(0.10), yc + Inches(0.70),
                    cw - Inches(0.50), Inches(0.35),
                    time, size=11, color=LIGHT_TEXT,
                    align=PP_ALIGN.CENTER)
        # Funding pill
        add_rect(s, cx + Inches(0.25), yc + Inches(1.10),
                 cw - Inches(0.75), Inches(0.30),
                 fill=WHITE, corner_radius=0.5)
        add_textbox(s, cx + Inches(0.25), yc + Inches(1.10),
                    cw - Inches(0.75), Inches(0.30),
                    fund, size=9.5, bold=True, color=color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom callout
    add_textbox(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
                "Additive by design  —  Standard contains Lite, Enterprise contains Standard. No client pays to redo earlier work.",
                size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 13, TOTAL)

    # ── Slide 14: Q&A — speech bubbles ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    header_minimal(s, "The questions you'll hear",
                   "Six openers cover ~90% of the first ten minutes.")

    questions = [
        ("Is this Microsoft's product?",                       GREEN),
        ("Can we buy APEX?",                                   GREEN),
        ("How is this different from Foundry?",                HEAD_NAVY),
        ("What if we're on AWS or Google?",                    HEAD_NAVY),
        ("How fast can we see something working?",             GREEN),
        ("What does the audit story look like?",               HEAD_NAVY),
    ]
    # 2 cols × 3 rows of speech-bubble-styled cards
    cw = Inches(5.95); ch = Inches(1.45); gx = Inches(0.35); gy = Inches(0.20)
    x0 = Inches(0.50); y0 = Inches(2.0)
    for i, (q, color) in enumerate(questions):
        col = i % 2; row = i // 2
        cx = x0 + Inches(col * (cw / Inches(1) + gx / Inches(1)))
        cy = y0 + Inches(row * (ch / Inches(1) + gy / Inches(1)))
        # Card
        add_rect(s, cx, cy, cw, ch,
                 fill=WHITE, line=color, corner_radius=0.08)
        # "Q" badge
        add_oval(s, cx + Inches(0.25), cy + Inches(0.30),
                 Inches(0.85), Inches(0.85), fill=color)
        add_textbox(s, cx + Inches(0.25), cy + Inches(0.30),
                    Inches(0.85), Inches(0.85),
                    "Q", size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Question
        add_textbox(s, cx + Inches(1.25), cy + Inches(0.30),
                    cw - Inches(1.50), Inches(0.85),
                    q, size=15, bold=True, color=STRONG_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
    # Footer line
    add_textbox(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
                "If stuck for 30 seconds — redirect:  “Let me show you in the live demo, it's faster.”",
                size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    page_footer(s, 14, TOTAL)

    # ── Slide 15: Closing — three actions + momentum gradient ──────────
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "hero_closing.png"),
                         Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Dark panel for text legibility
    add_rect(s, Inches(0), Inches(0.6), SLIDE_W, Inches(5.8), fill=DEEP_NAVY)
    add_rect(s, Inches(0), Inches(0.6), Inches(0.20), Inches(5.8), fill=GREEN)
    add_textbox(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.8),
                "This week",
                size=56, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.45),
                "Three actions  ·  every DMTSP seller  ·  every week.",
                size=16, color=LIGHT_TEXT, italic=True)

    actions = [
        ("01", "OPEN  the  demo",      "60 minutes  ·  one client meeting"),
        ("02", "BOOK  the  BVA",       "4 hours  ·  client-funded"),
        ("03", "SIGN  Pack Lite SOW",  "DCIF  ·  $150-250K  ·  4-6 wks"),
    ]
    y0 = Inches(2.55)
    for i, (n, t, sub) in enumerate(actions):
        cy = y0 + Inches(i * 1.25)
        # Big number
        add_textbox(s, Inches(0.7), cy, Inches(1.8), Inches(1.1),
                    n, size=72, bold=True, color=GREEN,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_textbox(s, Inches(2.7), cy + Inches(0.10),
                    Inches(10), Inches(0.55),
                    t, size=30, bold=True, color=WHITE)
        # Sub
        add_textbox(s, Inches(2.7), cy + Inches(0.70),
                    Inches(10), Inches(0.40),
                    sub, size=14, color=LIGHT_TEXT, italic=True)

    # URL strip at very bottom
    add_rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.55),
             fill=GREEN)
    add_textbox(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.55),
                "ca-visionkit-portal.gentlestone-9b49b099.eastus2.azurecontainerapps.io/architecture",
                size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 15, TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"size: {OUT.stat().st_size:,} bytes · {TOTAL} slides")


if __name__ == "__main__":
    main()
